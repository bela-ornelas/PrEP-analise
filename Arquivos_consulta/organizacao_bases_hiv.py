"""
Programa de Organização das Bases do HIV

Arquivo: organizacao_bases_hiv.py

Descrição: Um programa Python com as principais funções para organizar e analisar bases de monitoramento clínico do HIV.

Autor: Tiago Benoliel - Assessoria de Monitoramento e Avaliação (AMA)
Organização: Departamento de HIV/Aids, Tuberculose, Hepatites Virais e ISTs (DATHI) - Ministério da Saúde

Data de Criação: 15/05/2024
Versão: 1.0.0
Repositório GitHub: [https://github.com/tbenoliel/Monitoramento-Clinico-HIV]
"""

import numpy as np
import pandas as pd
import datetime
from IPython.display import display
import funcoes_gerais as fg
import funcoes_linkage as fl
import os
import openpyxl
import traceback
import re
import gc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import unicodedata
import glob



def carregar_bases(hoje: datetime.date, Cadastro: bool = False, Disp: bool = False, CV: bool = False, CD4: bool = False, Geno: bool = False, SIM: bool = False,   Sinan: bool = False, Sinan_G: bool = False, Sinan_cr:bool = False, Cod_uni: bool = False,
                   colsuso_padrao:bool = True, Cadastro_colsuso: list = None, Disp_colsuso: list = None, CV_colsuso: list = None, CD4_colsuso: list = None, Geno_colsuso:list = None,
                   Obito_colsuso:list = None, data_base:datetime.date = None, SIM_colsuso: list = None, SIM_total_colsuso:list = None, Sinan_colsuso: list = None, Sinan_cong_colsuso: list = None, SinanG_colsuso: list = None, Sinan_cr_colsuso:list = None,
                   caminho_colunas: str = "//SAP109/Bancos AMA/Arquivos Atuais/Programas Atuais/Python/Versoes_Bancos de dados.xlsx", col_codigo_paciente: str = "codigo_paciente",
                   col_codigo_pac_uni: str = "Cod_unificado"):

    """
    Carrega e organiza os bancos de dados do SICLOM/SISCEL conforme a data de análise especificada.

    Retorna os DataFrames das bases de dados solicitadas na seguinte ordem:
        Pac, Disp, CV, CD4, obitos_siclom, obito_rel

    Parameters
    ----------
    hoje: datetime.date
        Data de fechamento dos bancos de análise no formato datetime.date(aaaa, mm, dd).

    Cadastro: bool, opcional, default = False
        Indica se a base de cadastro de pacientes (tb_paciente_consolidado) deve ser carregada.

    Disp: bool, opcional, default = False
        Indica se a base de dispensação de medicamentos (tb_dispensas_esquemas_udm) deve ser carregada.

    CV: bool, opcional, default = False
        Indica se a base de carga viral (tb_carga_viral_consolidado) deve ser carregada.

    CD4: bool, opcional, default = False
        Indica se a base de contagem de CD4 (tb_cd4_consolidado) deve ser carregada.

    Obito: bool, opcional, default = False
        Indica se a base de óbitos (tb_obitos_consolidado) deve ser carregada.

    Obito_rel: bool, opcional, default = False
        Indica se a base de óbitos relacionada do SIM deve ser carregada.

    Cadastro_colsuso: list, opcional, default = True
        Lista de colunas a serem utilizadas da base de cadastro. Se não especificado, usa uma lista padrão.
            [
            "codigo_paciente",
            "codigo_pac_eleito",
            "data_nascimento",
            "sexo",
            "raca",
            "codigo_ibge_resid",
            "uf_residencia",
            "acomp_medico",
            "escolaridade",
            "data_ult_atu",
            "co_orientacao_sexual",
            "co_genero"
            ]

    Disp_colsuso: list, opcional, default = True
        Lista de colunas a serem utilizadas da base de dispensação. Se não especificado, usa uma lista padrão.
            [
            "codigo_pac_eleito",
            "tp_servico_atendimento",
            "cod_ibge_udm",
            "nm_udm",
            'st_pub_priv',
            'data_dispensa',
            'esquema',
            'esquema_forma',
            'duracao'
            ]

    CV_colsuso: list, opcional, default = True
        Lista de colunas a serem utilizadas da base de carga viral. Se não especificado, usa uma lista padrão.
            [
            "codigo_pac_eleito",
            "cod_ibge_solicitante_cv",
            "nm_inst_sol_cv",
            "tipo_inst_sol_cv",
            'data_hora_coleta_cv',
            'copias',
            'comentario_copias',
            "data_solicitacao_cv"
            ]

    CD4_colsuso: list, opcional, default = True
        Lista de colunas a serem utilizadas da base de contagem de CD4. Se não especificado, usa uma lista padrão.
            [
            "codigo_pac_eleito",
            "cod_ibge_solicitante_cd4",
            "nm_inst_sol_cd4",
            "tipo_inst_sol_cd4",
            'data_solicitacao_cd4',
            'data_hora_coleta_cd4',
            "contagem_cd4"
            ]

    Obito_colsuso: list, opcional, default = True
        Lista de colunas a serem utilizadas da base de óbitos. Se não especificado, usa uma lista padrão.
            [
            "codigo_paciente",
            "data_obito",
            "data_comunicacao"
            ]

    caminho_colunas: str, opcional, default = "V:/Monitoramento e Avaliação/DOCUMENTOS/Python/Versoes_Bancos de dados.xlsx"
        Caminho para o arquivo Excel que contém as definições das colunas válidas para cada base.

    Returns
    -------
    tuple
        DataFrames das bases de dados solicitadas na seguinte ordem:
        Pac, Disp, CV, CD4, obitos_siclom, obito_rel, Cod_uni

    Example
    -------
    >>> hoje = pd.to_datetime(2024-01-31)
    >>> Pac, Disp, CV, CD4, obitos_siclom, obito_rel, Cod_uni = carregar_bases(hoje, Cadastro = True, CD4 = True, CV = True, Disp = True, Obito=True, Obito_rel=True, Cod_uni=True)
        Ou
    >>> Pac, Cod_uni = carregar_bases(hoje, Cadastro = True, Cod_uni=True)
    """

    def calcular_contagem_mes(ano, mes):
        return (ano - 2021) * 12 + mes - 2

    def obter_nome_mes(mes):
        meses_do_ano = ["", "Janeiro", "Fevereiro", "Março", "Abril", "Maio", "Junho", "Julho", "Agosto", "Setembro", "Outubro", "Novembro", "Dezembro"]
        return meses_do_ano[mes]

    if data_base is None:
        data_base = hoje
    
    ano = data_base.year
    mes = data_base.month
    contagem_mes = calcular_contagem_mes(ano, mes)
    mes_nome = obter_nome_mes(mes)
    caminho = f"V:/{ano}/Monitoramento e Avaliação/COMPARTILHADO/AMA - Banco de Dados/Consolidado/{contagem_mes} - {mes_nome} {ano}"

    if not os.path.exists(caminho):
        raise FileNotFoundError(f"O caminho {caminho} não existe.")

    # Organiza as colunas que serão utilizadas em cada DataFrame de acordo com a data passada
    Dic_variaveis = pd.read_excel(caminho_colunas, sheet_name=None)
    Dic_colunas = {}
    for aba, df in Dic_variaveis.items():
        colunas_validas = [coluna for coluna in df.columns if coluna <= hoje]
        if not colunas_validas:
            raise ValueError(f"Nenhuma coluna válida encontrada para a aba {aba} até a data {hoje}.")
        ultima_coluna_valida = max(colunas_validas)
        lista_de_variaveis = df[ultima_coluna_valida].dropna()
        Dic_colunas[aba] = lista_de_variaveis.tolist()

    bases = {}
    Colunas_Uso = {}

    def adicionar_base(nome, base, cols):
        bases[nome] = base
        Colunas_Uso[nome] = cols

    if Cadastro:
        if colsuso_padrao is True:
            Cadastro_colsuso = [
                "codigo_paciente", "codigo_pac_eleito", "data_nascimento", "sexo",
                "raca", "codigo_ibge_resid", "uf_residencia", "acomp_medico",
                "escolaridade", "data_ult_atu","cd_pais", "st_estrangeiro", 
                "co_orientacao_sexual", "co_genero"
            ]
        adicionar_base("Pac","tb_paciente_consolidado", Cadastro_colsuso)

    if Disp:
        if colsuso_padrao is True:
            Disp_colsuso = [
                "codigo_paciente", "tp_servico_atendimento", "codigo_udm", "cod_ibge_udm","nm_udm",
                'st_pub_priv', 'data_dispensa', 'esquema', 'esquema_forma', 'duracao'
            ]
        adicionar_base("Disp","tb_dispensas_esquemas_udm", Disp_colsuso)

    if CV:
        if colsuso_padrao is True:
            CV_colsuso = [
                "codigo_paciente", "cod_ibge_solicitante_cv", "cd_inst_sol_cv", "nm_inst_sol_cv", "tipo_inst_sol_cv",
                'data_hora_coleta_cv', 'copias', 'comentario_copias', "data_solicitacao_cv"
            ]
        adicionar_base("CV","tb_carga_viral_consolidado", CV_colsuso)

    if CD4:
        if colsuso_padrao is True:
            CD4_colsuso = [
                "codigo_paciente", "cod_ibge_solicitante_cd4", "cd_inst_sol_cd4", "nm_inst_sol_cd4", "tipo_inst_sol_cd4",
                'data_solicitacao_cd4', 'data_hora_coleta_cd4', "contagem_cd4"
            ]
        adicionar_base("CD4","tb_cd4_consolidado", CD4_colsuso)

    if Geno:
        adicionar_base("Geno","tb_sisgeno", Geno_colsuso)

    if SIM:
        if colsuso_padrao is True:
            Obito_colsuso = [
                "codigo_paciente", "data_obito"
            ]
        adicionar_base("obito_siclom","tb_obitos_consolidado", Obito_colsuso)



    DFs = {}
    for (base,base), (nome, cols) in zip(bases.items(), Colunas_Uso.items()):
        colunas = Dic_colunas.get(base)
        if not colunas:
            raise ValueError(f"Colunas não encontradas para a base {base}.")
        df = pd.read_csv(f"{caminho}/{base}.txt", sep="\t", names=colunas,
                        usecols=cols, low_memory=True, index_col=False, engine="c", quoting=3,
                        encoding="latin-1", on_bad_lines="warn")
        DFs[nome] = df

    if SIM:
        if colsuso_padrao is True:
            SIM_colsuso = [
                "NOME","NOMEMAE","CODMUNRES","DTOBITO","DTNASC","SEXO","RACACOR","CODMUNOCOR","id_sim","Cod_unificado"
            ]

        sim_path = "B:/Bancos Atuais HIV/Relacionados/sim_hiv_uni.csv"
        if not os.path.exists(sim_path):
            raise FileNotFoundError(f"O caminho {sim_path} não existe.")
        sim = pd.read_csv(sim_path, sep=";", encoding="latin-1",
                        on_bad_lines="warn", low_memory = False,
                        usecols = SIM_colsuso) 
        sim[col_codigo_pac_uni] = sim[col_codigo_pac_uni].astype(int)
        # Adicionar zero inicial para datas com 7 dígitos e converte para o formato AAAA-MM-DD
        DFs["Sim"] = sim

        if colsuso_padrao is True:
            SIM_total_colsuso = [
                "DTOBITO",
                "dt_nasc",
                "sexo",
                "RACACOR",
                "ibge_res",
                "codigo_paciente"
            ]

        sim_total = pd.read_csv("B:/Bancos Atuais HIV/Relacionados/sim_completo_uni.csv",
                                sep=";", encoding="latin-1", usecols = SIM_total_colsuso,
                                on_bad_lines="warn", low_memory = True)
        
        DFs["Sim_total"] = sim_total


    if Sinan:
        if colsuso_padrao is True:
            Sinan_colsuso = [
                "NM_PACIENT","NM_MAE_PAC","DT_NOTIFIC","CS_RACA","CS_ESCOLAR","DT_DIAG","DT_NASC","CS_SEXO","ID_MN_RESI","ID_MUNICIP","CRITERIO","CS_ESCOL_N","DT_CONFIRM",
                "DT_SIN_PRI","DT_RAPIDO","id_sinan","Cod_unificado"
            ]

            Sinan_cong_colsuso = [
                'NM_PACIENT', 'NM_MAE_PAC', 'DT_NASC', 'ID_MN_RESI',
                'DT_DIAG', 'DT_NOTIFIC', 'ID_MUNICIP', 'CS_SEXO', 'CRITERIO', 'ANT_REL_CA', 'CS_RACA',
                'CS_ESCOLAR', 'ANODIAG','id_sinan', 'Cod_unificado'
            ]

        sinan = pd.read_csv("B:/Bancos Atuais HIV/Relacionados/sinan_adulto_uni.csv",
                                    sep=";", encoding="latin-1", usecols = Sinan_colsuso,
                                    on_bad_lines="warn", low_memory = True)
            
        DFs["Sinan"] = sinan

        sinan_cong = pd.read_csv("B:/Bancos Atuais HIV/Relacionados/sinan_congelado_uni.csv",
                            sep=";", encoding="latin-1", usecols = Sinan_cong_colsuso,
                            on_bad_lines="warn", low_memory = True)

        DFs["Sinan_cong"] = sinan_cong


    if Sinan_G:
        if colsuso_padrao is True:
            SinanG_colsuso  = [
                "NM_PACIENT","NM_MAE_PAC","DT_NOTIFIC","CS_RACA","CS_ESCOLAR","DT_DIAG","DT_NASC","CS_SEXO","DT_SIN_PRI",
                "ID_MN_RESI","ID_MUNICIP","CS_ESCOL_N","id_sinan","Cod_unificado"
            ]

        sinan_g = pd.read_csv("B:/Bancos Atuais HIV/Relacionados/sinan_gestante_uni.csv",
                                    sep=";", encoding="latin-1", usecols = SinanG_colsuso,
                                    on_bad_lines="warn", low_memory = True)
            
        DFs["Sinan_G"] = sinan_g

    
    if Sinan_cr:
        if colsuso_padrao is True:
            Sinan_cr_colsuso = [
                "NM_PACIENT","NM_MAE_PAC","DT_NOTIFIC","CS_RACA","CS_ESCOLAR","DT_DIAG","DT_NASC","CS_SEXO",
                "ID_MN_RESI","ID_MUNICIP","CRITERIO","CS_ESCOL_N","DT_SIN_PRI","DTCONFIRMA", "DTRAPIDO1","DT_LAB_HIV",
                "DT_PCR_1","DT_PCR_2","DT_PCR_3","id_sinan","Cod_unificado"]

        sinan_cr = pd.read_csv("B:/Bancos Atuais HIV/Relacionados/sinan_crianca_uni.csv",
                                    sep=";", encoding="latin-1", usecols = Sinan_cr_colsuso,
                                    on_bad_lines="warn", low_memory = True)
            
        DFs["Sinan_cr"] = sinan_cr


    if Cod_uni:
        cod_uni_path = "B:/Bancos Atuais HIV/Relacionados/cadastro_hiv_uni.csv"
        if not os.path.exists(cod_uni_path):
            raise FileNotFoundError(f"O caminho {cod_uni_path} não existe.")
        Cod_uni = pd.read_csv(cod_uni_path, sep=";", encoding="latin-1", on_bad_lines="warn") 
        Cod_uni[col_codigo_pac_uni] = Cod_uni[col_codigo_pac_uni].astype(int)
        DFs["Cod_uni"] = Cod_uni[[col_codigo_paciente,col_codigo_pac_uni,"nm_pac","nm_mae","data_nascimento","sexo","codigo_ibge_resid","num_cpf_pac"]]


    if len(DFs) == 1:  # Se houver apenas um DataFrame, retorne-o diretamente
        return next(iter(DFs.values()))

    else:
        return DFs.values()



def carregar_bases_PVHA(hoje: datetime.date, PVHA: bool = False, Cadastro: bool = False, Disp: bool = False, CV: bool = False, CD4: bool = False, Sinan: bool = False, Sinan_G: bool = False,
                        Sinan_cr:bool = False, Sim:bool = False, colsuso_padrao:bool = True, Cadastro_colsuso: list = None, Disp_colsuso: list = None,
                        CV_colsuso: list = None, CD4_colsuso: list = None, Sim_colsuso: list = None, Sinan_colsuso: list = None, Sinan_cong_colsuso: list = None, SinanG_colsuso: list = None,
                        Sinan_cr_colsuso:list = None, col_codigo_pac_uni: str = "Cod_unificado"):

    ano = hoje.year

    caminho = f"V:{ano}/Monitoramento e Avaliação/COMPARTILHADO/AMA - Banco de Dados/AMA-VIP/Bancos Compartilhados HIV"

    def carregar_csv(nome:str, cols:list):
        arquivo = f"{caminho}/{nome}.csv"
        if not os.path.exists(arquivo):
            raise FileNotFoundError(f"O caminho {arquivo} não existe.")
        DF = pd.read_csv(arquivo, sep=";", encoding="latin-1", on_bad_lines="warn", usecols=cols) 
        DF[col_codigo_pac_uni] = DF[col_codigo_pac_uni].astype(int)
        return DF
         
    DFs = {}

    if PVHA:
        PVHA = carregar_csv("PVHA",None)
        DFs["PVHA"] = PVHA

    if Cadastro:
        if colsuso_padrao is True:
            Cadastro_colsuso = [
                "codigo_paciente", col_codigo_pac_uni, "codigo_ibge_resid", "uf_residencia", "acomp_medico",
                "escolaridade", "data_ult_atu","cd_pais", 
                "co_orientacao_sexual", "co_genero","st_paciente"
            ]
        Pac = carregar_csv("Cadastro",Cadastro_colsuso)
        DFs["Cadastro"] = Pac

    if Disp:
        if colsuso_padrao is True:
            Disp_colsuso = [
                col_codigo_pac_uni, "tp_servico_atendimento", "codigo_udm", "cod_ibge_udm","nm_udm",
                'st_pub_priv', 'data_dispensa', 'esquema', 'esquema_forma', 'duracao', "cd_crm", "uf_crm"
            ]
        Disp = carregar_csv("Dispensacao",Disp_colsuso)
        DFs["Dispensacao"] = Disp

    if CV:
        if colsuso_padrao is True:
            CV_colsuso = [
                col_codigo_pac_uni, "cod_ibge_solicitante_cv", "cd_inst_sol_cv", "nm_inst_sol_cv", "tipo_inst_sol_cv",
                'data_hora_coleta_cv', 'copias', 'comentario_copias', "data_solicitacao_cv"
            ]
        CV = carregar_csv("CV",CV_colsuso)
        DFs["CV"] = CV

    if CD4:
        if colsuso_padrao is True:
            CD4_colsuso = [
                col_codigo_pac_uni, "cod_ibge_solicitante_cd4", "cd_inst_sol_cd4", "nm_inst_sol_cd4", "tipo_inst_sol_cd4",
                'data_solicitacao_cd4', 'data_hora_coleta_cd4', "contagem_cd4"
            ]
        CD4 = carregar_csv("CD4",CD4_colsuso)
        DFs["CD4"] = CD4

    if Sinan:
        if colsuso_padrao is True:
            Sinan_colsuso = [
                "DT_NOTIFIC","CS_ESCOLAR","ID_MN_RESI","ID_MUNICIP","CS_ESCOL_N","Cod_unificado"
            ]
            
            Sinan_cong_colsuso = [
                "DT_NOTIFIC","CS_ESCOLAR","ID_MN_RESI","ID_MUNICIP","Cod_unificado"
            ]

        sinan = carregar_csv("Sinan_hiv_adulto",Sinan_colsuso)
        DFs["Sinan_A"] = sinan

        sinan_cong = carregar_csv("Sinan_hiv_congelado",Sinan_cong_colsuso)
        DFs["Sinan_cong"] = sinan_cong


    if Sinan_G:
        if colsuso_padrao is True:
            SinanG_colsuso  = [
                "DT_NOTIFIC","CS_ESCOLAR","ID_MN_RESI","ID_MUNICIP","CS_ESCOL_N","Cod_unificado"
            ]

        sinan_g = carregar_csv("Sinan_hiv_gestante",SinanG_colsuso)
        DFs["Sinan_G"] = sinan_g

    
    if Sinan_cr:
        if colsuso_padrao is True:
            Sinan_cr_colsuso = [
                "DT_NOTIFIC","CS_ESCOLAR","ID_MN_RESI","ID_MUNICIP","CS_ESCOL_N","Cod_unificado"
                ]

        sinan_cr = carregar_csv("Sinan_hiv_crianca",Sinan_cr_colsuso)            
        DFs["Sinan_cr"] = sinan_cr


    if Sim:
        if colsuso_padrao is True:
            Sim_colsuso = [
                "CODMUNRES","CODMUNOCOR","CODMUNCART","Cod_unificado"
                ]

        sim = carregar_csv("SIM_aids",Sim_colsuso)            
        DFs["Sim"] = sim


    if len(DFs) == 1:  # Se houver apenas um DataFrame, retorne-o diretamente
        return next(iter(DFs.values()))

    else:
        return DFs.values()



def organizacao_cadastro(PVHA: pd.DataFrame, Pac: pd.DataFrame, hoje: datetime.date, col_escolaridade: str = "escolaridade", col_raca: str = "Raca_cat",
                        col_codigo_ibge_resid: str = "codigo_ibge_resid", col_data_ult_atu: str = "data_ult_atu", col_codigo_pac_uni: str = "Cod_unificado",
                        col_genero:str = "co_genero", col_orientacao:str = "co_orientacao_sexual", datas:list = ["data_min","data_nascimento", "data_obito"]):
    """
    Essa função organiza o banco de cadastro de pacientes do SICLOM/SISCEL (tb_paciente_consolidado).
    Retorna o dataframe Pac com as seguintes modificações:

        1) Criação de coluna ["Idade_cat_{ano}"] com idade na data de hoje em categorias.
        2) Padronização da Raça, Escolaridade e Sexo em categorias pré-definidas.
        3) Ajuste dos cadastros com uf de residência inválidos.
        4) Se usar_obito = True, une a data de óbito ao cadastro a partir das bases (obito_rel e obitos_siclom).

        
    Parameters
    ----------
    Pac: pd.DataFrame, default = Pac
        DataFrame original com a base de cadastros (tb_paciente_consolidado).

    hoje: datetime.date, default = hoje
        Data de fechamento dos bancos de análise no formato pd.to_datetime(aaaa-mm-dd).

    usar_obito: bool, opcional, default = False
        Indicar se deseja incluir a data de óbito na base.

    obito_rel: DataFrame, opcional, default = None
        Base de dados com as datas de óbito obtidas a partir do banco relacionado do SIM.
        Obrigatório se usar_obito = True.

    obitos_siclom: DataFrame, opcional, default = None
        Base de dados com as datas de óbitos do SICLOM.
        Obrigatório se usar_obito = True.

    col_escolaridade: str, opcional, default = "escolaridade"
        Nome da coluna que contém a informação sobre a escolaridade.

    col_raca: str, opcional, default = "raca"
        Nome da coluna que contém a informação sobre a raça.

    col_sexo: str, opcional, default = "sexo"
        Nome da coluna que contém a informação sobre o sexo.

    col_data_nascimento: str, opcional, default = "data_nascimento"
        Nome da coluna que contém a informação sobre a data de nascimento.

    col_uf_residencia: str, opcional, default = "uf_residencia"
        Nome da coluna que contém a informação sobre a UF de residência.

    col_codigo_ibge_resid: str, opcional, default = "codigo_ibge_resid"
        Nome da coluna que contém a informação sobre o código IBGE de residência.

    col_codigo_paciente: str, opcional, default = "codigo_paciente"
        Nome da coluna que contém a informação sobre o código do paciente.

    col_data_ult_atu: str, opcional, default = "data_ult_atu"
        Nome da coluna que contém a informação sobre a data da última atualização do cadastro.

    col_codigo_pac_uni: str, opcional, default = "codigo_pac_eleito"
        Nome da coluna que contém a informação sobre o código do paciente eleito.

    col_dtobito_sim: str, opcional, default = "dtobito_sim"
        Nome da coluna que contém a informação sobre a data de óbito do SIM.

    col_data_obito: str, opcional, default = "data_obito"
        Nome da coluna que contém a informação sobre a data de óbito do SICLOM.

    col_data_comunicacao: str, opcional, default = "data_comunicacao"
        Nome da coluna que contém a informação sobre a data de comunicação.


    Returns
    -------
    pd.DataFrame
        O DataFrame original com as modificações descritas acima.
        
    """

    Pac[col_codigo_ibge_resid] = pd.to_numeric(Pac[col_codigo_ibge_resid], errors='coerce')
    PVHA[col_codigo_pac_uni] = PVHA[col_codigo_pac_uni].astype(int)
    Pac[col_codigo_pac_uni] = Pac[col_codigo_pac_uni].astype(int)
    PVHA["Hoje"] = hoje

    # Ajusta as Datas
    datas_pac = [col_data_ult_atu]

    for col in datas_pac:
        Pac = fg.ajusta_data_linha_vetorizado(Pac, col, coluna_retorno = f"{col}_ajustada")
        Pac[col] = Pac[f"{col}_ajustada"]
        Pac = Pac.drop(columns = [f"{col}_ajustada"])


    Escol = [
        (Pac[col_escolaridade] == 1) | (Pac[col_escolaridade] == 2) | (Pac[col_escolaridade] == 10) | (Pac[col_escolaridade] == 11) | (Pac[col_escolaridade] == 12) | (Pac[col_escolaridade] == 17) | (Pac[col_escolaridade] == 18) | (Pac[col_escolaridade] == 19),
        (Pac[col_escolaridade] == 3) | (Pac[col_escolaridade] == 4) | (Pac[col_escolaridade] == 13) | (Pac[col_escolaridade] == 20),
        (Pac[col_escolaridade] == 5) | (Pac[col_escolaridade] == 6) | (Pac[col_escolaridade] == 7) | (Pac[col_escolaridade] == 8) | (Pac[col_escolaridade] == 9) | (Pac[col_escolaridade] == 14) | (Pac[col_escolaridade] == 21)
    ]
    
    Escol_escolha2 = [1,
                     2,
                     3]
     
    Pac["Escol_num"] = np.select(Escol, Escol_escolha2, default=None)
    display(Pac["Escol_num"].value_counts(dropna = False))
    print()

    Pac[col_genero] = Pac[col_genero].apply(lambda x:str(x).strip())
    Pac["Genero_cat"] = np.where(Pac[col_genero] == "-",None, Pac[col_genero])
    display(Pac["Genero_cat"].value_counts(dropna = False))
    print()

    Pac[col_orientacao] = Pac[col_orientacao].apply(lambda x:str(x).strip())
    Pac["Orientacao_cat"] = np.where(Pac[col_orientacao] == "-",None, Pac[col_orientacao])
    display(Pac["Orientacao_cat"].value_counts(dropna = False))
    print()

    Pac["st_paciente"] = Pac["st_paciente"].apply(lambda x:str(x).strip())
    stat = [
        ((Pac["st_paciente"] == "Paciente Ativo") |
        (Pac["st_paciente"] == "Paciente em abandono") |
        (Pac["st_paciente"] == "Conclusao do Tratamento") |
        (Pac["st_paciente"] == "Paciente em Tarv") |
        (Pac["st_paciente"] == "Fim de Gestação") |
        (Pac["st_paciente"] == "Suspensão de tratamento") |
        (Pac["st_paciente"] == "Interrupção de tratamento") |
        (Pac["st_paciente"] == "Paciente Transferido para unidade sem SICLOM")
        ),
        (Pac["st_paciente"] == "Mudança de pais"),
        (Pac["st_paciente"] == "Cadastro Indevido"),
        (Pac["st_paciente"] == "Paciente Duplicado")
    ]


    stat_escolha = [
        1,
        2,
        3,
        4
    ]
    
    Pac["st_num"] = np.select(stat, stat_escolha, default=1)
    display(Pac["st_num"].value_counts(dropna = False))
    print()

 
    # Preenche as colunas em branco com valores do outros registros
    cols = ["Escol_num", "Orientacao_cat", "Genero_cat","st_num"]
    funcoes = ["max","moda","moda","min"]

    # Primeiro, ordene os dados por 'Cod_unificado' e 'data_ref'

    for var,funcao in zip(cols,funcoes):
        Pac = fg.padronizar_variaveis_vetorizado(Pac, var = var, col_data_ref = col_data_ult_atu, funcao = funcao, mais_antigo = False)
    
    cond = [
        (Pac["Genero_cat"] == "Homem CIS") & (Pac["Orientacao_cat"] == "Heterossexual"),
        (Pac["Genero_cat"] == "Homem CIS") & ((Pac["Orientacao_cat"] == "Homossexual / Gay / Lésbica") | (Pac["Orientacao_cat"] == "Bissexual")),
        (Pac["Genero_cat"] == "Mulher Transexual") | (Pac["Genero_cat"] == "Travesti"),
        (Pac["Genero_cat"] == "Homem Transexual"),
        (Pac["Genero_cat"] == "Não binário"),
        (Pac["Genero_cat"] == "Mulher CIS") & (Pac["Orientacao_cat"] == "Heterossexual"),
        (Pac["Genero_cat"] == "Mulher CIS") & (Pac["Orientacao_cat"] == "Homossexual / Gay / Lésbica") | (Pac["Orientacao_cat"] == "Bissexual")
    ]

    resultados = [
        "Homem CIS hetero",
        "Homem CIS HSH",
        "Mulheres TRANS e Travestis",
        "Homens TRANS",
        "Pessoas não binárias",
        "Mulheres CIS hetero",
        "Mulheres CIS homo/bissexual"
    ]

    Pac["Pop_genero"] = np.select(cond, resultados, default = "Não Informado")
    display(Pac["Pop_genero"].value_counts())
    print()

    if "cd_pais" in Pac.columns:
        cod_pais = {
            4: 'AFEGANISTAO',
            8: 'ALBANIA',
            12: 'ARGELIA',
            16: 'SAMOAAMERICANA',
            20: 'ANDORRA',
            24: 'ANGOLA',
            28: 'ANTIGUAEBARBUDA',
            31: 'AZERBAIJAO',
            32: 'ARGENTINA',
            36: 'AUSTRALIA',
            38: 'INGLATERRA',
            40: 'AUSTRIA',
            44: 'BAHAMAS',
            48: 'BAREINE',
            50: 'BANGLADESH',
            51: 'ARMENIA',
            52: 'BARBADOS',
            56: 'BELGICA',
            60: 'BERMUDA',
            64: 'BUTAO',
            68: 'BOLIVIA',
            70: 'BOSNIA-HERZEGOVINA',
            72: 'BOTSUANA',
            76: 'BRASIL',
            84: 'BELIZE',
            90: 'ILHASSALOMAO',
            92: 'ILHASVIRGENSBRITANICAS',
            96: 'BRUNEI',
            100: 'BULGARIA',
            104: 'MIANMA',
            108: 'BURUNDI',
            112: 'BIELORUSSIA',
            116: 'CAMBOJA',
            120: 'REPUBLICADOSCAMAROES',
            124: 'CANADA',
            132: 'CABOVERDE',
            136: 'ILHASCAYMAN',
            140: 'REPCENTRO-AFRICANA',
            144: 'SRILANKA',
            148: 'CHADE',
            152: 'CHILE',
            156: 'CHINA',
            158: 'FORMOSA',
            170: 'COLOMBIA',
            174: 'COMORES',
            178: 'CONGO',
            180: 'REPDEMODOCONGO',
            184: 'ILHASCOOK',
            188: 'COSTARICA',
            191: 'CROACIA',
            192: 'CUBA',
            196: 'CHIPRE',
            203: 'REPUBLICATCHECA',
            204: 'BENIN',
            208: 'DINAMARCA',
            212: 'DOMINICA',
            214: 'REPUBLICADOMINICANA',
            218: 'EQUADOR',
            222: 'ELSALVADOR',
            226: 'GUINEEQUATORIAL',
            231: 'ETIOPIA',
            232: 'ERITREIA',
            233: 'ESTONIA',
            234: 'ILHASFAEROE',
            238: 'ILHASFALKLANDS',
            242: 'FIJI',
            246: 'FINLANDIA',
            250: 'FRANCA',
            254: 'GUIANAFRANCESA',
            258: 'POLINESIAFRANCESA',
            262: 'DJIBUTI',
            266: 'GABAO',
            268: 'GEORGIA',
            270: 'GAMBIA',
            274: 'FAIXADEGAZA',
            276: 'ALEMANHA',
            288: 'GANA',
            292: 'GIBRALTAR',
            296: 'QUIRIBATI',
            300: 'GRECIA',
            304: 'GROENLANDIA',
            308: 'GRANADA',
            312: 'GUADALUPE',
            316: 'GUAM',
            320: 'GUATEMALA',
            324: 'GUINE',
            328: 'GUIANA',
            332: 'HAITI',
            336: 'VATICANO',
            340: 'HONDURAS',
            344: 'HONGKONG',
            348: 'HUNGRIA',
            352: 'ISLANDIA',
            356: 'INDIA',
            360: 'INDONESIA',
            364: 'IRA',
            368: 'IRAQUE',
            372: 'IRLANDA',
            376: 'ISRAEL',
            380: 'ITALIA',
            384: 'COSTADOMARFIM',
            388: 'JAMAICA',
            392: 'JAPAO',
            398: 'CAZAQUISTAO',
            400: 'JORDANIA',
            404: 'QUENIA',
            408: 'REPDEMOPOPULARCOREIA',
            410: 'COREIA',
            414: 'KUAITE',
            417: 'QUIRGUIZIA',
            418: 'LAOS',
            422: 'LIBANO',
            426: 'LESOTO',
            428: 'LETONIA',
            430: 'LIBERIA',
            434: 'LIBIA',
            438: 'LIECHTENSTEIN',
            440: 'LITUANIA',
            442: 'LUXEMBURGO',
            446: 'MACAU',
            450: 'MADAGASCAR',
            454: 'MALAUI',
            458: 'MALASIA',
            462: 'MALDIVAS',
            466: 'MALI',
            470: 'MALTA',
            474: 'MARTINICA',
            478: 'MAURITANIA',
            480: 'MAURICIO',
            484: 'MEXICO',
            492: 'MONACO',
            496: 'MONGOLIA',
            498: 'MOLDAVIA',
            500: 'MONTSERRAT',
            504: 'MARROCOS',
            508: 'MOCAMBIQUE',
            512: 'OMA',
            516: 'NAMIBIA',
            520: 'NAURU',
            524: 'NEPAL',
            528: 'HOLANDA',
            530: 'ANTILHASHOLANDESAS',
            533: 'ARUBA',
            540: 'NOVACALEDONIA',
            548: 'VANUATU',
            554: 'NOVAZELANDIA',
            558: 'NICARAGUA',
            562: 'NIGER',
            566: 'NIGERIA',
            570: 'NIUE',
            574: 'ILHANORFOLK',
            578: 'NORUEGA',
            580: 'ILHASMARIANASDONORTE',
            583: 'MICRONESIA',
            584: 'ILHASMARSHALL',
            585: 'PALAU',
            586: 'PAQUISTAO',
            591: 'PANAMA',
            598: 'PAPUANOVAGUINE',
            600: 'PARAGUAI',
            604: 'PERU',
            608: 'FILIPINAS',
            612: 'ILHAPITCAIRN',
            616: 'POLONIA',
            620: 'PORTUGAL',
            624: 'GUINEBISSAU',
            626: 'TIMORLESTE',
            630: 'PORTORICO',
            634: 'CATAR',
            638: 'ILHASREUNIAO',
            642: 'ROMENIA',
            643: 'RUSSIA',
            646: 'RUANDA',
            654: 'SANTAHELENA',
            659: 'SAOCRISTOVAOENEVIS',
            660: 'ANGUILLA',
            662: 'SANTALUCIA',
            666: 'SAINTPIERREEMIQUELON',
            670: 'SAOVICENTEEGRANADINAS',
            674: 'SAOMARINO',
            678: 'SAOTOMEEPRINCIPE',
            682: 'ARABIASAUDITA',
            686: 'SENEGAL',
            690: 'SEICHELES',
            694: 'SERRALEOA',
            702: 'CINGAPURA',
            703: 'ESLOVAQUIA',
            704: 'VIETNA',
            705: 'ESLOVENIA',
            706: 'SOMALIA',
            710: 'AFRICADOSUL',
            716: 'ZIMBABUE',
            724: 'ESPANHA',
            732: 'SAARADOOESTE',
            736: 'SUDAO',
            740: 'SURINAME',
            744: 'ILHASSVALBARDEJANMAYEN',
            748: 'SUAZILANDIA',
            752: 'SUECIA',
            756: 'SUICA',
            760: 'SIRIA',
            762: 'TADJIQUISTAO',
            764: 'TAILANDIA',
            768: 'TOGO',
            772: 'TOKELAU',
            776: 'TONGA',
            780: 'TRINIDADETOBAGO',
            784: 'EMIRADOSARABESUNIDOS',
            788: 'TUNISIA',
            792: 'TURQUIA',
            795: 'TURCOMENISTAO',
            796: 'ILHATURKSECAICOS',
            798: 'TUVALU',
            800: 'UGANDA',
            804: 'UCRANIA',
            807: 'MACEDONIA',
            818: 'EGITO',
            826: 'REINOUNIDO',
            830: 'ILHASCANAL',
            833: 'ILHADEMAN',
            834: 'REPUNIDADATANZANIA',
            840: 'ESTADOSUNIDOS',
            850: 'ILHASVIRGENSAMERICANAS',
            854: 'BURKINAFASO',
            858: 'URUGUAI',
            860: 'UZBEQUISTAO',
            862: 'VENEZUELA',
            876: 'ILHASWALLISEFUTUNA',
            882: 'SAMOA',
            887: 'IEMEN',
            891: 'IUGOSLAVIA',
            894: 'ZAMBIA',
            895: 'BELARUS',
            999: 'PACIENTESEMCIDADESEMPAIS'
            }
        
        Pac["cd_pais"] = pd.to_numeric(Pac["cd_pais"], errors='coerce')
        Pac['pais'] =  Pac["cd_pais"].map(cod_pais)


    Pac.sort_values( by = col_data_ult_atu, ascending = False, inplace = True)
    Pac = Pac.drop_duplicates(col_codigo_pac_uni, keep = "first").copy()


    Rac = ((PVHA[col_raca] == "Branca") | (PVHA[col_raca] == "Amarela"))      
    Rac_escolha = "Branca/Amarela"
    PVHA["Raca_cat2"] = np.where(Rac, Rac_escolha, PVHA[col_raca])
    display(PVHA["Raca_cat2"].value_counts(dropna = False))
    print()

    for col in datas:
        PVHA = fg.ajusta_data_linha_vetorizado(PVHA, col, coluna_retorno = f"{col}_ajustada")
        PVHA[col] = PVHA[f"{col}_ajustada"]
        PVHA = PVHA.drop(columns = [f"{col}_ajustada"])

    PVHA["ano_min"] = PVHA["data_min"].dt.year
    PVHA["ano_obito"] = PVHA["data_obito"].dt.year
    
    # Calcula a idade de vinculação
    PVHA["Idade_vinc_cat"] = fg.idade_cat(DF = PVHA, data_ref="data_min", data_nasc="data_nascimento")
    display(PVHA["Idade_vinc_cat"].value_counts().sort_index())
    print()

    # Calcula a idade de vinculação nas faixas anteriores
    PVHA["Idade_vinc_cat_anterior"] = fg.idade_cat(DF = PVHA, data_ref="data_min", data_nasc="data_nascimento", faixas_etarias="MC_antigo")
    display(PVHA["Idade_vinc_cat_anterior"].value_counts().sort_index())
    print()

    # Calcula a idade atual
    PVHA["Idade_atual"] = fg.idade_cat(DF = PVHA, data_ref="Hoje", data_nasc="data_nascimento")
    display(PVHA["Idade_atual"].value_counts().sort_index())
    print()

    display(round(pd.crosstab(index = PVHA["ano_min"], columns = PVHA["origem_data_min"], normalize = 0)*100,1))
    print()
    print()

    display(pd.crosstab(index = PVHA["ano_min"], columns = PVHA["origem_data_min"]))

    # Une o cadasto ao código Unificado
    PVHA = pd.merge(PVHA,Pac, on = col_codigo_pac_uni, how = "left")

    return PVHA



def organizacao_disp(hoje: datetime.date, Disp: pd.DataFrame, Cod_tab: pd.DataFrame = None,  Prim: bool = False, var_prim:list = [], esquemas: bool = True,
                     col_data_dispensa: str = "data_dispensa", col_duracao: str = "duracao", col_codigo_pac_uni: str = "Cod_unificado", 
                     col_codigo_paciente: str = "codigo_paciente", col_esquema: str = "esquema", col_nm_udm:str = "nm_udm"):
    """
    Essa função organiza a base de dispensações de ARV do SICLOM.
    Retorna o dataframe Disp com as seguintes modificações:

    1) Padronização das datas de dispensação.
    2) Une a duração de dispensações feitas no mesmo dia [Duracao_sum] e exclui as dispensações duplicadas (no mesmo dia).
    3) Calcula o atraso em cada dispensação [Atraso] e cria o status da pessoa até a dispensação seguinte, ou até a data atual [Status_cadaDisp].
    4) Cria colunas [Abandono_sum,Atraso_sum] com o número de dias que a pessoa ficou em atraso e o número de eventos de abandono no ano.
    5) Cria colunas com a razão entre o número de dias em atraso e o número de dispensações no ano [Atraso_raz]
    e entre o número de eventos de abanono e o número de dispensações no ano [Abandono_raz].
    6) Categoria o número de dias em atraso no ano [Atraso_cat] e a duração das dispensações [Duracao_cat].
    7) Ordena a base de forma descendente a partir da data da dispensação.
    8) Se Prim = True, salva as variáveis na lista var_prim em novas colunas com os dados da primeira dispensação da vida ["coluna_Prim"].
    9) Se esquemas = True, padroniza os esquemas no [Esquema_AMA] e classifica as classes de medicamentos.
    10) Se status_ano = True, calcula se a pessoa estava em TARV ou Abandono ao final de cada ano [Status_ano].
 
    
    Parameters
    ----------
    Disp: pd.DataFrame, default = Disp
        DataFrame original com a base de dispensações de ARV (tb_dispensas_esquemas_udm).

    hoje: datetime.date, default = hoje
        Data de fechamento dos bancos de análise no formato pd.to_datetime(aaaa-mm-dd).

    Prim: bool, opcional, default = False
        Decide se os dados da primeira dispensação marcados em uma coluna própria com o sufixo "_Prim".

    var_prim: list, opcional, default = None
        Lista de variáveis da primeira dispensação da vida que devem ser salvas em colunas própias com o sufico "_Prim".
        Só deve ser utilizada quando Prim = True.

    esquemas: bool, opcional, default = True
        Decide se deve organizar os esquemas, criando o [Esquema_AMA], e as classificações de Classes.

    status_ano: bool, opcional, default = True
        Decide se deve criar o status TARV/Abandono por ano, calculado no último dia de cada ano da dispensação.

    col_data_dispensa: str, opcional, default = "data_dispensa"
        Nome da coluna que contém a data de dispensação.

    col_duracao: str, opcional, default = "duracao"
        Nome da coluna que contém a duração da dispensação.

    col_codigo_pac_eleito: str, opcional, default = "codigo_pac_eleito"
        Nome da coluna que contém o código do paciente eleito.

    col_esquema: str, opcional, default = "esquema"
        Nome da coluna que contém o esquema de medicação.
        

    """

    # Trabalhando as datas
    Disp[col_data_dispensa] = pd.to_datetime(Disp[col_data_dispensa], errors = "coerce")
    Disp[col_data_dispensa] = Disp[col_data_dispensa].dt.normalize()
    Disp = Disp[Disp[col_data_dispensa] <= hoje].copy()
    Disp["ano_disp"]= Disp[col_data_dispensa].dt.year
    Disp["mesN_disp"]= Disp[col_data_dispensa].dt.month
    ano = hoje.year
    display(Disp["ano_disp"].value_counts().sort_index())
    print()
    
    if Cod_tab is not None:
        Disp[col_codigo_paciente] = Disp[col_codigo_paciente].astype(int)
        # Une o cadasto ao código Unificado
        Disp = pd.merge(Disp,Cod_tab, on = col_codigo_paciente, how = "left")

    Disp[col_codigo_pac_uni] = Disp[col_codigo_pac_uni].astype("Int64")

    #Altera a abreviação do mês para o Português
    month_mapping = {
        "Jan": "Jan",
        "Feb": "Fev",
        "Mar": "Mar",
        "Apr": "Abr",
        "May": "Mai",
        "Jun": "Jun",
        "Jul": "Jul",
        "Aug": "Ago",
        "Sep": "Set",
        "Oct": "Out",
        "Nov": "Nov",
        "Dec": "Dez",
    }
    
    Disp["mes_disp"] = Disp[col_data_dispensa].dt.month_name().str[:3].map(month_mapping)
    
    Disp["mes_ano"] = Disp["mes_disp"] + "/" + Disp["ano_disp"].apply(str)
    Disp["mes_ano"]
    
    # Padroniza o nome das UDMs
    Disp[col_nm_udm] = Disp[col_nm_udm].apply(lambda x:str(x).upper().strip())
    Disp[col_nm_udm] = np.where((Disp[col_nm_udm] == "UDM INATIVADA") | (Disp[col_nm_udm] == "UDM TESTE") | (Disp[col_nm_udm] == "UDM TESTE III"), np.nan, Disp[col_nm_udm])
    
    Disp[col_duracao] = pd.to_numeric(Disp[col_duracao], errors="coerce")

    # Soma a duração de cod_pac e dt_disp iguais.
    Disp["duracao_sum"] = Disp.groupby([col_codigo_pac_uni, col_data_dispensa])[col_duracao].transform("sum")
    
    # ordenar por cod_pac e dt_disp
    Disp.sort_values([col_codigo_pac_uni, col_data_dispensa], ascending = [True, False], inplace = True)
    
    # Manter apenas uma entrada de dispensa por data com a duração correta
    Disp.drop_duplicates(subset=[col_codigo_pac_uni, col_data_dispensa], keep = "first", inplace = True)
    
    # Criação da data teórica da próxima dispensa
    Disp["dt_lim_prox_disp"] = Disp[col_data_dispensa] + pd.to_timedelta(Disp["duracao_sum"], unit="D")
    Disp["dt_lim_prox_disp"] = pd.to_datetime(Disp["dt_lim_prox_disp"], errors="coerce")
    Disp["dt_lim_prox_disp"] = Disp["dt_lim_prox_disp"].dt.normalize()
    
    ## A Coluna "Status" mostra a situação que o paciente ficou até a realização da próxima dispensação.
    # Criação da data real da dispensa seguinte
    Disp.reset_index(drop = True, inplace = True)
    Disp["dt_real_disp_seguinte"] = Disp.groupby(col_codigo_pac_uni)[col_data_dispensa].shift(1)
    
    # Preenchimento da data real da próxima dispensa da última com a data de hoje
    Disp["dt_real_disp_seguinte"] = Disp["dt_real_disp_seguinte"].fillna(hoje)
    Disp["dt_real_disp_seguinte"] = pd.to_datetime(Disp["dt_real_disp_seguinte"], errors="coerce")
    Disp["dt_real_disp_seguinte"] = Disp["dt_real_disp_seguinte"].dt.normalize()
    
    # Criação da coluna de atraso
    Disp["Atraso"] = (Disp["dt_real_disp_seguinte"] - Disp["dt_lim_prox_disp"]).dt.days
    
    # Criação da coluna de status
    Disp["Status_cadaDisp"] = np.where(Disp["Atraso"] <= 60, "Tarv", "Interrupção de Tarv")
    display(Disp["Status_cadaDisp"].value_counts())
    print()
    
    # Cria a coluna binária do Abandono
    Disp["Abandono"] = np.where(Disp["Status_cadaDisp"] == "Interrupção de Tarv", 1, 0)
    
    # Cria colunas para o somatório dos eventos no ano
    Disp["Abandono_sum"] = Disp.groupby([col_codigo_pac_uni,"ano_disp"])["Abandono"].transform("sum")
    Disp["N_dispensas"] = Disp.groupby([col_codigo_pac_uni,"ano_disp"])[col_codigo_pac_uni].transform("count")
    Disp["Atraso_sum"] = Disp.groupby([col_codigo_pac_uni,"ano_disp"])["Atraso"].transform("sum")
    

    # Cria uma máscara para selecionar registros dos últimos 12 meses
    corte_12mes = hoje - pd.DateOffset(months=12)
    mask_12_meses = Disp['data_dispensa'] >= corte_12mes

    # Atualiza os somatórios apenas para o ano vigente
    Disp.loc[Disp["ano_disp"] == ano, "Abandono_sum"] = Disp[(mask_12_meses)].groupby(col_codigo_pac_uni)["Abandono"].transform("sum")
    Disp.loc[Disp["ano_disp"] == ano, "Atraso_sum"] = Disp[(mask_12_meses)].groupby(col_codigo_pac_uni)["Atraso"].transform("sum")


    # Cria uma razão entre os eventos de abandono e o número de dispensações no ano e entre a soma dos dias em atraso e o número de dispensações no ano
    Disp["Abandono_raz"] = Disp["Abandono_sum"]/Disp["N_dispensas"]
    Disp["Atraso_raz"] = Disp["Atraso_sum"]/Disp["N_dispensas"]
    
    # Crias as categorias de estratificação para a soma do Atraso
    cond_status = [
        (Disp["Atraso_sum"] < 10),
        (Disp["Atraso_sum"] < 30),
        (Disp["Atraso_sum"] < 60),
        (Disp["Atraso_sum"] < 90),
        (Disp["Atraso_sum"] >= 90)
       ]
    
    escolha_status = ["Menos que 10 dias", "Entre 10 e 30 dias", "Entre 30 e 60 dias",
                     "Entre 60 e 90 dias","Mais de 90 dias"]
    
    Disp["Atraso_cat"] = np.select(cond_status, escolha_status, default="Erro")
    
    display(Disp["Atraso_cat"].value_counts())
    print()
    
    # Cria as categorias de ocorrência de interrupção de TARV no ano
    Disp["Abandono_ano"] = np.where(Disp["Abandono_sum"] == 0, "Sem interrupção no ano", "Com interrupção no ano")

    # Crias as categorias de estratificação para a soma da duração
    cond_status = [
        (Disp["duracao_sum"] == 30),
        (Disp["duracao_sum"] == 60),
        (Disp["duracao_sum"] == 90),
        (Disp["duracao_sum"] == 120)
       ]
    
    escolha_status = [
        "30 dias", "60 dias", "90 dias","120 dias"
        ]
    
    Disp["Duracao_cat"] = np.select(cond_status, escolha_status, default="Outros")
    
    display(Disp["Duracao_cat"].value_counts())
    print()


    if esquemas:

        def Esquema_AMA(x):
                valores = []
        
                if "3TC" in x or "AZL" in x or "TL" in x or "TLE" in x or "DTL" in x:
                        valores.append("3TC")
                if "ABC" in x:
                        valores.append("ABC")
                if "APV" in x: 
                        valores.append("APV")
                if "ATV" in x:
                        valores.append("ATV")
                if "AZT" in x or "AZL" in x:
                        valores.append("AZT")
                if "DRV" in x:
                        valores.append("DRV")
                if "d4T" in x:
                        valores.append("d4T")
                if "ddl" in x or "EC" in x:
                        valores.append("ddl")
                if "DTG" in x or "DTL" in x:
                        valores.append("DTG")
                if "EFZ" in x or "TLE" in x:
                        valores.append("EFZ")
                if "ETR" in x:
                        valores.append("ETR")
                if "FPV" in x:
                        valores.append("FPV")
                if "IDV" in x:
                        valores.append("IDV")
                if "LPV" in x:
                        valores.append("LPV")
                if "MVQ" in x:
                        valores.append("MVQ")
                if "NVP" in x:
                        valores.append("NVP")
                if "RAL" in x:
                        valores.append("RAL")
                if "RTV" in x or "LPV" in x:
                        valores.append("RTV")
                if "SQV" in x:
                        valores.append("SQV")
                if "T20" in x:
                        valores.append("T20")
                if "TDF" in x or " TL" in x or "+TL" in x or "TLE" in x:
                        valores.append("TDF")
                if "TPV" in x:
                        valores.append("TPV")
                if "NFV" in x:
                        valores.append("NFV")
                if "TAF" in x:
                        valores.append("TAF")
        
                return "+".join(valores)
        
        Disp["esquema_AMA"] = Disp[col_esquema].apply(lambda x: Esquema_AMA(x))
        display(Disp["esquema_AMA"].value_counts())
        
        print()        

        def Classe_Med(x):
                valores = []
        
                if "3TC" in x or "ABC" in x or "AZT" in x or "d4T" in x or "ddl" in x or "TDF" in x or "TAF" in x:
                        valores.append("ITRN")
                if "EFZ" in x or "ETR" in x or "NVP" in x:
                        valores.append("ITRNN")
                if "APV" in x or "ATV" in x or "DRV" in x or "FPV" in x or "IDV" in x or "LPV" in x or "RTV" in x or "SQV" in x or "TPV" in x or "NFV" in x: 
                        valores.append("IP")
                if "DTG" in x or "RAL" in x:
                        valores.append("INI")
                if "MVQ" in x:
                        valores.append("Inibidor CCR5")
                if "T20" in x:
                        valores.append("Inibidor de Fusão")
        
                return "+".join(valores)
        
        Disp["Classe_Med"] = Disp["esquema_AMA"].apply(lambda x: Classe_Med(x))
        display(Disp["Classe_Med"].value_counts())
        
        print()
        
        def Classe_Med_count(x):
                class_counts = {
                "ITRN": 0,
                "ITRNN": 0,
                "IP": 0,
                "INI": 0,
                "Inibidor CCR5": 0,
                "Inibidor de Fusão": 0
                }
        
                # Dividir o esquema em medicamentos individuais
                medicamentos = x.split("+")
        
                # Verificar cada medicamento e contar as classes
                for medicamento in medicamentos:
                        if "3TC" in medicamento or "ABC" in medicamento or "AZT" in medicamento or "d4T" in medicamento or "ddl" in medicamento or "TDF" in medicamento or "TAF" in medicamento:
                                class_counts["ITRN"] += 1
                        if "EFZ" in medicamento or "ETR" in medicamento or "NVP" in medicamento:
                                class_counts["ITRNN"] += 1
                        if "APV" in medicamento or "ATV" in medicamento or "DRV" in medicamento or "FPV" in medicamento or "IDV" in medicamento or "LPV" in medicamento or "RTV" in medicamento or "SQV" in medicamento or "TPV" in medicamento or "NFV" in medicamento: 
                                class_counts["IP"] += 1
                        if "DTG" in medicamento or "RAL" in medicamento:
                                class_counts["INI"] += 1
                        if "MVQ" in medicamento:
                                class_counts["Inibidor CCR5"] += 1
                        if "T20" in medicamento:
                                class_counts["Inibidor de Fusão"] += 1
        
                class_str = "+".join([f"{cls}({count})" for cls, count in class_counts.items() if count > 0])
        
                return class_str
        
        # Aplicar a função a coluna "esquema_AMA"
        Disp["Classe_Med_Count"] = Disp["esquema_AMA"].apply(lambda x: Classe_Med_count(x))
        display(Disp["Classe_Med_Count"].value_counts())
        
        print()

        # Criar um dicionário de mapeamento para substituição com os 5 esquemas mais frequente no ano vigente
        categorias_principais = Disp[Disp["ano_disp"] == ano]["esquema_AMA"].value_counts().head(5).index

        # Criar a nova coluna com apenas os 5 esquemas
        Disp['esquema_cat'] = np.where(Disp['esquema_AMA'].isin(categorias_principais), Disp['esquema_AMA'], "Outro Esquema")
    

    if Prim:
        Disp.sort_values([col_data_dispensa], ascending = [True], inplace = True)
        Disp.reset_index(drop = True, inplace = True)

        if esquemas:
            Disp["esquemaAMA_Prim"] = Disp.groupby([col_codigo_pac_uni])["esquema_AMA"].transform("first")
            Disp["data_esquemaAMA_Prim"] = Disp.groupby([col_codigo_pac_uni,"esquema_AMA"])[col_data_dispensa].transform("first")

        
        Disp["data_dispensa_Prim"] = Disp.groupby([col_codigo_pac_uni])[col_data_dispensa].transform("first")
        for col in var_prim:
            Disp[f"{col}_Prim"] = Disp.groupby([col_codigo_pac_uni])[col].transform("first")

        Disp["ano_disp_Prim"] = Disp["data_dispensa_Prim"].dt.year


    return Disp

    """
    # Bloco para criação do Atraso_cat_ano, onde há o somatório dos dias em atraso de dispensação apenas do ano
    # Criação da coluna de atraso no ano

    Disp["Ult_dia_ano_disp"] = pd.to_datetime(Disp["ano_disp"].astype(str) + "-12-31")
    Disp["Ult_dia_ano_disp"] = Disp["Ult_dia_ano_disp"].dt.normalize()

    Disp["Atraso_no_Ano"] = (Disp["dt_real_disp_seguinte"] - Disp["dt_lim_prox_disp"]).dt.days
    Disp["Atraso_ultdia_Ano"] = (Disp["Ult_dia_ano_disp"] - Disp["dt_lim_prox_disp"]).dt.days
    Disp["Atraso_dif_ano"] = Disp["Atraso_no_Ano"] - Disp["Atraso_ultdia_Ano"]
    Disp["Atraso_dif_ano"] = Disp.groupby(col_codigo_pac_uni)["Atraso_dif_ano"].shift(-1)



    Disp.loc[Disp["dt_real_disp_seguinte"].dt.year > Disp["ano_disp"], "Atraso_no_Ano"] = (Disp["Ult_dia_ano_disp"] - Disp["dt_lim_prox_disp"]).dt.days
    
    Disp["Atraso_ano_sum"] = Disp.groupby([col_codigo_pac_uni,"ano_disp"])["Atraso_no_Ano"].transform("sum")
    Disp.loc[Disp["dt_real_disp_seguinte"].dt.year > Disp["ano_disp"], "Atraso_ano_sum"] =  Disp.groupby([col_codigo_pac_uni,"ano_disp"])["Atraso_no_Ano"].transform("sum") + Disp["Atraso_dif_ano"]
    

    # Cria uma máscara para selecionar registros dos últimos 12 meses
    corte_12mes = hoje - pd.DateOffset(months=12)
    mask_12_meses = Disp['data_dispensa'] >= corte_12mes

    # Atualiza os somatórios apenas para o ano vigente
    Disp.loc[Disp["ano_disp"] == ano, "Abandono_sum"] = Disp[(mask_12_meses)].groupby(col_codigo_pac_uni)["Abandono"].transform("sum")
    Disp.loc[Disp["ano_disp"] == ano, "Atraso_sum"] = Disp[(mask_12_meses)].groupby(col_codigo_pac_uni)["Atraso"].transform("sum")


    # Cria uma razão entre os eventos de abandono e o número de dispensações no ano e entre a soma dos dias em atraso e o número de dispensações no ano
    Disp["Abandono_raz"] = Disp["Abandono_sum"]/Disp["N_dispensas"]
    Disp["Atraso_raz"] = Disp["Atraso_sum"]/Disp["N_dispensas"]
    
    # Crias as categorias de estratificação para a soma do Atraso
    cond_status = [
        (Disp["Atraso_sum"] < 10),
        (Disp["Atraso_sum"] < 30),
        (Disp["Atraso_sum"] < 60),
        (Disp["Atraso_sum"] < 90),
        (Disp["Atraso_sum"] >= 90)
       ]
    
    escolha_status = ["Menos que 10 dias", "Entre 10 e 30 dias", "Entre 30 e 60 dias",
                     "Entre 60 e 90 dias","Mais de 90 dias"]
    
    Disp["Atraso_cat"] = np.select(cond_status, escolha_status, default="Erro")
    
    display(Disp["Atraso_cat"].value_counts())
    print()
    
    # Cria as categorias de ocorrência de interrupção de TARV no ano
    Disp["Abandono_ano"] = np.where(Disp["Abandono_sum"] == 0, "Sem interrupção no ano", "Com interrupção no ano")

    """



def organizacao_cv(hoje: datetime.date, CV: pd.DataFrame, Cod_tab: pd.DataFrame = None, categorizar:bool = True, Prim: bool = False,
                col_data_hora_coleta_cv: str = "data_hora_coleta_cv", col_data_solicitacao_cv: str = "data_solicitacao_cv",
                col_codigo_pac_uni: str = "Cod_unificado", col_codigo_paciente: str = "codigo_paciente", col_copias: str = "copias", col_comentario_copias: str = "comentario_copias",
                col_nome_inst_sol_cv:str = "nm_inst_sol_cv", col_cod_ibge_solicitante_cv:str = "cod_ibge_solicitante_cv"):
    
    """
    Essa função organiza a base de exames de carga viral do SISCEL.
    Retorna o dataframe CV com as seguintes modificações:

    1) Padronização das datas de coleta e solicitação de CV.
    2) Cálcula o número de exames de CV realizados por pessoa em cada ano.
    3) Se categorizar = True, gera 3 classificações de categorias do número de cópias.
    4) Ordena a base de forma descendente a partir da data da coleta.
    5) Se Prim = True, salva as variáveis de cópias e categoriza em novas colunas os dados do primeiro exame de CV da vida.
    
    Parameters
    ----------
    CV: pd.DataFrame, default = CV
        DataFrame original com a base de carga viral (tb_carga_viral_consolidado).

    hoje: datetime.date, default = hoje
        Data de fechamento dos bancos de análise no formato pd.to_datetime(aaaa-mm-dd).

    Prim: bool, opcional, default = False
        Decide se os dados da primeira carga viral serão salvos em uma coluna própria com o sufixo "_Prim".

    categorizar: bool, opcional, default = True
        Decide se serão criadas categorias para o resultado das cópias de carga viral.

    col_data_hora_coleta_cv: str, opcional, default = "data_hora_coleta_cv"
        Nome da coluna que contém a data e hora da coleta de carga viral.

    col_data_solicitacao_cv: str, opcional, default = "data_solicitacao_cv"
        Nome da coluna que contém a data de solicitação de carga viral.

    col_codigo_pac_eleito: str, opcional, default = "codigo_pac_eleito"
        Nome da coluna que contém o código do paciente eleito.

    col_copias: str, opcional, default = "copias"
        Nome da coluna que contém o número de cópias de carga viral.

    col_comentario_copias: str, opcional, default = "comentario_copias"
        Nome da coluna que contém o comentário sobre as cópias de carga viral.
        
    """

    CV[col_data_hora_coleta_cv] = pd.to_datetime(CV[col_data_hora_coleta_cv],errors="coerce").dt.normalize()
    CV[col_data_solicitacao_cv] = pd.to_datetime(CV[col_data_solicitacao_cv],errors="coerce").dt.normalize()
    CV["ano_coleta_cv"]= CV[col_data_hora_coleta_cv].dt.year
    CV["mesN_coleta_cv"]= CV[col_data_hora_coleta_cv].dt.month
    CV["mes_coleta_cv"]= CV[col_data_hora_coleta_cv].dt.month_name().str[:3]
    display(CV["ano_coleta_cv"].value_counts().sort_index())
    print()

    if Cod_tab is not None:
        CV[col_codigo_paciente] = CV[col_codigo_paciente].astype(int)
        # Une o cadasto ao código Unificado
        CV = pd.merge(CV,Cod_tab, on = col_codigo_paciente, how = "left")

    CV[col_codigo_pac_uni] = CV[col_codigo_pac_uni].astype("Int64")

    CV = CV[CV[col_data_hora_coleta_cv] <= hoje].copy()
    
    # Conta o número de exames realizados no ano
    CV["N_exames_ano_cv"] = CV.groupby(by=[col_codigo_pac_uni, "ano_coleta_cv"])["mes_coleta_cv"].transform("count")

    # Padroniza o nome das instituições solicitantes de CV
    CV[col_nome_inst_sol_cv] = CV[col_nome_inst_sol_cv].apply(lambda x:str(x).upper().strip())
    CV[col_nome_inst_sol_cv] = np.where(CV[col_nome_inst_sol_cv] == "APAGAR - SEM USO", np.nan, CV[col_nome_inst_sol_cv])
    CV.loc[CV[col_nome_inst_sol_cv] == "LACEN/PARANÁ - UNIDADE DE FRONTEIRA", col_cod_ibge_solicitante_cv] = 4108304

    # Contagem de quantas CV detectável na vida
    CV["CV_detec_count"] = CV[
        ((CV[col_copias].isnull() == True) & (CV[col_comentario_copias] == 2)) | (CV[col_copias] >= 50)
    ].groupby(col_codigo_pac_uni)[col_data_hora_coleta_cv].transform("count")

    # Contagem de quantas CV >5000 na vida
    CV["CV_5000_count"] = CV[
        ((CV[col_copias].isnull() == True) & (CV[col_comentario_copias] == 2)) | (CV[col_copias] >= 5000)
    ].groupby(col_codigo_pac_uni)[col_data_hora_coleta_cv].transform("count")
    
    if categorizar:
        CV[col_copias] = pd.to_numeric(CV[col_copias], errors='coerce')
        CV[col_comentario_copias] = pd.to_numeric(CV[col_comentario_copias], errors='coerce')
        
        cond_cv = [
                (((CV[col_copias].isnull() == True) & (CV[col_comentario_copias] == 0)) |
                ((CV[col_copias].isnull() == True) & (CV[col_comentario_copias] == 1)) |
                (CV[col_copias] < 50)),
                (((CV[col_copias].isnull() == True) & (CV[col_comentario_copias] == 2)) |
                (CV[col_copias] >= 10000)),
                (CV[col_copias] >= 50) & (CV[col_copias] < 200),
                (CV[col_copias] >= 200) & (CV[col_copias] < 1000),
                (CV[col_copias] >= 1000) & (CV[col_copias] < 10000)
        ]
        
        escolhas_cv = [
                "<50","10.000+","50-199","200-999","1.000-9.999"
        ]
        
        CV["CV_cat"] = np.select(cond_cv,escolhas_cv, default = "Não disponível")
        display(CV["CV_cat"].value_counts())
        print()
        
        cond_cv = [
                (CV["CV_cat"] == "1.000-9.999") | (CV["CV_cat"] == "10.000+"),
                (CV["CV_cat"] == "<50") | (CV["CV_cat"] == "50-199") | (CV["CV_cat"] == "200-999")   
        ]
        
        escolhas_cv = [
                ">=1.000",
                "<1.000"
        ]
        
        CV["CV_cat1000"] = np.select(cond_cv,escolhas_cv, default = "Não disponível")
        display(CV["CV_cat1000"].value_counts())
        print()
        
        cond_cv = [
                (CV["CV_cat"] == "<50"),
                (CV["CV_cat"] == "1.000-9.999") | (CV["CV_cat"] == "10.000+") |
                (CV["CV_cat"] == "50-199") | (CV["CV_cat"] == "200-999")   
        ]
        
        escolhas_cv = [
                "<50",
                ">=50"
        ]
        
        CV["CV_cat50"] = np.select(cond_cv,escolhas_cv, default = "Não disponível")
        display(CV["CV_cat50"].value_counts())
        print()

        cond_cv = [
                (CV["CV_cat"] == "<50") | (CV["CV_cat"] == "50-199"),
                (CV["CV_cat"] == "1.000-9.999") | (CV["CV_cat"] == "10.000+") |
                (CV["CV_cat"] == "200-999")   
        ]
        
        escolhas_cv = [
                "<200",
                ">=200"
        ]
        
        CV["CV_cat200"] = np.select(cond_cv,escolhas_cv, default = "Não disponível")
        display(CV["CV_cat200"].value_counts())
        print()

        cond_cv = [
                (((CV["copias"].isnull() == True) & (CV["comentario_copias"] == 0)) |
                ((CV["copias"].isnull() == True) & (CV["comentario_copias"] == 1)) |
                (CV["copias"] < 500)),
                (((CV["copias"].isnull() == True) & (CV["comentario_copias"] == 2)) |
                (CV["copias"] >= 500))
        ]

        escolhas_cv = [
                "<500",">=500"
        ]

        CV["CV_cat500"] = np.select(cond_cv,escolhas_cv, default = "Não disponível")
        display(CV["CV_cat500"].value_counts())
        print()


    if Prim:    
        
        # Cria a primeira data de coleta
        CV.sort_values([col_data_hora_coleta_cv], ascending = [True], inplace = True)
        CV.reset_index(drop = True, inplace = True)
        CV["data_coleta_cv_Prim"] = CV.groupby([col_codigo_pac_uni])[col_data_hora_coleta_cv].transform("first")
        CV["copia_cv_Prim"] = CV.groupby([col_codigo_pac_uni])[col_copias].transform("first")
        CV["comentario_copias_Prim"] = CV.groupby([col_codigo_pac_uni])[col_comentario_copias].transform("first")
        CV["data_sol_cv_Prim"] = CV.groupby([col_codigo_pac_uni])[col_data_solicitacao_cv].transform("first")

        CV["ano_coleta_cv_Prim"] = CV["data_coleta_cv_Prim"].dt.year

        # Cria a classificação da primeira CV
        CV["copia_cv_Prim"] = pd.to_numeric(CV["copia_cv_Prim"], errors='coerce')
        CV["comentario_copias_Prim"] = pd.to_numeric(CV["comentario_copias_Prim"], errors='coerce')
        
        cond_cv = [
                (((CV["copia_cv_Prim"].isnull() == True) & (CV["comentario_copias_Prim"] == 0)) |
                ((CV["copia_cv_Prim"].isnull() == True) & (CV["comentario_copias_Prim"] == 1)) |
                (CV["copia_cv_Prim"] < 50)),
                (((CV["copia_cv_Prim"].isnull() == True) & (CV["comentario_copias_Prim"] == 2)) |
                (CV["copia_cv_Prim"] >= 10000)),
                (CV["copia_cv_Prim"] >= 50) & (CV["copia_cv_Prim"] < 200),
                (CV["copia_cv_Prim"] >= 200) & (CV["copia_cv_Prim"] < 1000),
                (CV["copia_cv_Prim"] >= 1000) & (CV["copia_cv_Prim"] < 10000)
        ]
        
        escolhas_cv = [
                "<50","10.000+","50-199","200-999","1.000-9.999"
        ]
        
        CV["CV_cat_Prim"] = np.select(cond_cv,escolhas_cv, default = "Não disponível")
        display(CV["CV_cat_Prim"].value_counts())

    return CV



def organizacao_cd4(hoje: datetime.date, CD4: pd.DataFrame, Cod_tab:pd.DataFrame = None, categorizar:bool = True, Prim: bool = False,
                col_data_hora_coleta_cd4: str = "data_hora_coleta_cd4", col_data_solicitacao_cd4: str = "data_solicitacao_cd4",
                col_codigo_pac_uni: str = "Cod_unificado", col_codigo_paciente: str = "codigo_paciente", col_contagem_cd4: str = "contagem_cd4",
                col_nome_inst_sol_cd4:str = "nm_inst_sol_cd4", col_cod_ibge_solicitante_cd4:str = "cod_ibge_solicitante_cd4"):
    """
    Essa função organiza a base de exames de CD4 do SISCEL.
    Retorna o dataframe CD4 com as seguintes modificações:

    1) Padronização das datas de coleta e solicitação de CD4.
    2) Cálcula o número de exames de CD4 realizados por pessoa em cada ano.
    3) Ordena a base de forma descendente a partir da data da coleta.
    4) Se categorizar = True, gera 2 classificações de categorias da contagem de CD4.
    5) Se Prim = True, salva as variáveis de contagem de CD4 e categoriza em novas colunas os dados do primeiro exame de CD4 da vida.
    
    Parameters
    ----------
    CD4: pd.DataFrame, default = CD4
        DataFrame original com a base de CD4 (tb_cd4_consolidado).

    hoje: datetime.date, default = hoje
        Data de fechamento dos bancos de análise no formato pd.to_datetime(aaaa-mm-dd).

    Prim: bool, opcional, default = False
        Decide se os dados do primeiro CD4 serão salvos em uma coluna própria com o sufixo "_Prim".

    categorizar: bool, opcional, default = True
        Decide se serão criadas categorias para o resultado da contagem de CD4.

    col_data_hora_coleta_cd4: str, opcional, default = "data_hora_coleta_cd4"
        Nome da coluna que contém a data e hora da coleta de CD4.

    col_data_solicitacao_cd4: str, opcional, default = "data_solicitacao_cd4"
        Nome da coluna que contém a data de solicitação de CD4.

    col_codigo_pac_eleito: str, opcional, default = "codigo_pac_eleito"
        Nome da coluna que contém o código do paciente eleito.

    col_contagem_cd4: str, opcional, default = "contagem_cd4"
        Nome da coluna que contém a contagem de CD4.


    """
    CD4[col_data_hora_coleta_cd4] = pd.to_datetime(CD4[col_data_hora_coleta_cd4],errors="coerce")
    CD4[col_data_hora_coleta_cd4] = CD4[col_data_hora_coleta_cd4].dt.normalize()
    CD4[col_data_solicitacao_cd4] = pd.to_datetime(CD4[col_data_solicitacao_cd4],errors="coerce")
    CD4[col_data_solicitacao_cd4] = CD4[col_data_solicitacao_cd4].dt.normalize()
    CD4["ano_coleta_cd4"]= CD4[col_data_hora_coleta_cd4].dt.year
    CD4["mesN_coleta_cd4"]= CD4[col_data_hora_coleta_cd4].dt.month
    CD4["mes_coleta_cd4"]= CD4[col_data_hora_coleta_cd4].dt.month_name().str[:3]
    display(CD4["ano_coleta_cd4"].value_counts().sort_index())
    print()

    if Cod_tab is not None:
        CD4[col_codigo_paciente] = CD4[col_codigo_paciente].astype(int)
        # Une o cadasto ao código Unificado
        CD4 = pd.merge(CD4,Cod_tab, on = col_codigo_paciente, how = "left")

    CD4[col_codigo_pac_uni] = CD4[col_codigo_pac_uni].astype("Int64")
    
    CD4 = CD4[CD4[col_data_hora_coleta_cd4] <= hoje].copy()
    
    # Conta o número de exames realizados no ano
    CD4["N_exames_ano_cd4"] = CD4.groupby([col_codigo_pac_uni,"ano_coleta_cd4"])["mes_coleta_cd4"].transform("count")


    # Padroniza o nome das instituições solicitantes de CV
    CD4[col_nome_inst_sol_cd4] = CD4[col_nome_inst_sol_cd4].apply(lambda x:str(x).upper().strip())
    CD4[col_nome_inst_sol_cd4] = np.where(CD4[col_nome_inst_sol_cd4] == "APAGAR - SEM USO", np.nan, CD4[col_nome_inst_sol_cd4])
    CD4.loc[CD4[col_nome_inst_sol_cd4] == "LACEN/PARANÁ - UNIDADE DE FRONTEIRA", col_cod_ibge_solicitante_cd4] = 4108304

    if categorizar is True:
        CD4[col_contagem_cd4] = pd.to_numeric(CD4[col_contagem_cd4], errors="coerce")

        # criando variável de categorias do valor do CD4
        Cond_CD4cat = [
            (CD4[col_contagem_cd4] > 0) & (CD4[col_contagem_cd4] < 200),
            (CD4[col_contagem_cd4] >= 200) & (CD4[col_contagem_cd4] < 350),
            (CD4[col_contagem_cd4] >= 350) & (CD4[col_contagem_cd4] < 500),
            (CD4[col_contagem_cd4] >= 500) & (CD4[col_contagem_cd4] < 2500)
        ]

        Escolha_CD4cat = ["0-199", "200-349", "350-499", "500+"]

        CD4["CD4_cat"] = np.select(Cond_CD4cat, Escolha_CD4cat, default="Sem resultado de CD4")

        display(CD4["CD4_cat"].value_counts(dropna = False).sort_index())
        print()

        # criando variável de categorias do valor do CD4 menor que 200
        Cond_CD4cat2 = [
                (CD4[col_contagem_cd4] > 0) & (CD4[col_contagem_cd4] < 200),
                (CD4[col_contagem_cd4] >= 200) & (CD4[col_contagem_cd4] < 2500)
        ]

        Escolha_CD4cat2 = ["<200 células / mm³", "≥200 células / mm³"]

        CD4["CD4_cat200"] = np.select(Cond_CD4cat2, Escolha_CD4cat2, default="Sem resultado de CD4")

        display(CD4["CD4_cat200"].value_counts(dropna = False).sort_index())
        print()


        # criando variável de categorias do valor do CD4 menor que 200
        Cond_CD4cat2 = [
                (CD4["contagem_cd4"] > 0) & (CD4["contagem_cd4"] < 350),
                (CD4["contagem_cd4"] >= 350) & (CD4["contagem_cd4"] < 2500)
        ]

        Escolha_CD4cat2 = ["<350 células / mm³", "≥350 células / mm³"]

        CD4["CD4_cat350"] = np.select(Cond_CD4cat2, Escolha_CD4cat2, default="Sem resultado de CD4")

        display(CD4["CD4_cat350"].value_counts(dropna = False).sort_index())
        print()


    if Prim is True:
        
        # Cria a primeira data de coleta e solicitação
        CD4.sort_values([col_data_hora_coleta_cd4], ascending = [True], inplace = True)
        CD4.reset_index(drop = True, inplace = True)
        CD4["data_coleta_cd4_Prim"] = CD4.groupby([col_codigo_pac_uni])[col_data_hora_coleta_cd4].transform("first")
        CD4[f"{col_contagem_cd4}_Prim"] = CD4.groupby([col_codigo_pac_uni])[col_contagem_cd4].transform("first")
        CD4["data_sol_cd4_Prim"] = CD4.groupby([col_codigo_pac_uni])[col_data_solicitacao_cd4].transform("first")

        CD4[f"{col_contagem_cd4}_Prim"] = pd.to_numeric(CD4[f"{col_contagem_cd4}_Prim"], errors="coerce")

        CD4["ano_coleta_cd4_Prim"] = CD4["data_coleta_cd4_Prim"].dt.year


        # criando variável de categorias do valor do primeiro CD4
        Cond_CD4cat = [
                (CD4[f"{col_contagem_cd4}_Prim"] > 0) & (CD4[f"{col_contagem_cd4}_Prim"] < 200),
                (CD4[f"{col_contagem_cd4}_Prim"] >= 200) & (CD4[f"{col_contagem_cd4}_Prim"] < 350),
                (CD4[f"{col_contagem_cd4}_Prim"] >= 350) & (CD4[f"{col_contagem_cd4}_Prim"] < 500),
                (CD4[f"{col_contagem_cd4}_Prim"] >= 500) & (CD4[f"{col_contagem_cd4}_Prim"] < 2500)
        ]

        Escolha_CD4cat = ["0-199", "200-349", "350-499", "500+"]

        CD4["CD4_cat_Prim"] = np.select(Cond_CD4cat, Escolha_CD4cat, default="Sem resultado de CD4")

        display(CD4["CD4_cat_Prim"].value_counts(dropna = False).sort_index())
        print()

        # criando variável de categorias do valor do CD4 primeiro CD4 menor que 200
        Cond_CD4cat2 = [
                (CD4[f"{col_contagem_cd4}_Prim"] > 0) & (CD4[f"{col_contagem_cd4}_Prim"] < 200),
                (CD4[f"{col_contagem_cd4}_Prim"] >= 200) & (CD4[f"{col_contagem_cd4}_Prim"] < 2500)
        ]

        Escolha_CD4cat2 = ["<200 células / mm³", "≥200 células / mm³"]

        CD4["CD4_cat200_Prim"] = np.select(Cond_CD4cat2, Escolha_CD4cat2, default="Sem resultado de CD4")

        display(CD4["CD4_cat200_Prim"].value_counts(dropna = False).sort_index())
        print()

    return CD4



def organizacao_sim(SIM:pd.DataFrame, SIM_total:pd.DataFrame, Cod_tab:pd.DataFrame, datas_sim:list = ["DTOBITO_sim","DTNASC_sim"], lista_raca:list = ["RACACOR_sim"],
                    cols_filna:list = ["NOME_sim","NOMEMAE_sim","DTNASC_sim", "SEXO_sim","RACACOR_sim","CODMUNRES_sim"], col_codigo_pac_uni: str = "Cod_unificado"):

    Cod_tab["codigo_paciente"] = Cod_tab["codigo_paciente"].astype(int)
    Sim_total["codigo_paciente"] = Sim_total["codigo_paciente"].astype(int)
    Sim_total = pd.merge(Sim_total,Cod_tab, on = "codigo_paciente", how = "left")
    
    SIM_total.rename(columns = {"dt_nasc":"DTNASC", "sexo":"SEXO","ibge_res":"CODMUNOCOR"}, inplace = True)

    linhas_total = SIM.shape[0] + SIM_total.shape[0]

    colunas_total = SIM.shape[1]

    SIM = pd.concat([SIM,SIM_total])

    if (SIM.shape[1] != colunas_total) | (SIM.shape[0] != linhas_total):
        print("Erro ao unir o Sim_AIDS com o Sim_Completo")
        return KeyError

    
    SIM[col_codigo_pac_uni] = SIM[col_codigo_pac_uni].astype("Int64")

    # Identificar os registros com valores em branco no ID_Sim
    Sem_cod_prob = SIM["id_sim"].isnull()

    # Definir o valor inicial para os códigos faltantes
    start_value = 210000000

    # Gerar valores únicos para os registros faltantes,
    SIM.loc[Sem_cod_prob, "id_sim"] = np.arange(start_value, start_value + Sem_cod_prob.sum())

    SIM = SIM.rename(columns={col: f"{col}_sim" for col in SIM.columns[:-2]})

    for col in datas_sim:
        SIM = fg.ajusta_data_linha_vetorizado(SIM, col, coluna_retorno = f"{col}_ajustada")
        SIM[col] = SIM[f"{col}_ajustada"]
        SIM = SIM.drop(columns = [f"{col}_ajustada"])

    for col_raca in lista_raca:

        SIM[col_raca] = SIM[col_raca].apply(lambda x:str(x).strip())
        
        Rac_cond = [((SIM[col_raca] == "1.0") | (SIM[col_raca] == "3.0")),
                    (SIM[col_raca] == "2.0"),
                    (SIM[col_raca] == "4.0"),
                    (SIM[col_raca] == "5.0")
                    ]
                
        Rac_escolha = ["Branca/Amarela","Preta","Parda","Indígena"]
        
        SIM[col_raca] = np.select(Rac_cond,Rac_escolha, default = None)


    SIM = SIM.sort_values([col_codigo_pac_uni, "DTOBITO_sim"], ascending = True)
    SIM[cols_filna] = SIM.groupby(col_codigo_pac_uni)[cols_filna].ffill()
    SIM[cols_filna] = SIM.groupby(col_codigo_pac_uni)[cols_filna].bfill()

    return SIM



def organizacao_sinan(Sinan:pd.DataFrame, tipo:str, col_codigo_pac_uni: str = "Cod_unificado", col_data:str = "DT_NOTIFIC"):
     
    Sinan[col_codigo_pac_uni] = Sinan[col_codigo_pac_uni].astype("Int64")

    Sinan = fg.ajusta_data_linha_vetorizado(Sinan, col_data, coluna_retorno = f"{col_data}_ajustada")
    Sinan[col_data] = Sinan[f"{col_data}_ajustada"]
    Sinan = Sinan.drop(columns = [f"{col_data}_ajustada"])

    
    # Padroniza a Escolaridade
    lista_escol = [
        "CS_ESCOLAR",
        "CS_ESCOL_N"
    ]

    for cols_escol in lista_escol:

        Sinan[cols_escol] = Sinan[cols_escol].apply(lambda x:str(x).strip())
        
        condicoes_escol = [
            ((Sinan[cols_escol] == "1.0") | (Sinan[cols_escol] == "01") |
            (Sinan[cols_escol] == "0 -Analfabeto") | (Sinan[cols_escol] == "1.0") |
            (Sinan[cols_escol] == "2.0") | (Sinan[cols_escol] == "3.0") |
            (Sinan[cols_escol] == "02") | (Sinan[cols_escol] == "03") |
            (Sinan[cols_escol] == "3 -5ª à 8ª série incompleta do EF (antigo ginásio ou 1º grau)") | 
            (Sinan[cols_escol] == "1 -1ª a 4ª série incompleta do EF (antigo primário ou 1º grau)") |
            (Sinan[cols_escol] == "2 -4ª série completa do EF (antigo primário ou 1º grau)")
            ),
            ((Sinan[cols_escol] == "4.0") | (Sinan[cols_escol] == "04") |
            (Sinan[cols_escol] == "4 -Ensino fundamental completo (antigo ginásio ou 1º grau)") |
            (Sinan[cols_escol] == "5 -Ensino médio incompleto (antigo colegial ou 2º grau )")
            ),
            ((Sinan[cols_escol] == "5.0") | (Sinan[cols_escol] == "05") |
            (Sinan[cols_escol] == "6 -Ensino médio completo (antigo colegial ou 2º grau )") |
            (Sinan[cols_escol] == "8 -Educação superior completa") |
            (Sinan[cols_escol] == "7 -Educação superior incompleta") | (Sinan[cols_escol] == "05")
            )
        ]
        
        resultado_escol = [
           1, 2, 3
        ]
        
        Sinan[cols_escol] = np.select(condicoes_escol,resultado_escol, default = None)


    Sinan['Escol_num'] = Sinan[['CS_ESCOLAR','CS_ESCOL_N']].max(axis=1)

    # Preenche as colunas em branco com valores do outros registros
    cols = ["Escol_num"]
    funcoes = ["max"]

    # Primeiro, ordene os dados por 'Cod_unificado' e 'data_ref'
    for var,funcao in zip(cols,funcoes):
        Sinan = fg.padronizar_variaveis_vetorizado(Sinan, var = var, col_data_ref = col_data, funcao = funcao, mais_antigo = False)

    Sinan = Sinan.drop(columns=lista_escol)

    Sinan = Sinan.sort_values([col_codigo_pac_uni, f"DT_NOTIFIC"], ascending = True)
    Sinan = Sinan.rename(columns={col: f"{col}_{tipo}" for col in Sinan.columns if col != col_codigo_pac_uni})

    return Sinan



def ajuste_sinan_var_padrao(Sinan:pd.DataFrame,tipo:str,col_codigo_pac_uni: str = "Cod_unificado"):


    datas_sinan = [
        f"DT_NASC_{tipo}",
        f"DT_NOTIFIC_{tipo}",
        f"DT_DIAG_{tipo}",
        f"DT_CONFIRM_{tipo}",
        f"DT_RAPIDO_{tipo}",
        f"DT_SIN_PRI_{tipo}",
        f"DT_RAPIDO_{tipo}",
        f"DT_LAB_HIV_{tipo}",
        f"DTCONFIRMA_{tipo}",
        f"DTRAPIDO1_{tipo}",
        f"DT_PCR_1_{tipo}",
        f"DT_PCR_2_{tipo}",
        f"DT_PCR_3_{tipo}"
    ]
    
    Sinan = Sinan.rename(columns={col: f"{col}_{tipo}" for col in Sinan.columns[:-2]})
    Sinan[col_codigo_pac_uni] = Sinan[col_codigo_pac_uni].astype("Int64")

    # Ajusta as datas
    for col in datas_sinan:
        if col in Sinan.columns:
            Sinan = fg.ajusta_data_linha_vetorizado(Sinan, col, coluna_retorno = f"{col}_ajustada")
            Sinan[col] = Sinan[f"{col}_ajustada"]
            Sinan = Sinan.drop(columns = [f"{col}_ajustada"])

    Sinan["data_nascimento"] = Sinan[f"DT_NASC_{tipo}"]
    Sinan["data_ref"] = Sinan[f"DT_NOTIFIC_{tipo}"]
    Sinan["nome"] = Sinan[f"NM_PACIENT_{tipo}"]
    Sinan["nome_mae"] = Sinan[f"NM_MAE_PAC_{tipo}"]
    Sinan["ibge6_res"] = Sinan[f"ID_MN_RESI_{tipo}"].apply(lambda x:str(x)[:6])

    # Padroniza a Raca
    lista_raca = [
        f"CS_RACA_{tipo}"
    ]
    for col_raca in lista_raca:
        Sinan[col_raca] = Sinan[col_raca].apply(lambda x:str(x).strip())
        Rac_cond = [((Sinan[col_raca] == "1.0") | (Sinan[col_raca] == "1 -Branca") | (Sinan[col_raca] == "1")),
                    ((Sinan[col_raca] == "3.0") | (Sinan[col_raca] == "3 -Amarela") | (Sinan[col_raca] == "3")),
                    ((Sinan[col_raca] == "2.0") | (Sinan[col_raca] == "2 -Preta") | (Sinan[col_raca] == "2")),
                    ((Sinan[col_raca] == "4.0") | (Sinan[col_raca] == "4 -Parda") | (Sinan[col_raca] == "4")),
                    ((Sinan[col_raca] == "5.0") | (Sinan[col_raca] == "5 -Indigena") | (Sinan[col_raca] == "5"))
                    ]
                
        Rac_escolha = ["Branca","Amarela","Preta","Parda","Indígena"]
        Sinan["Raca_cat"] = np.select(Rac_cond,Rac_escolha, default = None)


    # Ajusta o Sexo
    Sex = [(Sinan[f"CS_SEXO_{tipo}"] == "M"),(Sinan[f"CS_SEXO_{tipo}"] == "F")]
             
    Sex_escolha = ["Homem","Mulher"]
    
    Sinan["Sexo_cat"] = np.select(Sex, Sex_escolha, default=None)
    display(Sinan["Sexo_cat"].value_counts(dropna = False))
    print()
    return Sinan



def organizar_bases(hoje: datetime.date, Pac: pd.DataFrame, Disp: pd.DataFrame, CV: pd.DataFrame, CD4: pd.DataFrame, Cod_uni:pd.DataFrame,
                    obito:dict = {"SIM":None, "obitos_siclom":None}, Usar_para_Prim:bool = False):
    """
    Essa função deve ser utilizada para organizar simultaneamente as bases de 
    Cadastro, Dispensação, Carga Viral e CD4 do HIV, com o intuito de uni-las em uma base única posterioremente.
    Todas as bases são obrigatórias, exceto a base de óbitos que deve ser passada como um dicionário com as chaves {obito_rel, obitos_siclom},
    caso deseje incluir a data do óbito na base de cadastro.

    Essa função unifica as funções:

    organizacao_cadastro()
    organizacao_disp()
    organizacao_cv()
    organizacao_cd4()

    E retorna obrigatoriamente 4 DataFrames na seguinte ordem:

    Pac, Disp, CV e CD4

    
    Parameters
    ----------

    hoje: datetime.date, default = hoje
        Data de fechamento dos bancos de análise no formato pd.to_datetime(aaaa-mm-dd).

    Pac: pd.DataFrame, default = Pac
        DataFrame original com a base de cadastro de pacientes (tb_paciente_consolidado).

    Disp: pd.DataFrame, default = Disp
        DataFrame original com a base de dispensação de ARV (tb_dispensas_esquemas_udm).

    CV: pd.DataFrame, default = CV
        DataFrame original com a base de exames de carga viral (tb_carga_viral_consolidado).

    CD4: pd.DataFrame, default = CD4
        DataFrame original com a base de exames de CD4 (tb_cd4_consolidado).

    obito: Dicionário, opcional, default = None
        Dicionário com os dataframes de óbito, deve ser passada no formato: {obito_rel:df_rel, obitos_siclom:df_siclom}.

    Usar_para_Prim: bool, opcional, default = False
        Define se quer gerar as bases para unir no DF_Prim.
        Se usar para fazer a base DF_ult_ano, manter False.

    """
    
    if obito["SIM"] is None:
        usar_obito = False
    else:
        usar_obito = True       


    if Usar_para_Prim is True:
          
        Pac, Cod_tab = organizacao_cadastro(Pac, Cod_uni, hoje,  usar_obito = usar_obito, SIM = obito["SIM"], obitos_siclom = obito["obitos_siclom"])

        Disp = organizacao_disp(Disp, Cod_tab, hoje, Prim=False)

        CV = organizacao_cv(CV, Cod_tab, hoje, categorizar=True, Prim=False)

        CD4 = organizacao_cd4(CD4, Cod_tab, hoje, categorizar=True, Prim=False)
          
    else:
        Pac, Cod_tab = organizacao_cadastro(Pac, Cod_uni, hoje, usar_obito = usar_obito, SIM = obito["SIM"], obitos_siclom = obito["obitos_siclom"])

        Disp = organizacao_disp(Disp, Cod_tab, hoje, Prim=True, var_prim = ["esquema_AMA"])

        CV = organizacao_cv(CV, Cod_tab, hoje, categorizar=False, Prim=True)

        CD4 = organizacao_cd4(CD4, Cod_tab, hoje, categorizar=False, Prim=True)

    return Pac, Disp, CV, CD4



def bases_ult_ano(hoje: datetime.date, DF_lista: list, DF_nome:list, col_data:list, col_codigo_pac_uni: str = "Cod_unificado"):
    """
    Essa função cria um DataFrame onde cada unidade de análise (por exemplo pessoa) fica apenas com 
    a última linha por ano.
    Para funcionar corretamente, o DataFrame original DEVE ESTAR ORGANIZADO PELA DATA COM ASCENDING = FALSE).
    As funções de organização desse programa já distribui as bases dessa forma.

    Essa função retorna os DataFrames com o sufixo "_ano" na mesma ordem em que foram passados na DF_lista.

    Parameters
    ----------

    hoje: datetime.date, default = hoje
        Data de fechamento dos bancos de análise no formato pd.to_datetime(aaaa-mm-dd).

    ano_corte = int, default = ano_corte
        Ano de referência para iniciar a contagem das análises (utilizar ano - x, para automatizar o ano)

    DF_lista = list, default = None
        Lista de DataFrames que deseja transformar para o formato anual.

    DF_nome = list, default = None
        Lista com o nome dos DataFrames em String, deve seguir a mesma ordem da DF_lista.

    col_data = list, default = None
        Lista com o nome da coluna com a referência de ano de cada DataFrame passado, deve seguir a mesma ordem da DF_lista.

    col_codigo_pac_eleito: str, opcional, default = "codigo_pac_eleito"
        Nome da coluna que contém o código do paciente eleito.

    """
    resultados = {}
    for DF, nome, ano_data in zip(DF_lista, DF_nome, col_data):

        # Criando um DataFrame com todas as combinações de ano e código da pessoa
        ano = hoje.year
        periodo = range(2000, ano + 1)
        codigos = DF[col_codigo_pac_uni].unique()


        # Listas para armazenar os dados
        dados = {"ano": [], col_codigo_pac_uni: []}
        
        # Preenchendo as listas com todas as combinações de ano e código da pessoa
        for ano_analise in periodo:
            for codigo in codigos:
                dados["ano"].append(ano_analise)
                dados[col_codigo_pac_uni].append(codigo)
        
        # Criando o DataFrame a partir das listas
        df = pd.DataFrame(dados)
        
        lista_dfs = []
        for ano_analise in periodo:
            ult_linha = DF[(DF[ano_data] <= ano_analise)].drop_duplicates(subset = [col_codigo_pac_uni], keep = "first")
            df_temp = pd.merge(df[df["ano"] == ano_analise], ult_linha, on = [col_codigo_pac_uni], how = "left")
            lista_dfs.append(df_temp)
                
        DF_ano = pd.concat(lista_dfs)
        
        DF_ano = DF_ano[(DF_ano["ano"] >= DF_ano[ano_data])].copy()

        resultados[f"{nome}_ano"] = DF_ano

    return resultados
    


def bases_ult_meses(hoje: datetime.date,
                      DF_lista: list,
                      DF_nome: list,
                      col_data: list,
                      col_codigo_pac_uni: str = "Cod_unificado",
                      meses:int = 24):
    """
    Cria DataFrames contendo a última linha de cada pessoa por mês
    (último dia do mês) nos últimos X meses.

    Parâmetros
    ----------
    hoje : pd.Timestamp
        Data de referência.
    DF_lista : list
        Lista de DataFrames de entrada.
    DF_nome : list
        Lista com os nomes (string) dos DataFrames, na mesma ordem.
    col_data : list
        Lista com os nomes das colunas de data em cada DataFrame.
    col_codigo_pac_uni : str, default = "Cod_unificado"
        Coluna identificadora única do paciente.
    meses : int, default = 24
        Quantidade de meses retroativos a considerar.

    Retorna
    -------
    dict : {nome_base_mes: DataFrame}
    """

    resultados = {}

    # Gera lista de datas (último dia de cada mês)
    data_corte = hoje - pd.DateOffset(months=meses)
    meses_analise = pd.date_range(start=data_corte, end=hoje, freq='M')

    for DF, nome, col_dt in zip(DF_lista, DF_nome, col_data):
        # Garante que a coluna de data é datetime
        DF[col_dt] = pd.to_datetime(DF[col_dt], errors='coerce')

        # Ordena por data decrescente (importante para drop_duplicates funcionar direito)
        DF = DF.sort_values(by=col_dt, ascending=False).copy()

        lista_dfs = []

        for data_ref in meses_analise:
            ult_linha = DF[DF[col_dt] <= data_ref].drop_duplicates(
                subset=[col_codigo_pac_uni],
                keep='first'
            )
            ult_linha["mes_ano"] = data_ref.strftime("%m-%Y")
            lista_dfs.append(ult_linha)

        DF_meses = pd.concat(lista_dfs, ignore_index=True)
        resultados[f"{nome}_mes"] = DF_meses

    return resultados



def bases_prim(DF_lista: list, DF_nome:list, col_data:list, col_codigo_pac_uni: str = "Cod_unificado"):
    """
    Filtra e organiza as bases de dados para obter as primeiras ocorrências de cada paciente a partir de um ano de corte especificado.

    Retorna um dicionário com DataFrames das bases de dados filtradas, cada uma contendo apenas a primeira ocorrência de cada paciente a partir do ano de corte.

    Parameters
    ----------
    ano_corte: int
        Ano de corte para filtrar os dados. Apenas registros a partir deste ano serão considerados.

    DF_lista: list
        Lista de DataFrames a serem processados.

    DF_nome: list
        Lista de nomes das bases de dados correspondentes aos DataFrames em DF_lista. Usado como chave no dicionário de resultados.

    col_data: list
        Lista de nomes das colunas de data em cada DataFrame, usadas para ordenar e filtrar os registros.

    col_codigo_pac_eleito: str, opcional, default = "codigo_pac_eleito"
        Nome da coluna que contém o código do paciente eleito. Usada para identificar registros únicos de pacientes.

    Returns
    -------
    dict
        Um dicionário com DataFrames filtrados, onde cada chave é o nome da base de dados acrescido de "_Prim" e o valor é o DataFrame correspondente com as primeiras ocorrências a partir do ano de corte.

    Example
    -------
    >>> DF_lista = [Disp, CV, CD4]
    >>> DF_nome = ["Disp", "CV", "CD4"]
    >>> col_data = ["data_dispensa", "data_hora_coleta_cv", "data_hora_coleta_cd4"]
    >>> Disp_Prim, CV_Prim, CD4_Prim = bases_prim(ano_corte, DF_lista, DF_nome, col_data).values()

    """

    resultados = {}
    for DF, nome, data in zip(DF_lista, DF_nome, col_data):
        DF.sort_values([data], ascending = [True], inplace = True)
        DF.reset_index(drop = True, inplace = True)
        df_prim = DF.drop_duplicates([col_codigo_pac_uni], keep = "first")
        display(df_prim[data].dt.year.value_counts().sort_index())
        resultados[f"{nome}_Prim"] = df_prim

    return resultados



def df_ult_ano(hoje: datetime.date, ano_corte:int, Pac: pd.DataFrame, Disp_ano: pd.DataFrame, CV_ano: pd.DataFrame, CD4_ano: pd.DataFrame, col_codigo_pac_uni: str = "Cod_unificado",
                col_ano:str = "ano", col_ano_disp:str = "ano_disp",  col_ano_cv:str = "ano_coleta_cv",  col_ano_cd4:str = "ano_coleta_cd4",
                col_dt_obito:str = "data_obito_total", cols_datas:list = ["data_dispensa","data_hora_coleta_cv", "data_hora_coleta_cd4"],
                col_data_nascimento:str = "data_nascimento", cols_data_min:list = ["data_coleta_cd4_Prim", "data_coleta_cv_Prim","data_sol_cd4_Prim","data_sol_cv_Prim","data_dispensa_Prim"],
                col_ibge_resid:str = "codigo_ibge_resid", cols_ibge:list = ["cod_ibge_udm","cod_ibge_solicitante_cv","cod_ibge_solicitante_cd4"],
                col_cod_ibge_udm:str = "cod_ibge_udm", col_cod_ibge_solicitante_cv:str = "cod_ibge_solicitante_cv",
                col_cod_ibge_solicitante_cd4 = "cod_ibge_solicitante_cd4", col_nome_inst_sol_cv:str = "nm_inst_sol_cv",col_nome_inst_sol_cd4:str = "nm_inst_sol_cd4",
                col_copias: str = "copias", col_comentario_copias: str = "comentario_copias", col_contagem_cd4:str = "contagem_cd4",
                obito_admin:bool = True, anos_obito_adm:int = 10):
    
    """
    A função df_ult_ano realiza a integração de diversas fontes de dados relacionadas a pacientes, dispensações, carga viral e contagem de CD4,
    e retorna um DataFrame onde cada pessoa possui uma linha por ano, com as informações do último evento daquele ano.
    Em colunas identificadas com o sufixo _Prim, constam informações sobre o primeiro evento da vida, e para utilizá-las, deve-se filtrar a base para [Base_Prim] == 1.
    Diversas categorias e outras métricas que envolvam mais de um DataFrame também são calculadas.

    
    Parameters
    ----------
    hoje: datetime.date, default = hoje
        Data de fechamento dos bancos de análise no formato pd.to_datetime(aaaa-mm-dd).

    Pac: pd.DataFrame, default = Pac
        DataFrame original com a base de cadastro de pacientes (tb_paciente_consolidado).

    Disp_ano: pd.DataFrame, default = Disp_ano
        DataFrame original com a base de dispensação de ARV (tb_dispensas_esquemas_udm) por ano.

    CV_ano: pd.DataFrame, default = CV_ano
        DataFrame original com a base de exames de carga viral (tb_carga_viral_consolidado) por ano.

    CD4_ano: pd.DataFrame, default = CD4_ano
        DataFrame original com a base de exames de CD4 (tb_cd4_consolidado) por ano.

    col_codigo_pac_eleito: str, default = "codigo_pac_eleito"
        Nome da coluna que contém o código do paciente.

    col_ano: str, default = "ano"
        Nome da coluna que contém o ano.

    col_ano_disp: str, default = "ano_disp"
        Nome da coluna que contém o ano de dispensação.

    col_ano_cv: str, default = "ano_coleta_cv"
        Nome da coluna que contém o ano de coleta de carga viral.

    col_ano_cd4: str, default = "ano_coleta_cd4"
        Nome da coluna que contém o ano de coleta de CD4.

    col_dt_obito: str, default = "data_obito_total"
        Nome da coluna que contém a data de óbito.

    cols_datas: list, default = ["data_dispensa", "data_hora_coleta_cv", "data_hora_coleta_cd4"]
        Lista de colunas de datas de cada banco.

    col_data_nascimento: str, default = "data_nascimento"
        Nome da coluna que contém a data de nascimento.

    cols_data_min: list, default = ["data_coleta_cd4_Prim", "data_coleta_cv_Prim", "data_sol_cd4_Prim", "data_sol_cv_Prim", "data_dispensa_Prim"]
        Lista de colunas para cálculo da data mínima.

    col_ibge_resid: str, default = "codigo_ibge_resid"
        Nome da coluna que contém o código IBGE de residência.

    cols_ibge: list, default = ["cod_ibge_udm", "cod_ibge_solicitante_cv", "cod_ibge_solicitante_cd4"]
        Lista de colunas IBGE para complementar o IBGE de residência.

    col_contagem_cd4: str, default = "contagem_cd4"
        Nome da coluna que contém a contagem de cd4.

    obito_admin: bool, default = True
        Indica se óbitos presumidos (pessoas há mais de 10 anos sem aparecer na base) devem ser excluídos.

        
        
    Returns
    -------
    pd.DataFrame
        DataFrame consolidado com informações detalhadas sobre tratamentos e exames dos pacientes, onde cada pessoa possui uma linha por ano, com as informações do último evento daquele ano.
    

    """
    ano = hoje.year

    Disp_ano = Disp_ano.drop(columns=["codigo_paciente"])
    CD4_ano = CD4_ano.drop(columns=["codigo_paciente"])
    CV_ano = CV_ano.drop(columns=["codigo_paciente"])   
    
    DF_ult_ano = pd.merge(Disp_ano,CD4_ano, on = [col_ano,col_codigo_pac_uni], how = "outer")
    DF_ult_ano = pd.merge(DF_ult_ano,CV_ano, on = [col_ano,col_codigo_pac_uni], how = "outer")
    DF_ult_ano = pd.merge(DF_ult_ano,Pac, on = [col_codigo_pac_uni], how = "left")

    # Preenche as colunas do ano que estavam vazias.
    DF_ult_ano["ano_obito"]= DF_ult_ano[col_dt_obito].dt.year
    DF_ult_ano[col_ano_disp] = DF_ult_ano[col_ano_disp].fillna(np.nan)
    DF_ult_ano[col_ano_disp] = DF_ult_ano[col_ano_disp].astype("Int64")
    DF_ult_ano[col_ano_cv] = DF_ult_ano[col_ano_cv].fillna(np.nan)
    DF_ult_ano[col_ano_cv] = DF_ult_ano[col_ano_cv].astype("Int64")
    DF_ult_ano[col_ano_cd4] = DF_ult_ano[col_ano_cd4].fillna(np.nan)
    DF_ult_ano[col_ano_cd4] = DF_ult_ano[col_ano_cd4].astype("Int64")


    # Calcula a idade em todas as datas passadas no cols_idade.
    for cols in cols_datas:
        
        DF_ult_ano[f"Idade_{cols}_cat"] = fg.idade_cat(DF = DF_ult_ano, data_ref=cols, data_nasc=col_data_nascimento)
        display(DF_ult_ano[f"Idade_{cols}_cat"].value_counts().sort_index())
        print()

    # Calcula a data mínima de realização de exames (1o da vida)
    colunas_data_min_exame = [col for col in cols_data_min if "disp" not in col]

    DF_ult_ano = fg.encontrar_menor_data_vetorizado(df=DF_ult_ano, colunas_datas=colunas_data_min_exame, nome_dt_min = "data_min_exame")
    DF_ult_ano["data_min_exame"] = pd.to_datetime(DF_ult_ano["data_min_exame"],errors="coerce").dt.normalize()
    DF_ult_ano["ano_min_exame"]= DF_ult_ano["data_min_exame"].dt.year
    display(DF_ult_ano["ano_min_exame"].value_counts(dropna = False).sort_index())
    print()

    # Calcula a data mínima de vinculação (1o da vida)
    DF_ult_ano = fg.encontrar_menor_data_vetorizado(df=DF_ult_ano, colunas_datas=cols_data_min)
    DF_ult_ano["data_min"] = pd.to_datetime(DF_ult_ano["data_min"],errors="coerce").dt.normalize()
    DF_ult_ano["ano_min"]= DF_ult_ano["data_min"].dt.year
    display(DF_ult_ano["ano_min"].value_counts(dropna = False).sort_index())
    print()

    # Calcula a idade de vinculação
    DF_ult_ano["Idade_vinc_cat"] = fg.idade_cat(DF = DF_ult_ano, data_ref="data_min", data_nasc=col_data_nascimento)
    display(DF_ult_ano["Idade_vinc_cat"].value_counts().sort_index())
    print()

    # Calcula a idade de vinculação em faixas etárias antigas
    DF_ult_ano["Idade_vinc_cat_anterior"] = fg.idade_cat(DF = DF_ult_ano, data_ref="data_min", data_nasc=col_data_nascimento, faixas_etarias="MC_antigo")
    display(DF_ult_ano["Idade_vinc_cat_anterior"].value_counts().sort_index())
    print()

    # Crie uma série de datas com o último dia do ano
    DF_ult_ano["Ult_dia_ano"] = pd.to_datetime(DF_ult_ano[col_ano].astype(str) + "-12-31")

    # Ajusta o último dia do ano para o dia de fechamento da base no ano atual
    DF_ult_ano.loc[DF_ult_ano[col_ano] == ano, "Ult_dia_ano"] = hoje

    # Normalize as datas para garantir que apenas a data seja mantida (sem informações de hora)
    DF_ult_ano["Ult_dia_ano"] = DF_ult_ano["Ult_dia_ano"].dt.normalize()

    # Calcula a idade no último dia de cada ano (importante para Cascata e BI)
    DF_ult_ano["Idade_cascata"] = fg.idade_cat(DF = DF_ult_ano, data_ref="Ult_dia_ano", data_nasc=col_data_nascimento)
    display(DF_ult_ano["Idade_cascata"].value_counts().sort_index())
    print()

    # Calcula a idade no último dia de cada ano nas categorias do spectrum
    DF_ult_ano["Idade_Spec_crianca"] = fg.idade_cat(DF = DF_ult_ano, data_ref="Ult_dia_ano", data_nasc=col_data_nascimento, faixas_etarias="spectrum_crianca")
    display(DF_ult_ano["Idade_Spec_crianca"].value_counts().sort_index())
    print()   
    DF_ult_ano["Idade_Spec1"] = fg.idade_cat(DF = DF_ult_ano, data_ref="Ult_dia_ano", data_nasc=col_data_nascimento, faixas_etarias="spectrum1")
    display(DF_ult_ano["Idade_Spec1"].value_counts().sort_index())
    print()   
    DF_ult_ano["Idade_Spec2"] = fg.idade_cat(DF = DF_ult_ano, data_ref="Ult_dia_ano", data_nasc=col_data_nascimento, faixas_etarias="spectrum2")
    display(DF_ult_ano["Idade_Spec2"].value_counts().sort_index())
    print()   

   
    data_disp = next((col for col in cols_datas if "disp" in col.lower()), None)

    data_cv = next((col for col in cols_datas if "cv" in col.lower() or "carga" in col.lower()), None)

    data_cd4 = next((col for col in cols_datas if "cd4" in col.lower()), None)

    data_disp_prim = next((col for col in cols_data_min if "disp" in col.lower()), None)


    # Filtra excluindo: crianças menores de 2 anos sem TARV, linhas de anos anteriores à vinculação e linhas de anos a partir do ano do óbito.
    DF_ult_ano = DF_ult_ano[((DF_ult_ano["Idade_vinc_cat_anterior"] != "Menos de 2 anos") | (DF_ult_ano[data_disp_prim].isnull() == False) | (DF_ult_ano["CV_5000_count"] >= 2)) &
                            (DF_ult_ano[col_ano] >= DF_ult_ano["ano_min"]) &
                            ((DF_ult_ano["ano_obito"].isnull() == True) |
                            (DF_ult_ano[col_ano] < DF_ult_ano["ano_obito"]))].copy()
        
    # Organiza o IBGE de residência
    DF_ult_ano = fg.ibge_resid(DF=DF_ult_ano, col_ibge_resid=col_ibge_resid, cols_ibge=cols_ibge)

    # Organiza o IBGE das instituições Solicitantes
    DF_ult_ano = fg.ibge_inst_sol_exames(DF=DF_ult_ano, col_cod_ibge_udm = col_cod_ibge_udm, col_cod_ibge_solicitante_cv = col_cod_ibge_solicitante_cv,
                    col_cod_ibge_solicitante_cd4 = col_cod_ibge_solicitante_cd4, col_nome_inst_sol_cv = col_nome_inst_sol_cv,
                    col_nome_inst_sol_cd4 = col_nome_inst_sol_cd4, col_ibge_resid = col_ibge_resid)

    ## A Coluna "Status_ano" mostra a situação do paciente em relação ao último dia do ano em que a dispensação foi feita, ou em relação a data atual
    # Calcula o atraso em relação ao último dia do ano da dispensação
    DF_ult_ano["Atraso_ano"] = (DF_ult_ano["Ult_dia_ano"] - DF_ult_ano["dt_lim_prox_disp"]).dt.days

    cond_status = [
        (DF_ult_ano["Atraso_ano"] <= 60),
        (DF_ult_ano["Atraso_ano"] > 60)
        ]

    escolha_status = ["Tarv", "Interrupção de Tarv"]

    DF_ult_ano["Status_ano"] = np.select(cond_status, escolha_status, default= None)

    DF_ult_ano["Status_ano"] = DF_ult_ano["Status_ano"].fillna("Gap de tratamento")

    display("Status no último dia do ano, para o ano vigente, considera-se a data de fechamento do banco",
            pd.DataFrame(DF_ult_ano[[col_ano,"Status_ano"]].value_counts(dropna = False).sort_index(ascending = False)))
    print()

    # Cria uma variável com o vínculo no ano
    Cond_vin = [
        (DF_ult_ano[col_ano] == ano) &
        (((DF_ult_ano["Hoje"] - DF_ult_ano[data_cv]).dt.days <= 365) |
        ((DF_ult_ano["Hoje"] - DF_ult_ano[data_cd4]).dt.days <= 365) |
        ((DF_ult_ano["Hoje"] - DF_ult_ano[data_disp]).dt.days <= 365)),
        
        ((DF_ult_ano[col_ano] == DF_ult_ano[data_cv].dt.year) |
        (DF_ult_ano[col_ano] == DF_ult_ano[data_cd4].dt.year) |
        (DF_ult_ano[col_ano] == DF_ult_ano[data_disp].dt.year)),
    ]

    escolha_vin = [
        "Vinculado", "Vinculado"
    ]

    DF_ult_ano["Vinculado_ano"]  = np.select(Cond_vin,escolha_vin, default = "Não vinculado")
    display("Vinculados por ano, considera os últimos 12 meses para o ano vigente",
            pd.DataFrame(DF_ult_ano[["ano","Vinculado_ano"]].value_counts(dropna = False).sort_index(ascending = False)))
    print()

    data_disp_prim = next((col for col in cols_data_min if "disp" in col.lower()), None)

    # Conta o número de pessoas que fazem acompanhamento (exame CV) em cada instituição.
    DF_ult_ano["N_PVHA_CV_ano_Instituicao"] = DF_ult_ano[DF_ult_ano["nome_inst_exames"] != "SEM EXAMES NO SUS"].groupby(["ano","nome_inst_exames"])[col_codigo_pac_uni].transform("count")

    cond_inst = [
            (DF_ult_ano["N_PVHA_CV_ano_Instituicao"] >= 1000),
            (DF_ult_ano["N_PVHA_CV_ano_Instituicao"] >= 500),
            (DF_ult_ano["N_PVHA_CV_ano_Instituicao"] >= 1)
    ]

    escolhas_inst = [
            "Instituições acompanhando mais de 1000 pessoas",
            "Instituições acompanhando entre 500 e 1000 pessoas",
            "Instituições acompanhando menos de 500 pessoas"
    ]

    DF_ult_ano["Grupo_instituicao"] = np.select(cond_inst,escolhas_inst, default = "Sem pedidos de exames")
    display(DF_ult_ano["Grupo_instituicao"].value_counts())
    print()


    # Categoriza as faixas de Carga Viral
    DF_ult_ano[col_copias] = pd.to_numeric(DF_ult_ano[col_copias], errors='coerce')
    DF_ult_ano[col_comentario_copias] = pd.to_numeric(DF_ult_ano[col_comentario_copias], errors='coerce')
    
    cond_cv = [
            (((DF_ult_ano[col_ano_cv] == DF_ult_ano[col_ano]) |
            ((DF_ult_ano[col_ano] == ano) & ((pd.to_datetime(hoje) - DF_ult_ano[data_cv]).dt.days <= 365))) &
            ((DF_ult_ano[col_copias].isnull() == True) & (DF_ult_ano[col_comentario_copias] == 0)) |
            ((DF_ult_ano[col_copias].isnull() == True) & (DF_ult_ano[col_comentario_copias] == 1)) |
            (DF_ult_ano[col_copias] < 50)
            ),
            (
            ((DF_ult_ano[col_ano_cv] == DF_ult_ano[col_ano]) |
            ((DF_ult_ano[col_ano] == ano) & ((pd.to_datetime(hoje) - DF_ult_ano[data_cv]).dt.days <= 365))) &
            (((DF_ult_ano[col_copias].isnull() == True) & (DF_ult_ano[col_comentario_copias] == 2)) |
            (DF_ult_ano[col_copias] >= 10000))
            ),
            (
            ((DF_ult_ano[col_ano_cv] == DF_ult_ano[col_ano]) |
            ((DF_ult_ano[col_ano] == ano) & ((pd.to_datetime(hoje) - DF_ult_ano[data_cv]).dt.days <= 365))) &
            (DF_ult_ano[col_copias] >= 50) & (DF_ult_ano[col_copias] < 200)
            ),
            (
            ((DF_ult_ano[col_ano_cv] == DF_ult_ano[col_ano]) |
            ((DF_ult_ano[col_ano] == ano) & ((pd.to_datetime(hoje) - DF_ult_ano[data_cv]).dt.days <= 365))) &
            (DF_ult_ano[col_copias] >= 200) & (DF_ult_ano[col_copias] < 1000)
            ),
            (
            ((DF_ult_ano[col_ano_cv] == DF_ult_ano[col_ano]) |
            ((DF_ult_ano[col_ano] == ano) & ((pd.to_datetime(hoje) - DF_ult_ano[data_cv]).dt.days <= 365))) &        
            (DF_ult_ano[col_copias] >= 1000) & (DF_ult_ano[col_copias] < 10000)
            )
    ]
    
    escolhas_cv = [
            "<50","10.000+","50-199","200-999","1.000-9.999"
    ]
    
    DF_ult_ano["CV_cat"] = np.select(cond_cv,escolhas_cv, default = "Não disponível")
    display(DF_ult_ano["CV_cat"].value_counts())
    print()

    # Cria a categoria com 1000 copias de CV   
    cond_cv = [
            (DF_ult_ano["CV_cat"] == "1.000-9.999") | (DF_ult_ano["CV_cat"] == "10.000+"),
            (DF_ult_ano["CV_cat"] == "<50") | (DF_ult_ano["CV_cat"] == "50-199") | (DF_ult_ano["CV_cat"] == "200-999")   
    ]
    
    escolhas_cv = [
            ">=1.000",
            "<1.000"
    ]
    
    DF_ult_ano["CV_cat1000"] = np.select(cond_cv,escolhas_cv, default = "Não disponível")
    display(DF_ult_ano["CV_cat1000"].value_counts())
    print()

    # Cria a categoria com 200 copias de CV
    cond_cv = [
            (DF_ult_ano["CV_cat"] == "<50") | (DF_ult_ano["CV_cat"] == "50-199"),
            (DF_ult_ano["CV_cat"] == "1.000-9.999") | (DF_ult_ano["CV_cat"] == "10.000+") |
            (DF_ult_ano["CV_cat"] == "200-999")   
    ]
    
    escolhas_cv = [
            "<200",
            ">=200"
    ]
    
    DF_ult_ano["CV_cat200"] = np.select(cond_cv,escolhas_cv, default = "Não disponível")
    display(DF_ult_ano["CV_cat200"].value_counts())
    print()

    # Cria a categoria com 50 copias de CV   
    cond_cv = [
            (DF_ult_ano["CV_cat"] == "<50"),
            (DF_ult_ano["CV_cat"] == "1.000-9.999") | (DF_ult_ano["CV_cat"] == "10.000+") |
            (DF_ult_ano["CV_cat"] == "50-199") | (DF_ult_ano["CV_cat"] == "200-999")   
    ]
    
    escolhas_cv = [
            "<50",
            ">=50"
    ]
    
    DF_ult_ano["CV_cat50"] = np.select(cond_cv,escolhas_cv, default = "Não disponível")
    display(DF_ult_ano["CV_cat50"].value_counts())
    print()


    # Cria uma variável com a indetecção viral em 50 cópias
    Cond_sup = [
        (
        (DF_ult_ano["Status_ano"] == "Tarv") & 
        ((DF_ult_ano[data_cv] - DF_ult_ano[data_disp_prim]).dt.days > 180) &
        (DF_ult_ano["CV_cat50"] == "<50")
        ),
        (
        (DF_ult_ano["Status_ano"] == "Tarv") & 
        ((DF_ult_ano[data_cv] - DF_ult_ano[data_disp_prim]).dt.days > 180) &
        (DF_ult_ano["CV_cat50"] == ">=50")
        )
    ]

    escolha_sup = [
       "CV indetectável","CV detectável"
    ]

    DF_ult_ano["Deteccao_CV50"]  = np.select(Cond_sup,escolha_sup, default = "Sem dados de CV")
    display("Indetecção Viral pós TARV na última CV do ano (<50 cp)",
            DF_ult_ano[[col_ano,"Deteccao_CV50"]].value_counts(dropna = False).sort_index())
    

    # Cria uma variável com a supressão viral em 200 cópias
    Cond_sup = [
        (
        (DF_ult_ano["Status_ano"] == "Tarv") & 
        ((DF_ult_ano[data_cv] - DF_ult_ano[data_disp_prim]).dt.days > 180) &
        (DF_ult_ano["CV_cat200"] == "<200")
        ),
        (
        (DF_ult_ano["Status_ano"] == "Tarv") & 
        ((DF_ult_ano[data_cv] - DF_ult_ano[data_disp_prim]).dt.days > 180) &
        (DF_ult_ano["CV_cat200"] == ">=200")
        )
    ]

    escolha_sup = [
        "CV suprimida","CV não suprimida"
    ]

    DF_ult_ano["Deteccao_CV200"]  = np.select(Cond_sup,escolha_sup, default = "Sem dados de CV")
    display("Supressão Viral pós TARV na última CV do ano (<200 cp)",
            DF_ult_ano[[col_ano,"Deteccao_CV200"]].value_counts(dropna = False).sort_index())
    

    # Cria uma variável com a supressão viral em 1000 cópias
    Cond_sup = [
        (
        (DF_ult_ano["Status_ano"] == "Tarv") & 
        ((DF_ult_ano[data_cv] - DF_ult_ano[data_disp_prim]).dt.days > 180) &
        (DF_ult_ano["CV_cat1000"] == "<1.000")
        ),
        (
        (DF_ult_ano["Status_ano"] == "Tarv") & 
        ((DF_ult_ano[data_cv] - DF_ult_ano[data_disp_prim]).dt.days > 180) &
        (DF_ult_ano["CV_cat1000"] == ">=1.000")
        )
    ]

    escolha_sup = [
        "CV suprimida","CV não suprimida"
    ]

    DF_ult_ano["Deteccao_CV1000"]  = np.select(Cond_sup,escolha_sup, default = "Sem dados de CV")
    display("Supressão Viral pós TARV na última CV do ano (<1000 cp)",
            DF_ult_ano[[col_ano,"Deteccao_CV1000"]].value_counts(dropna = False).sort_index())


    DF_ult_ano[col_contagem_cd4] = pd.to_numeric(DF_ult_ano[col_contagem_cd4], errors="coerce")

    # criando variável de categorias do valor do CD4
    Cond_CD4cat = [
        (DF_ult_ano[col_contagem_cd4] > 0) & (DF_ult_ano[col_contagem_cd4] < 200),
        (DF_ult_ano[col_contagem_cd4] >= 200) & (DF_ult_ano[col_contagem_cd4] < 350),
        (DF_ult_ano[col_contagem_cd4] >= 350) & (DF_ult_ano[col_contagem_cd4] < 500),
        (DF_ult_ano[col_contagem_cd4] >= 500) & (DF_ult_ano[col_contagem_cd4] < 2500)
        ]

    Escolha_CD4cat = ["0-199", "200-349", "350-499", "500+"]

    DF_ult_ano["CD4_cat"] = np.select(Cond_CD4cat, Escolha_CD4cat, default="Sem resultado de CD4")

    display(DF_ult_ano["CD4_cat"].value_counts(dropna = False).sort_index())
    print()

    # criando variável de categorias do valor do CD4
    Cond_CD4cat2 = [
        (DF_ult_ano[col_contagem_cd4] > 0) & (DF_ult_ano[col_contagem_cd4] < 200),
        (DF_ult_ano[col_contagem_cd4] >= 200) & (DF_ult_ano[col_contagem_cd4] < 2500)
        ]

    Escolha_CD4cat2 = ["<200 células / mm³", "≥200 células / mm³"]

    DF_ult_ano["CD4_cat200"] = np.select(Cond_CD4cat2, Escolha_CD4cat2, default="Sem resultado de CD4")

    display(DF_ult_ano["CD4_cat200"].value_counts(dropna = False).sort_index())
    print()

    # Calcular o último ano de registro de evento de cada pessoa
    ultimo_evento = DF_ult_ano.groupby(col_codigo_pac_uni).agg({
        col_ano_disp: 'max',
        col_ano_cv: 'max',
        col_ano_cd4: 'max'
    }).max(axis=1).reset_index()
    ultimo_evento.columns = [col_codigo_pac_uni, 'ultimo_evento']

    # Mesclar essa informação de último evento com o DataFrame original
    DF_ult_ano = DF_ult_ano.merge(ultimo_evento, on=col_codigo_pac_uni)

    # Retira os óbitos presumidos (pessoas há mais de 10 anos sem aparecer na base)
    if obito_admin is True:

        # Filtrar os dados onde a diferença entre o ano atual e o último evento é menor ou igual a 10 anos
        DF_ult_ano = DF_ult_ano[DF_ult_ano[col_ano] - DF_ult_ano['ultimo_evento'] <= anos_obito_adm]

    # Adiciona nova coluna para os vinculados nos últimos 5 anos
    DF_ult_ano['Vinculado_5anos'] = np.where(DF_ult_ano["ano"] - DF_ult_ano['ultimo_evento'] <= 5, "Vinculado", "Não Vinculado")

    # Adiciona nova coluna para os vinculados nos últimos 5 anos
    DF_ult_ano['Vinculado_10anos'] = np.where(DF_ult_ano["ano"] - DF_ult_ano['ultimo_evento'] <= 10, "Vinculado", "Não Vinculado")
    
    # Filtra o ano de corte e as pessoas em GAP de tratamento sem CV detectável ou exame de CD4 realizado
    DF_ult_ano = DF_ult_ano[(DF_ult_ano[col_ano] >= ano_corte) & (DF_ult_ano[col_codigo_pac_uni] < 900000000) &
                          ((DF_ult_ano["Status_ano"] != "Gap de tratamento") | (DF_ult_ano["CV_detec_count"] >= 1) | (DF_ult_ano[data_cd4].isnull() == False))]


    # Classifica os Municípios com menos de 50 pessoas por ano em uma categoria unificada
    condicao_mun = (DF_ult_ano.groupby(by = ["ibge_inst_exames", "ano"])[col_codigo_pac_uni].transform("count") < 50).fillna(False)

    DF_ult_ano["Nome_mun_exame_50+"] = np.where(condicao_mun,
                                            "Município com menos de 50 pessoas no cuidado",
                                            DF_ult_ano["Nome_mun_exame"])
    
    DF_ult_ano["Cod_ibge_insti_sol_50+"] = np.where(condicao_mun,
                                        (DF_ult_ano["ibge_inst_exames"].apply(lambda x: str(x)[:2] + "99999")),
                                        DF_ult_ano["ibge_inst_exames"])

    DF_ult_ano["Cod_ibge_insti_sol_50+"] = pd.to_numeric(DF_ult_ano["Cod_ibge_insti_sol_50+"], errors = "coerce").astype("Int64")

    # Classifica os Serviços com menos de 50 pessoas por ano em uma categoria unificada
    condicao_inst = (DF_ult_ano.groupby(by = ["nome_inst_exames","ibge_inst_exames", "ano"])[col_codigo_pac_uni].transform("count") < 50).fillna(False)

    DF_ult_ano["nome_inst_exames_50+"] = np.where(condicao_inst,
                                            "Serviços com menos de 50 pessoas no cuidado",
                                            DF_ult_ano["nome_inst_exames"])


    return DF_ult_ano



def df_prim(Pac: pd.DataFrame, Disp_Prim: pd.DataFrame, CV_Prim: pd.DataFrame, CD4_Prim: pd.DataFrame,
            cols_ano:list = ["ano_disp","ano_coleta_cv","ano_coleta_cd4"],
            col_codigo_pac_uni: str = "Cod_unificado",
            col_data_disp_Prim:str = "data_dispensa",
            col_data_cv_Prim:str = "data_hora_coleta_cv",
            col_data_cd4_Prim:str = "data_hora_coleta_cd4",
            col_data_nascimento:str = "data_nascimento",
            cols_data_min:list = ["data_hora_coleta_cd4", "data_hora_coleta_cv","data_solicitacao_cd4","data_solicitacao_cv","data_dispensa"],
            col_ibge_resid:str = "codigo_ibge_resid",
            cols_ibge:list = ["cod_ibge_udm","cod_ibge_solicitante_cv","cod_ibge_solicitante_cd4"],
            col_cod_ibge_udm:str = "cod_ibge_udm",
            col_cod_ibge_solicitante_cv:str = "cod_ibge_solicitante_cv",
            col_cod_ibge_solicitante_cd4 = "cod_ibge_solicitante_cd4",
            col_nome_inst_sol_cv:str = "nm_inst_sol_cv",
            col_nome_inst_sol_cd4:str = "nm_inst_sol_cd4",
            col_dt_obito:str = "data_obito_total"):
    
    """
    FUNÇÃO EM DESUSO - UTILIZAR O GERAR PVHA_PRIM

    Combina, organiza e filtra os dados das bases de dados do SICLOM/SISCEL sobre dispensações de TARV, exames de carga viral e CD4,
    retornando um DataFrame único que possui apenas uma linha por pessoa com as informações sobre o primeiro evento da vida da pessoa em cada base.

    Parameters
    ----------
    Pac: pd.DataFrame
        DataFrame contendo informações de cadastro no SICLOM/SISCEL.

    Disp_Prim: pd.DataFrame
        DataFrame contendo informações das primeiras dispensações de TARV.

    CV_Prim: pd.DataFrame
        DataFrame contendo informações do primeiro exame de carga viral do HIV realizado.

    CD4_Prim: pd.DataFrame
        DataFrame contendo informações do primeiro exame de contagem CD4.

    cols_ano: list, opcional
        Lista com os nomes das colunas que representam os anos da primeira Dispensação de TARV, exame de Carga Viral e exame de CD4. Default é ["ano_disp", "ano_coleta_cv", "ano_coleta_cd4"].

    col_codigo_pac_eleito: str, opcional
        Nome da coluna que contém o código do paciente. Default é "codigo_pac_eleito".

    col_data_disp_Prim: str, opcional
        Nome da coluna que contém a data da primeira dispensação da vida. Default é "data_dispensa".

    col_data_nascimento: str, opcional
        Nome da coluna que contém a data de nascimento do paciente. Default é "data_nascimento".

    cols_data_min: list, opcional
        Lista com os nomes das colunas que contêm datas para encontrar a menor data. Default é ["data_hora_coleta_cd4", "data_hora_coleta_cv", "data_solicitacao_cd4", "data_solicitacao_cv", "data_dispensa"].

    col_ibge_resid: str, opcional
        Nome da coluna que contém o código IBGE de residência. Default é "codigo_ibge_resid".

    cols_ibge: list, opcional
        Lista com os nomes das colunas que contêm os códigos IBGE adicionais. Default é ["cod_ibge_udm", "cod_ibge_solicitante_cv", "cod_ibge_solicitante_cd4"].

    Returns
    -------
    pd.DataFrame
        Um DataFrame contendo os dados organizados e filtrados com apenas uma linha por pessoa, e os primeiros da vida.

    Example
    -------
    >>> Combinar com bases_prim() para obter os DataFrames de input.
    >>> DF_Prim = df_prim(Pac, Disp_Prim, CV_Prim, CD4_Prim)
    """
    
    Disp_Prim.drop(columns=["codigo_paciente"], inplace=True)
    CD4_Prim.drop(columns=["codigo_paciente"], inplace=True)
    CV_Prim.drop(columns=["codigo_paciente"], inplace=True)

    # Merge dos DataFrames
    DF_Prim = pd.merge(Disp_Prim, CD4_Prim, on=col_codigo_pac_uni, how="outer")
    DF_Prim = pd.merge(DF_Prim, CV_Prim, on=col_codigo_pac_uni, how="outer")
    DF_Prim = pd.merge(DF_Prim, Pac, on=col_codigo_pac_uni, how="left")

    # Cria o ano do óbito
    DF_Prim["ano_obito"]= DF_Prim[col_dt_obito].dt.year

    # Preenchimento de valores nulos e conversão de tipos
    for coluna_ano in cols_ano:
        DF_Prim[coluna_ano] = DF_Prim[coluna_ano].fillna(np.nan).astype('Int64')

    # Encontra a menor data em um vetor de colunas e adiciona informações de idade
    DF_Prim = fg.encontrar_menor_data_vetorizado(DF_Prim, colunas_datas=cols_data_min)
    DF_Prim["data_min"] = pd.to_datetime(DF_Prim["data_min"], errors='coerce').dt.normalize()
    DF_Prim['ano_min'] = DF_Prim["data_min"].dt.year
    DF_Prim['mes_min'] = DF_Prim["data_min"].dt.month

    # Calcula a data mínima de realização de exames (1o da vida)
    colunas_data_min_exame = [col for col in cols_data_min if "disp" not in col]

    DF_Prim = fg.encontrar_menor_data_vetorizado(DF_Prim, colunas_datas=colunas_data_min_exame, nome_dt_min = "data_min_exame")
    DF_Prim["data_min_exame"] = pd.to_datetime(DF_Prim["data_min_exame"], errors='coerce').dt.normalize()
    DF_Prim['ano_min_exame'] = DF_Prim["data_min_exame"].dt.year
    DF_Prim['mes_min_exame'] = DF_Prim["data_min_exame"].dt.month

    # Calcula a idade de vinculação
    DF_Prim["Idade_vinc_cat"] = fg.idade_cat(DF = DF_Prim, data_ref="data_min", data_nasc=col_data_nascimento)
    display(DF_Prim["Idade_vinc_cat"].value_counts().sort_index())
    print()

    # Calcula a idade de vinculação nas faixas anteriores
    DF_Prim["Idade_vinc_cat_anterior"] = fg.idade_cat(DF = DF_Prim, data_ref="data_min", data_nasc=col_data_nascimento, faixas_etarias="MC_antigo")
    display(DF_Prim["Idade_vinc_cat_anterior"].value_counts().sort_index())
    print()

    # Calcula a idade na primeira dispensação
    DF_Prim["Idade_disp_prim"] = fg.idade_cat(DF = DF_Prim, data_ref="data_dispensa", data_nasc=col_data_nascimento)
    display(DF_Prim["Idade_disp_prim"].value_counts().sort_index())
    print()

    # Calcula a idade na primeira dispensação nas categorias do spectrum
    DF_Prim["Idade_Spec_crianca"] = fg.idade_cat(DF = DF_Prim, data_ref="data_dispensa", data_nasc=col_data_nascimento, faixas_etarias="spectrum_crianca")
    display(DF_Prim["Idade_Spec_crianca"].value_counts().sort_index())
    print()   
    DF_Prim["Idade_Spec1"] = fg.idade_cat(DF = DF_Prim, data_ref="data_dispensa", data_nasc=col_data_nascimento, faixas_etarias="spectrum1")
    display(DF_Prim["Idade_Spec1"].value_counts().sort_index())
    print()   
    DF_Prim["Idade_Spec2"] = fg.idade_cat(DF = DF_Prim, data_ref="data_dispensa", data_nasc=col_data_nascimento, faixas_etarias="spectrum2")
    display(DF_Prim["Idade_Spec2"].value_counts().sort_index())
    print()   

   # Cria a coluna com a quantidade de dias entre o diagnóstico e a primeira dispensação
    DF_Prim["Dias_diag_TARV"] = (DF_Prim[col_data_disp_Prim] - DF_Prim["data_min_exame"]).dt.days

    cond = [
        DF_Prim["data_min_exame"].isnull() == True,
        DF_Prim[col_data_disp_Prim].isnull() == True,
        DF_Prim["Dias_diag_TARV"] < -15,
        DF_Prim["Dias_diag_TARV"] < 1,
        DF_Prim["Dias_diag_TARV"] < 8,
        DF_Prim["Dias_diag_TARV"] < 30,
        DF_Prim["Dias_diag_TARV"] < 180,
        DF_Prim["Dias_diag_TARV"] >= 180
    ]

    retorno = [
        "Não tem exames disponíveis",
        "Não iniciou TARV",
        "Iniciou TARV antes de exames no SUS",
        "No mesmo dia",
        "1-7 dias",
        "8-30 dias",
        "1-6 meses",
        "Mais de 6 meses"  
    ]

    DF_Prim["Dias_diag_TARV_cat"] = np.select(cond,retorno, default = "Erro")
    display("Tempo entre vínculo e TARV ",DF_Prim[["ano_min","Dias_diag_TARV_cat"]].value_counts(dropna = False).sort_index())
    print()

    # Cria uma variável com a indetecção viral em 50 cópias
    Cond_sup = [
        (DF_Prim[col_data_disp_Prim] - DF_Prim[col_data_cv_Prim]).dt.days < -15,
        (DF_Prim["CV_cat50"] == "<50"),
        (DF_Prim["CV_cat50"] == ">=50")
    ]

    escolha_sup = [
        "Sem dados de CV antes de TARV","CV indetectável","CV detectável"
    ]

    DF_Prim["Deteccao_CV"]  = np.select(Cond_sup,escolha_sup, default = "Sem dados de CV")
    display("Indetecção Viral antes de TARV",
            DF_Prim["Deteccao_CV"].value_counts(dropna = False).sort_index())


    # Organiza o IBGE de residência
    DF_Prim = fg.ibge_resid(DF=DF_Prim, col_ibge_resid=col_ibge_resid, cols_ibge=cols_ibge)

    # Organiza o IBGE das instituições Solicitantes
    DF_Prim = fg.ibge_inst_sol_exames(DF=DF_Prim, col_cod_ibge_udm = col_cod_ibge_udm, col_cod_ibge_solicitante_cv = col_cod_ibge_solicitante_cv,
                    col_cod_ibge_solicitante_cd4 = col_cod_ibge_solicitante_cd4, col_nome_inst_sol_cv = col_nome_inst_sol_cv,
                    col_nome_inst_sol_cd4 = col_nome_inst_sol_cd4, col_ibge_resid = col_ibge_resid)

    print(DF_Prim.shape)

    # Filtra as crianças menores de 2 anos sem TARV ou sem 2 CV > 5000, os pessoas sem dispensação de TARV e nenhuma CV detectável ou exame de CD4 coletado
    DF_Prim = DF_Prim[((DF_Prim["Idade_vinc_cat_anterior"] != "Menos de 2 anos") | (DF_Prim[col_data_disp_Prim].isnull() == False) | (DF_Prim["CV_5000_count"] >= 2)) &
                      ((DF_Prim[col_data_disp_Prim].isnull() == False) | (DF_Prim["CV_detec_count"] >= 1) | (DF_Prim[col_data_cd4_Prim].isnull() == False)) &
                      (DF_Prim[col_codigo_pac_uni] < 900000000)].copy()
    
    print(DF_Prim.shape)
    
    # Classifica os Municípios com menos de 50 pessoas por ano em uma categoria unificada
    condicao_mun = (DF_Prim.groupby(by = ["ibge_inst_exames", "ano_min"])[col_codigo_pac_uni].transform("count") < 50).fillna(False)

    DF_Prim["Nome_mun_exame_50+"] = np.where(condicao_mun,
                                            "Município com menos de 50 pessoas no cuidado",
                                            DF_Prim["Nome_mun_exame"])
    
    DF_Prim["Cod_ibge_insti_sol_50+"] = np.where(condicao_mun,
                                        (DF_Prim["ibge_inst_exames"].apply(lambda x: str(x)[:2] + "99999")),
                                        (DF_Prim["ibge_inst_exames"]))

    DF_Prim["Cod_ibge_insti_sol_50+"] = pd.to_numeric(DF_Prim["Cod_ibge_insti_sol_50+"], errors = "coerce").astype("Int64")



    # Classifica os Serviços com menos de 50 pessoas por ano em uma categoria unificada
    condicao_inst = (DF_Prim.groupby(by = ["nome_inst_exames","ibge_inst_exames", "ano_min"])[col_codigo_pac_uni].transform("count") < 50).fillna(False)

    DF_Prim["nome_inst_exames_50+"] = np.where(condicao_inst,
                                            "Serviços com menos de 50 pessoas no cuidado",
                                            DF_Prim["nome_inst_exames"])

    return DF_Prim



def gerar_base_ult(hoje: datetime.date, ano_corte:int, colsuso_padrao:bool = True,
                    Cadastro_colsuso: list = None, Disp_colsuso: list = None, CV_colsuso: list = None, obito_admin:bool = True,
                    CD4_colsuso: list = None, SIM_colsuso: list = None, anos_obito_adm:int = 10, data_base: datetime.date = None):
                    
    """
    Carrega, organiza e filtra os bancos de dados do SICLOM/SISCEL conforme a data de análise e o ano de corte especificados.

    Retorna um DataFrame que contém os todos os dados de dispensação de TARV, exames de CD4 e carga viral em uma linha por ano por paciente.
    As informações do último evento de cada ano são mantidas na linha.
    Colunas com o sufixo _Prim são criadas para trazer as informações do primeiros evento da vida. Essas colunas devem ser utilizados com o filtro ["Base_Prim"] == 1.

    Parameters
    ----------
    hoje: datetime.date
        Data de fechamento dos bancos de análise no formato datetime.date(aaaa, mm, dd).

    ano_corte: int
        Ano de corte para filtrar os dados das bases de dispensação, carga viral e contagem de CD4.
        Deve ser um ano válido entre 2000 e o ano da data de fechamento dos bancos.

    colsuso_padrao: bool, optional
        Indica se as colunas padrão devem ser usadas. O padrão é True.
        Se False, utilizará todas as colunas da base, ou as colunas especificadas nas listas de [Base]_colsuso.

    Cadastro_colsuso: list, optional
        Lista de colunas a serem usadas da base de cadastro. colsuso_padrao deve ser passada False para usar a lista especificada.

    Disp_colsuso: list, optional
        Lista de colunas a serem usadas da base de dispensação. colsuso_padrao deve ser passada False para usar a lista especificada.

    CV_colsuso: list, optional
        Lista de colunas a serem usadas da base de carga viral. colsuso_padrao deve ser passada False para usar a lista especificada.

    CD4_colsuso: list, optional
        Lista de colunas a serem usadas da base de contagem de CD4. colsuso_padrao deve ser passada False para usar a lista especificada.

    Obito_colsuso: list, optional
        Lista de colunas a serem usadas da base de óbito. colsuso_padrao deve ser passada False para usar a lista especificada.

    Returns
    -------
    pd.DataFrame
        Um DataFrame contendo os dados organizados e filtrados dos pacientes.

    Raises
    ------
    ValueError
        Se o 'ano_corte' não for um ano válido ou se 'hoje' não for uma instância de datetime.date.


    Example
    -------
    >>> hoje = pd.to_datetime(2024, 1, 31)
    >>> ano_corte = 2024 - 10
    >>> df_completo = gerar_base_completa(hoje, ano_corte)
    """

    # Verificação dos parâmetros de entrada
    if not isinstance(hoje, datetime.date):
        raise ValueError("O parâmetro 'hoje' deve ser uma instância de datetime.date.")
    if not isinstance(ano_corte, int) or ano_corte < 2000 or ano_corte > hoje.year:
        raise ValueError("O parâmetro 'ano_corte' deve ser um ano válido.")

    try:
        # Carregar as bases de dados
        Pac, Disp, CV, CD4, obitos_siclom, SIM, SIM_Total, Cod_uni = carregar_bases(hoje, Cadastro=True, Disp=True, CV=True, CD4=True, SIM=True, Cod_uni=True,
                                                                    colsuso_padrao = colsuso_padrao, Cadastro_colsuso = Cadastro_colsuso, Disp_colsuso = Disp_colsuso,
                                                                    CV_colsuso = CV_colsuso, CD4_colsuso = CD4_colsuso, SIM_colsuso = SIM_colsuso, data_base=data_base)
        
        # Organizar as bases
        Pac, Disp, CV, CD4 = organizar_bases(hoje, Pac, Disp, CV, CD4, Cod_uni, obito={"SIM": SIM, "obitos_siclom": obitos_siclom})
        
        # Organizar as bases pelo ultimo evento de cada ano
        Disp_ano, CV_ano, CD4_ano = bases_ult_ano(hoje, DF_lista=[Disp, CV, CD4], DF_nome=["Disp", "CV", "CD4"], col_data=["ano_disp", "ano_coleta_cv", "ano_coleta_cd4"]).values()
        
        # Criar o DataFrame final
        DF_ult_ano = df_ult_ano(hoje, ano_corte, Pac, Disp_ano, CV_ano, CD4_ano, anos_obito_adm =anos_obito_adm, obito_admin = obito_admin)
        
        return DF_ult_ano

    except Exception as e:
        print(f"Erro ao gerar a base completa: {e}")
        return pd.DataFrame()  # Retorna um DataFrame vazio em caso de erro



def gerar_PVHA_compartilhada(hoje: datetime.date, col_codigo_pac_uni: str = "Cod_unificado",
                            PVHA_cols:list = [
                                "Cod_unificado",
                                "data_min",
                                "origem_data_min",
                                "data_obito",
                                "data_nascimento",
                                "Sexo_cat",
                                "Raca_cat",
                                "data_registro_min",
                                "origem_data_registro_min",
                                "PVHA"
                                ],
                            Sinan_cong_drop:list = [
                                'Unnamed: 0',
                                'id_sinan',
                                'Cod_orig_uni',
                                "NM_PACIENT",
                                "NM_MAE_PAC",
                                "DT_NASC",
                                "CS_SEXO",
                                "CS_RACA"
                                ],
                            Sinan_cr_drop:list = [
                                'Unnamed: 0',
                                'id_sinan',
                                'Cod_orig_uni',
                                "NM_PACIENT",
                                "NM_MAE_PAC",
                                "DT_NASC",
                                "CS_SEXO",
                                "CS_RACA"
                                ],
                            Sinan_G_drop:list = [
                                'Unnamed: 0',
                                'id_sinan',
                                'Cod_orig_uni',
                                "NM_PACIENT",
                                "NM_MAE_PAC",
                                "DT_NASC",
                                "CS_SEXO",
                                "CS_RACA"
                                ],
                            Sinan_A_drop:list = [
                                'Unnamed: 0',
                                'id_sinan',
                                'Cod_orig_uni',
                                "NM_PACIENT",
                                "NM_MAE_PAC",
                                "DT_NASC",
                                "CS_SEXO",
                                "CS_RACA"
                                ],
                            Sim_drop:list = [
                                'Unnamed: 0',
                                'id_sim',
                                'Cod_orig_uni',
                                "NOME",
                                "NOMEMAE",
                                "DTNASC",
                                "SEXO",
                                "RACACOR"
                                ],
                            Sim_total_drop:list = [
                                'Unnamed: 0',
                                'NUMSUS',
                                'NOME',
                                'NOMEPAI',
                                'NOMEMAE',
                                "DTNASC",
                                "SEXO",
                                "RACACOR"
                                ],
                            CD4_cols:list = [
                                "Cod_unificado",
                                'cod_ibge_solicitante_cd4',
                                'cd_inst_sol_cd4',
                                'nm_inst_sol_cd4',
                                'tipo_inst_sol_cd4',
                                'data_solicitacao_cd4',
                                'data_hora_coleta_cd4',
                                'contagem_cd4',
                                'ano_coleta_cd4',
                                'CD4_cat',
                                'paciente_gestante',
                                'nu_idade_gestacional',
                                "nm_metodo_cd4",
                                "nm_kit_cd4"
                                ],
                            CV_cols:list = [
                                "Cod_unificado",
                                'cod_ibge_solicitante_cv',
                                'cd_inst_sol_cv',
                                'nm_inst_sol_cv',
                                'tipo_inst_sol_cv',
                                'data_solicitacao_cv',
                                'data_hora_coleta_cv',
                                'copias',
                                'comentario_copias',
                                'ano_coleta_cv',
                                'CV_cat',
                                'paciente_gestante',
                                'nu_idade_gestacional'
                            ],
                            Disp_cols:list = [
                                "Cod_unificado",
                                'codigo_udm',
                                'nm_udm',
                                'tp_servico_atendimento',
                                'cod_ibge_udm',
                                'st_pub_priv',
                                'data_dispensa',
                                'esquema',
                                'esquema_forma',
                                'duracao',
                                'ano_disp',
                                'categoria_usuario',
                                'categoria_crianca',
                                'periodo_gest',
                                "cd_crm",
                                "uf_crm",
                                'st_paciente_gestante'
                            ],
                            Pac_cols:list = [
                                'codigo_paciente',
                                "Cod_unificado",
                                'codigo_ibge_resid',
                                'uf_residencia',
                                'acomp_medico',
                                'escolaridade',
                                'data_ult_atu',
                                'cd_pais',
                                'co_orientacao_sexual',
                                'co_genero',
                                "st_paciente"
                            ]):

    ano = hoje.year
    caminho = f"V:/{ano}/Monitoramento e Avaliação/COMPARTILHADO/AMA - Banco de Dados/AMA-VIP/Bancos Compartilhados HIV"
     
    Pac, Disp, CV, CD4, obitos_siclom, Sim, Sim_total, Sinan_A, Sinan_cong, Sinan_G, Sinan_cr, Cod_uni = carregar_bases(
        hoje, 
        Cadastro = True,
        Disp = True,
        CV = True,
        CD4=True,
        SIM=True,
        Sinan=True,
        Sinan_G=True,
        Sinan_cr=True,
        Cod_uni = True,
        colsuso_padrao = False,
        Cadastro_colsuso = [
                    "codigo_paciente", "codigo_pac_eleito", "nm_pac", "nm_mae", "data_nascimento", "sexo",
                    "raca", "codigo_ibge_resid", "uf_residencia", "acomp_medico",
                    "escolaridade", "data_ult_atu","cd_pais", "st_estrangeiro", 
                    "co_orientacao_sexual", "co_genero","st_paciente"],
        Disp_colsuso = [
                    "codigo_paciente", "tp_servico_atendimento", "codigo_udm", "cod_ibge_udm","nm_udm",
                    'st_pub_priv', 'data_dispensa', 'esquema', 'esquema_forma', 'duracao','categoria_usuario',
                    'categoria_crianca','periodo_gest','st_paciente_gestante',"cd_crm","uf_crm"
    ],
        CV_colsuso = [
                    "codigo_paciente", "cod_ibge_solicitante_cv", "cd_inst_sol_cv", "nm_inst_sol_cv", "tipo_inst_sol_cv",
                    'data_hora_coleta_cv', 'copias', 'comentario_copias', "data_solicitacao_cv",'paciente_gestante','nu_idade_gestacional'
    ],
        CD4_colsuso = [
                    "codigo_paciente", "cod_ibge_solicitante_cd4", "cd_inst_sol_cd4", "nm_inst_sol_cd4", "tipo_inst_sol_cd4",
                    'data_solicitacao_cd4', 'data_hora_coleta_cd4', "contagem_cd4",'paciente_gestante','nu_idade_gestacional',
                    "nm_metodo_cd4","nm_kit_cd4"
    ],
        Obito_colsuso = [
                    "codigo_paciente", "data_obito"]
    )

    # Organização dos Bancos

    ## Padronização de Variáveis Comuns

    ### Siclom\SISCEL

    ## VERIFICAR OS CODIGO_PACIENTE 811146-811149

    Pac["raca"] = pd.to_numeric(Pac["raca"], errors='coerce')
    Pac["sexo"] = pd.to_numeric(Pac["sexo"], errors='coerce')
    Pac["codigo_paciente"] = Pac["codigo_paciente"].astype(int)
    Cod_uni["codigo_paciente"] = Cod_uni["codigo_paciente"].astype(int)
    Cod_uni[col_codigo_pac_uni] = Cod_uni[col_codigo_pac_uni].astype(int)
    Pac["Hoje"] = hoje

    # Une o cadasto ao código Unificado
    Pac = pd.merge(Pac,Cod_uni[["codigo_paciente",col_codigo_pac_uni]], on = "codigo_paciente", how = "left")

    # Identificar os registros com valores em branco no Cod Unificado
    Sem_cod_prob = Pac[col_codigo_pac_uni].isnull()

    # Definir o valor inicial para os códigos únicos faltantes
    start_value = 800000000

    # Gerar valores únicos para os registros faltantes,
    Pac.loc[Sem_cod_prob, col_codigo_pac_uni] = np.arange(start_value, start_value + Sem_cod_prob.sum())

    Cod_tab = Pac[["codigo_paciente",col_codigo_pac_uni]].copy()

    # Ajusta as Datas
    datas_siclom = ["data_nascimento",
                    "data_ult_atu"]

    for col in datas_siclom:
        Pac = fg.ajusta_data_linha_vetorizado(Pac, col, coluna_retorno = f"{col}_ajustada")
        Pac[col] = Pac[f"{col}_ajustada"]
        Pac = Pac.drop(columns = [f"{col}_ajustada"])

    Pac["data_ref"] = Pac["data_ult_atu"]

    Pac.rename(columns={"nm_pac":"nome","nm_mae":"nome_mae"}, inplace=True)
    Pac["ibge6_res"] = Pac["codigo_ibge_resid"].apply(lambda x:str(x)[:6])

    # Ajusta raça
    Rac = [(Pac["raca"] == 1),(Pac["raca"] == 2),(Pac["raca"] == 3),
            (Pac["raca"] == 4),(Pac["raca"] == 5)]
            
    Rac_escolha = ["Branca","Preta","Amarela","Parda","Indígena"]

    Pac["Raca_cat"] = np.select(Rac, Rac_escolha, default=None)

    # Ajusta sexo
    Sex = [(Pac["sexo"] == 1),(Pac["sexo"] == 2)]
            
    Sex_escolha = ["Homem","Mulher"]

    Pac["Sexo_cat"] = np.select(Sex, Sex_escolha, default=None)

    Pac[[col_codigo_pac_uni,"data_nascimento","Raca_cat","Sexo_cat","data_ref"]]

    ### SIM
    Sim_total[col_codigo_pac_uni] = Sim_total["COD_UNIFICADO"].astype(int)

    SIM_colsuso = [
        "NOME",
        "NOMEMAE",
        "CODMUNRES",
        "DTOBITO",
        "DTNASC",
        "SEXO",
        "RACACOR",
        "id_sim",
        col_codigo_pac_uni
    ]

    SIM_total_colsuso = [
        "DTOBITO",
        "DTNASC",
        "SEXO",
        "RACACOR",
        "CODMUNRES",
        col_codigo_pac_uni
    ]

    # Une o SIM Total
    linhas_total = Sim[SIM_colsuso].shape[0] + Sim_total[SIM_total_colsuso].shape[0]

    colunas_total = Sim[SIM_colsuso].shape[1]

    SIM = pd.concat([Sim[SIM_colsuso],Sim_total[SIM_total_colsuso]], ignore_index = True)

    if (SIM.shape[1] != colunas_total) | (SIM.shape[0] != linhas_total):
        print("Erro ao unir o Sim_AIDS com o Sim_Completo")

    SIM[col_codigo_pac_uni] = SIM[col_codigo_pac_uni].astype("Int64")

    # Identificar os registros com valores em branco no ID_Sim
    Sem_cod_prob = SIM["id_sim"].isnull()

    # Definir o valor inicial para os códigos faltantes
    start_value = 210000000

    # Gerar valores únicos para os registros faltantes,
    SIM.loc[Sem_cod_prob, "id_sim"] = np.arange(start_value, start_value + Sem_cod_prob.sum())

    # Ajusta o nome das colunas
    SIM = SIM.rename(columns={col: f"{col}_sim" for col in SIM.columns[:-2]})

    # Padroniza o formato das datas
    datas_sim = ["DTOBITO_sim","DTNASC_sim"]

    for col in datas_sim:
        SIM = fg.ajusta_data_linha_vetorizado(SIM, col, coluna_retorno = f"{col}_ajustada")
        SIM[col] = SIM[f"{col}_ajustada"]
        SIM = SIM.drop(columns = [f"{col}_ajustada"])

    SIM["data_nascimento"] = SIM["DTNASC_sim"]
    SIM["data_ref"] = SIM["DTOBITO_sim"]
    SIM["nome"] = SIM["NOME_sim"]
    SIM["nome_mae"] = SIM["NOMEMAE_sim"]
    SIM["ibge6_res"] = SIM["CODMUNRES_sim"].apply(lambda x:str(x)[:6])

    # Ajusta a Raça/Cor
    SIM["RACACOR_sim"] = SIM["RACACOR_sim"].apply(lambda x:str(x).strip())
    Rac_cond = [(SIM["RACACOR_sim"] == "1.0"),
                (SIM["RACACOR_sim"] == "3.0"),
                (SIM["RACACOR_sim"] == "2.0"),
                (SIM["RACACOR_sim"] == "4.0"),
                (SIM["RACACOR_sim"] == "5.0")
                ]
            
    Rac_escolha = ["Branca", "Amarela","Preta","Parda","Indígena"]
    SIM["Raca_cat"] = np.select(Rac_cond,Rac_escolha, default = None)

    # Ajusta o Sexo
    Sex = [(SIM["SEXO_sim"] == "M"),(SIM["SEXO_sim"] == "F")]
            
    Sex_escolha = ["Homem","Mulher"]

    SIM["Sexo_cat"] = np.select(Sex, Sex_escolha, default=None)

    ### Sinan

    Sinan_colsuso = [
                    "NM_PACIENT", "NM_MAE_PAC", "DT_NOTIFIC","CS_RACA","DT_DIAG","DT_NASC","CS_SEXO","ID_MN_RESI","ID_MUNICIP","CRITERIO","DT_CONFIRM",
                    "DT_SIN_PRI","DT_RAPIDO","DT_LAB_HIV","id_sinan",col_codigo_pac_uni]

    Sinan_cong_colsuso = [
                    'NM_PACIENT', 'NM_MAE_PAC', 'DT_NASC', 'ID_MN_RESI',
                    'DT_DIAG', 'DT_NOTIFIC', 'ID_MUNICIP', 'CS_SEXO', 'CRITERIO', 'ANT_REL_CA', 'CS_RACA',
                    'id_sinan', col_codigo_pac_uni]

    SinanG_colsuso  = [
                    "NM_PACIENT", "NM_MAE_PAC","DT_NOTIFIC","CS_RACA","DT_DIAG","DT_NASC","CS_SEXO","DT_SIN_PRI",
                    "ID_MN_RESI","ID_MUNICIP","id_sinan",col_codigo_pac_uni]

    Sinan_cr_colsuso = [
                    "NM_PACIENT", "NM_MAE_PAC","DT_NOTIFIC","CS_RACA","DT_DIAG","DT_NASC","CS_SEXO",
                    "ID_MN_RESI","ID_MUNICIP","CRITERIO","DT_SIN_PRI", "DTCONFIRMA", "DTRAPIDO1","DT_LAB_HIV",
                    "DT_PCR_1","DT_PCR_2","DT_PCR_3","id_sinan",col_codigo_pac_uni]

    Sinan_A_org = ajuste_sinan_var_padrao(Sinan_A[Sinan_colsuso],"sinan")

    Sinan_G_org = ajuste_sinan_var_padrao(Sinan_G[SinanG_colsuso],"gest")

    Sinan_cr_org = ajuste_sinan_var_padrao(Sinan_cr[Sinan_cr_colsuso],"cr")

    Sinan_cong_org = ajuste_sinan_var_padrao(Sinan_cong[Sinan_cong_colsuso],"cong")


    ### União das variáveis padrão

    var_padrao = [
        col_codigo_pac_uni,
        "nome",
        "nome_mae",
        "ibge6_res",
        "data_nascimento",
        "Raca_cat",
        "Sexo_cat",
        "data_ref"
    ]

    DF_var = pd.concat([
        Pac[var_padrao],
        SIM[var_padrao],
        Sinan_A_org[var_padrao],
        Sinan_G_org[var_padrao],
        Sinan_cr_org[var_padrao],
        Sinan_cong_org[var_padrao]
    ], ignore_index = True)


    # Lista de variáveis a serem padronizadas
    var_list = [
        "nome",
        "nome_mae",
        "ibge6_res",
        "data_nascimento",
        "Raca_cat",
        "Sexo_cat"
    ]

    # Aplique a padronização para cada variável da lista
    for var in var_list:
        if var in ["nome","nome_mae"]:
                DF_var[var] = DF_var[var].apply(lambda x: str(x).strip())
                DF_var[var] = DF_var[var].apply(lambda x: unicodedata.normalize("NFD", str(x)))   
                DF_var[var] = DF_var[var].apply(lambda x: x.encode("ascii", "ignore")) 
                DF_var[var] = DF_var[var].apply(lambda x: x.decode("utf-8")) 
                DF_var[var] = DF_var[var].apply(lambda x: x.lower().strip())
                DF_var[var] = DF_var[var].apply(fl.limpa_texto)
        DF_var = fg.padronizar_variaveis_vetorizado(DF_var, var)

    DF_var.drop_duplicates(subset = [col_codigo_pac_uni,"data_nascimento","Raca_cat","Sexo_cat"], inplace = True)
    DF_var[col_codigo_pac_uni] = DF_var[col_codigo_pac_uni].astype("Int64")
    DF_var = DF_var.drop(columns = ["data_ref"])

    ### Definindo Data mínima de cada banco

    #### SICLOM/SISCEL

    # Organiza a dispensação e ajusta o Cod_unificado
    Disp = organizacao_disp(
        hoje, Disp, Cod_tab, 
        esquemas = False,
        Prim = False
    )
    Disp["duracao"] = Disp["duracao_sum"]

    CV = organizacao_cv(
        hoje, CV, Cod_tab,
        categorizar = True,
        Prim = False
    )

    CD4 = organizacao_cd4(
        hoje, CD4, Cod_tab, 
        categorizar = True,
        Prim = False
    )

    Disp_Prim = Disp[[col_codigo_pac_uni,"data_dispensa"]].copy()
    Disp_Prim = Disp_Prim.sort_values(["data_dispensa"], ascending = True)
    Disp_Prim = Disp_Prim.reset_index(drop = True)
    Disp_Prim = Disp_Prim.drop_duplicates([col_codigo_pac_uni], keep = "first")

    CV_Prim = CV[[col_codigo_pac_uni,"data_hora_coleta_cv","data_solicitacao_cv","CV_5000_count","CV_detec_count"]].copy()
    CV_Prim = CV_Prim.sort_values(["data_hora_coleta_cv"], ascending = True)
    CV_Prim = CV_Prim.reset_index(drop = True)
    CV_Prim = CV_Prim.drop_duplicates([col_codigo_pac_uni], keep = "first")

    CD4_Prim = CD4[[col_codigo_pac_uni,"data_hora_coleta_cd4","data_solicitacao_cd4"]].copy()
    CD4_Prim = CD4_Prim.sort_values(["data_hora_coleta_cd4"], ascending = True)
    CD4_Prim = CD4_Prim.reset_index(drop = True)
    CD4_Prim = CD4_Prim.drop_duplicates([col_codigo_pac_uni], keep = "first")

    Disp_Prim = pd.merge(Disp_Prim, DF_var[[col_codigo_pac_uni,"data_nascimento"]], on = col_codigo_pac_uni, how = "left")
    mask_disp = ((Disp_Prim["data_dispensa"] > pd.to_datetime(hoje)) |
                (Disp_Prim["data_dispensa"] < pd.to_datetime("2000-01-01")) |
                (Disp_Prim["data_dispensa"] < Disp_Prim["data_nascimento"]))  
    Disp_Prim.loc[mask_disp,"data_dispensa"] = pd.NaT

    CD4_Prim = pd.merge(CD4_Prim, DF_var[[col_codigo_pac_uni,"data_nascimento"]], on = col_codigo_pac_uni, how = "left")
    mask_cd4 = ((CD4_Prim["data_solicitacao_cd4"] > pd.to_datetime(hoje)) |
                (CD4_Prim["data_solicitacao_cd4"] < pd.to_datetime("2000-01-01")) |
                (CD4_Prim["data_solicitacao_cd4"] < CD4_Prim["data_nascimento"]))  
    CD4_Prim.loc[mask_cd4,"data_solicitacao_cd4"] = pd.NaT

    mask_cd4b = ((CD4_Prim["data_hora_coleta_cd4"] > pd.to_datetime(hoje)) |
                (CD4_Prim["data_hora_coleta_cd4"] < pd.to_datetime("2000-01-01")) |
                (CD4_Prim["data_hora_coleta_cd4"] < CD4_Prim["data_nascimento"]))  
    CD4_Prim.loc[mask_cd4b,"data_hora_coleta_cd4"] = pd.NaT

    CV_Prim = pd.merge(CV_Prim, DF_var[[col_codigo_pac_uni,"data_nascimento"]], on = col_codigo_pac_uni, how = "left")
    mask_cv = ((CV_Prim["data_solicitacao_cv"] > pd.to_datetime(hoje)) |
                (CV_Prim["data_solicitacao_cv"] < pd.to_datetime("2000-01-01")) |
                (CV_Prim["data_solicitacao_cv"] < CV_Prim["data_nascimento"]))  
    CV_Prim.loc[mask_cv,"data_solicitacao_cv"] = pd.NaT

    mask_cv2 = ((CV_Prim["data_hora_coleta_cv"] > pd.to_datetime(hoje)) |
                (CV_Prim["data_hora_coleta_cv"] < pd.to_datetime("2000-01-01")) |
                (CV_Prim["data_hora_coleta_cv"] < CV_Prim["data_nascimento"]))  
    CV_Prim.loc[mask_cv2,"data_hora_coleta_cv"] = pd.NaT

    Disp_Prim = Disp_Prim.drop(columns=["data_nascimento"])
    CD4_Prim = CD4_Prim.drop(columns=["data_nascimento"])
    CV_Prim = CV_Prim.drop(columns=["data_nascimento"])

    ## Adicionando o Cod_uni aos óbitos do siclom
    obitos_siclom = pd.merge(obitos_siclom,Cod_uni[["codigo_paciente",col_codigo_pac_uni]], on = "codigo_paciente", how = "left")
    obitos_siclom = fg.ajusta_data_linha_vetorizado(obitos_siclom, "data_obito", coluna_retorno = "data_obito_ajustada")
    obitos_siclom["data_obito"] = obitos_siclom["data_obito_ajustada"]
    obitos_siclom[col_codigo_pac_uni] = obitos_siclom[col_codigo_pac_uni].astype("Int64")

    obitos_siclom = pd.merge(obitos_siclom, DF_var[[col_codigo_pac_uni,"data_nascimento"]], on = col_codigo_pac_uni, how = "left")
    mask_obito = ((obitos_siclom["data_obito"] > pd.to_datetime(hoje)) |
                (obitos_siclom["data_obito"] < pd.to_datetime("2000-01-01")) |
                (obitos_siclom["data_obito"] < obitos_siclom["data_nascimento"]))  
    obitos_siclom.loc[mask_obito,"data_obito"] = pd.NaT

    obitos_siclom = obitos_siclom.drop(columns = ["codigo_paciente","data_obito_ajustada","data_nascimento"])

    #### SIM

    SIM.drop(columns = ["data_nascimento"], inplace = True)
    SIM = pd.merge(SIM, DF_var[[col_codigo_pac_uni,"data_nascimento"]], on = col_codigo_pac_uni, how = "left")
    SIM["data_obito"] = SIM["DTOBITO_sim"]

    mask_obito = ((SIM["data_obito"] > pd.to_datetime(hoje)) |
                (SIM["data_obito"] < pd.to_datetime("2000-01-01")) |
                (SIM["data_obito"] < SIM["data_nascimento"]))

    SIM.loc[mask_obito,"data_obito"] = pd.NaT

    SIM.drop(columns = ["nome","nome_mae","ibge6_res","data_nascimento","Raca_cat","Sexo_cat","data_ref"], inplace = True)

    #### SINAN_A

    datas_sinan_min = [
        "DT_NOTIFIC_sinan",
        "DT_DIAG_sinan",
        "DT_CONFIRM_sinan",
        "DT_RAPIDO_sinan",
        "DT_SIN_PRI_sinan",
        "DT_LAB_HIV_sinan"
    ]

    Sinan_A_org = Sinan_A_org[Sinan_A_org["CRITERIO_sinan"] != 900].copy()

    Sinan_A_org.drop(columns = ["data_nascimento"], inplace = True)
    Sinan_A_org = pd.merge(Sinan_A_org, DF_var[[col_codigo_pac_uni,"data_nascimento"]], on = col_codigo_pac_uni, how = "left")

    # Encontra a menor data em um vetor de colunas
    for col in datas_sinan_min:

        mask = ((Sinan_A_org[col] > pd.to_datetime(hoje)) |
                (Sinan_A_org[col] < pd.to_datetime("1991-01-01")) |
                (Sinan_A_org[col] < Sinan_A_org["data_nascimento"]))        
        
        Sinan_A_org.loc[mask,col] = pd.NaT

    # Encontra a menor data em um vetor de colunas
    Sinan_A_org = fg.encontrar_menor_data_vetorizado(Sinan_A_org, 
                                            colunas_datas = datas_sinan_min,
                                            nome_dt_min = "data_min_sinan")

    Sinan_A_org.drop(columns = ["nome","nome_mae","ibge6_res","data_nascimento","Raca_cat","Sexo_cat","data_ref"], inplace = True)


    #### SINAN_Cong

    Sinan_cong_org["data_min_congelado"] = Sinan_cong_org["DT_DIAG_cong"] 

    Sinan_cong_org.drop(columns = ["nome","nome_mae","ibge6_res","data_nascimento","Raca_cat","Sexo_cat","data_ref"], inplace = True)


    #### SINAN_G

    datas_sinan_min = [
        "DT_NOTIFIC_gest",
        "DT_DIAG_gest",
        "DT_SIN_PRI_gest"
    ]

    Sinan_G_org.drop(columns = ["data_nascimento"], inplace = True)
    Sinan_G_org = pd.merge(Sinan_G_org, DF_var[[col_codigo_pac_uni,"data_nascimento"]], on = col_codigo_pac_uni, how = "left")

    # Encontra a menor data em um vetor de colunas
    for col in datas_sinan_min:

        mask = ((Sinan_G_org[col] > pd.to_datetime(hoje)) |
                (Sinan_G_org[col] < pd.to_datetime("2000-01-01")) |
                (Sinan_G_org[col] < Sinan_G_org["data_nascimento"]))

        Sinan_G_org.loc[mask,col] = pd.NaT

    # Encontra a menor data em um vetor de colunas
    Sinan_G_org = fg.encontrar_menor_data_vetorizado(Sinan_G_org, 
                                            colunas_datas = datas_sinan_min,
                                            nome_dt_min = "data_min_gest")

    Sinan_G_org.drop(columns = ["nome","nome_mae","ibge6_res","data_nascimento","Raca_cat","Sexo_cat","data_ref"], inplace = True)

    #### SINAN_cr

    datas_sinan_min = [
        "DT_NOTIFIC_cr",
        "DT_DIAG_cr",
        "DT_SIN_PRI_cr",
        "DTCONFIRMA_cr",
        "DTRAPIDO1_cr",
        "DT_LAB_HIV_cr",
        "DT_PCR_1_cr",
        "DT_PCR_2_cr",
        "DT_PCR_3_cr"
    ]

    Sinan_cr_org = Sinan_cr_org[Sinan_cr_org["CRITERIO_cr"] != 900].copy()

    Sinan_cr_org.drop(columns = ["data_nascimento"], inplace = True)
    Sinan_cr_org = pd.merge(Sinan_cr_org, DF_var[[col_codigo_pac_uni,"data_nascimento"]], on = col_codigo_pac_uni, how = "left")

    # Encontra a menor data em um vetor de colunas
    for col in datas_sinan_min:

        mask = ((Sinan_cr_org[col] > pd.to_datetime(hoje)) |
                (Sinan_cr_org[col] < pd.to_datetime("1991-01-01")) |
                (Sinan_cr_org[col] < Sinan_cr_org["data_nascimento"]))

        Sinan_cr_org.loc[mask,col] = pd.NaT

    # Encontra a menor data em um vetor de colunas
    Sinan_cr_org = fg.encontrar_menor_data_vetorizado(Sinan_cr_org, 
                                            colunas_datas = datas_sinan_min,
                                            nome_dt_min = "data_min_cr")

    Sinan_cr_org.drop(columns = ["nome","nome_mae","ibge6_res","data_nascimento","Raca_cat","Sexo_cat","data_ref"], inplace = True)

    ### Definição de data mínima global

    # Pega apenas a entrada mais antiga
    SIM = SIM.sort_values(["data_obito"], ascending = True)
    Sim_Prim = SIM.drop_duplicates(subset = [col_codigo_pac_uni], keep="first")
    Sim_Prim = Sim_Prim.reset_index(drop = True)

    Sinan_A_org = Sinan_A_org.sort_values(["data_min_sinan"], ascending = True)
    Sinan_Prim = Sinan_A_org.drop_duplicates(subset = [col_codigo_pac_uni], keep="first")
    Sinan_Prim = Sinan_Prim.reset_index(drop = True)

    Sinan_G_org = Sinan_G_org.sort_values(["data_min_gest"], ascending = True)
    Sinan_G_Prim = Sinan_G_org.drop_duplicates(subset = [col_codigo_pac_uni], keep="first")
    Sinan_G_Prim = Sinan_G_Prim.reset_index(drop = True)

    Sinan_cr_org = Sinan_cr_org.sort_values(["data_min_cr"], ascending = True)
    Sinan_cr_Prim = Sinan_cr_org.drop_duplicates(subset = [col_codigo_pac_uni], keep="first")
    Sinan_cr_Prim = Sinan_cr_Prim.reset_index(drop = True)

    obitos_siclom = obitos_siclom.sort_values(["data_obito"], ascending = [True])
    obitos_siclom_Prim = obitos_siclom.drop_duplicates([col_codigo_pac_uni], keep = "first")
    obitos_siclom_Prim = obitos_siclom_Prim.reset_index(drop = True)

    DF = pd.merge(DF_var, Disp_Prim, on = col_codigo_pac_uni, how = "left")
    DF = pd.merge(DF, CD4_Prim, on = col_codigo_pac_uni, how = "left")
    DF = pd.merge(DF, CV_Prim, on = col_codigo_pac_uni, how = "left")
    DF = pd.merge(DF, Sinan_Prim, on = col_codigo_pac_uni, how = "left")
    DF = pd.merge(DF, Sinan_cong_org, on = col_codigo_pac_uni, how = "left", suffixes = ("", "_cong"))
    DF = pd.merge(DF, Sinan_G_Prim, on = col_codigo_pac_uni, how = "left", suffixes = ("", "_gest"))
    DF = pd.merge(DF, Sinan_cr_Prim, on = col_codigo_pac_uni, how = "left", suffixes = ("", "_cr"))
    DF = pd.merge(DF, Sim_Prim, on = col_codigo_pac_uni, how = "left")
    DF = pd.merge(DF, obitos_siclom_Prim, on = col_codigo_pac_uni, how = "left", suffixes = ("", "_siclom"))

    ### Organização de Variáveis

    # Unifica a data de obito
    DF["data_obito"] = DF["data_obito"].fillna(DF["data_obito_siclom"])
    DF["ano_obito"] = DF["data_obito"].dt.year

    colunas_datamin = [
        "data_min_congelado",
        "data_min_sinan",
        "data_min_gest",
        "data_min_cr",
        "data_obito",
        "data_solicitacao_cv",
        "data_hora_coleta_cv",
        "data_solicitacao_cd4",
        "data_hora_coleta_cd4",
        "data_dispensa"
    ]

    # Encontra a menor data em um vetor de colunas e adiciona informações de idade
    DF = fg.encontrar_menor_data_vetorizado(DF, 
                                            colunas_datas = colunas_datamin,
                                            nome_dt_min = "data_min")
    DF["ano_min"] = DF["data_min"].dt.year

    colunas_data_registro = [
        "data_min_congelado",
        "DT_NOTIFIC_sinan",
        "DT_NOTIFIC_gest",
        "DT_NOTIFIC_cr",
        "data_obito",
        "data_hora_coleta_cv",
        "data_hora_coleta_cd4",
        "data_dispensa"
    ]

    # Encontra a menor data em um vetor de colunas e adiciona informações de idade
    DF = fg.encontrar_menor_data_vetorizado(DF, 
                                            colunas_datas = colunas_data_registro,
                                            nome_dt_min = "data_registro_min")


    ### Limpeza da Base de PVHA

    DF1 = DF[~DF["ano_min"].isnull()].copy()

    # Calcula a idade de vinculação nas faixas anteriores
    DF1["Idade_Diag"] = fg.idade_cat(DF = DF1, data_ref="data_min", data_nasc="data_nascimento", faixas_etarias="MC_antigo")


    DF2 = DF1[((DF1["Idade_Diag"] != "Menos de 2 anos") | (DF1["data_dispensa"].isnull() == False) | (DF1["CV_5000_count"] >= 2) | (DF1["data_min_cr"].isnull() == False) | (DF1["data_min_congelado"].isnull() == False)) &
        ((DF1["data_dispensa"].isnull() == False) | (DF1["CV_detec_count"] >= 1) | (DF1["data_min_sinan"].isnull() == False) | (DF1["data_min_congelado"].isnull() == False) | (DF1["data_min_gest"].isnull() == False) | (DF1["data_min_cr"].isnull() == False) | ((DF1["data_obito"].isnull() == False) & (DF1["id_sim"] < 210000000))) &
        (DF1[col_codigo_pac_uni] < 900000000)].copy()

    DF2["PVHA"] = "Sim"

    print()

    display(DF2[PVHA_cols])
    print()
    print()

    display(round(DF2["origem_data_min"].value_counts(dropna = False, normalize = True)*100,1))
    print()

    display(round(DF2["origem_data_registro_min"].value_counts(dropna = False, normalize = True)*100,1))
    print()

    Sinan_cr = Sinan_cr.drop(columns = Sinan_cr_drop)

    Sinan_G = Sinan_G.drop(columns = Sinan_G_drop)

    Sinan_A = Sinan_A.drop(columns = Sinan_A_drop)

    Sinan_cong = Sinan_cong.drop(columns = Sinan_cong_drop)

    Sim = Sim.drop(columns = Sim_drop)

    Sim_total = Sim_total.drop(columns = Sim_total_drop)

    # Teste de integridade das bases:

    assert not Sinan_cr.empty, "Sinan_cr está vazio"
    assert not Sinan_G.empty, "Sinan_G está vazio"
    assert not Sinan_A.empty, "Sinan_A está vazio"
    assert not Sinan_cong.empty, "Sinan_cong está vazio"
    assert not Sim.empty, "Sim está vazio"
    assert not Sim_total.empty, "Sim_total está vazio"
    assert not CD4.empty, "CD4 está vazio"
    assert not CV.empty, "CV está vazio"
    assert not Disp.empty, "Disp está vazio"
    assert not Pac.empty, "Pac está vazio"
    assert not DF2.empty, "PVHA está vazio"

    assert len(DF2) > 1500000, "PVHA está com número menor do que o esperado"


    # FIXAR VARIÁVEIS NUMÉRICAS COMO CATEGÓRICAS (NU_NOTIF, ID_DO, COD_PAC)
    Pac["codigo_paciente"] = Pac["codigo_paciente"].apply(lambda x: str(int(x)).zfill(8)).astype(str)
    Sim["ID_DO"] = Sim["ID_DO"].apply(lambda x: str(x).strip().zfill(8) if str(x).strip().isdigit() else str(x).strip()).astype(str)
    Sinan_A["NU_NOTIFIC"] = Sinan_A["NU_NOTIFIC"].apply(lambda x: str(int(str(x).replace(" ", ""))).strip().zfill(7)).astype(str)
    Sinan_G["NU_NOTIFIC"] = Sinan_G["NU_NOTIFIC"].apply(lambda x: str(int(str(x).replace(" ", ""))).strip().zfill(7)).astype(str)
    Sinan_cr["NU_NOTIFIC"] = Sinan_cr["NU_NOTIFIC"].apply(lambda x: str(int(str(x).replace(" ", ""))).strip().zfill(7)).astype(str)


    # Exportação das bases
    Sinan_cr.to_csv(f"{caminho}/Sinan_hiv_crianca.csv",
                            sep = ";", encoding = "latin-1", index = False)
    Sinan_G.to_csv(f"{caminho}/Sinan_hiv_gestante.csv",
                            sep = ";", encoding = "latin-1", index = False)
    Sinan_A.to_csv(f"{caminho}/Sinan_hiv_adulto.csv",
                            sep = ";", encoding = "latin-1", index = False)
    Sinan_cong.to_csv(f"{caminho}/Sinan_hiv_congelado.csv",
                            sep = ";", encoding = "latin-1", index = False)
    Sim.to_csv(f"{caminho}/SIM_aids.csv",
                            sep = ";", encoding = "latin-1", index = False)
    Sim_total.to_csv(f"{caminho}/SIM_outros.csv",
                            sep = ";", encoding = "latin-1", index = False)
    CD4[CD4_cols].to_csv(f"{caminho}/CD4.csv",
                            sep = ";", encoding = "latin-1", index = False)
    CV[CV_cols].to_csv(f"{caminho}/CV.csv",
                            sep = ";", encoding = "latin-1", index = False)
    Disp[Disp_cols].to_csv(f"{caminho}/Dispensacao.csv",
                            sep = ";", encoding = "latin-1", index = False)
    Pac[Pac_cols].to_csv(f"{caminho}/Cadastro.csv",
                            sep = ";", encoding = "latin-1", index = False)
    DF2[PVHA_cols].to_csv(f"{caminho}/PVHA.csv",
                            sep = ";", encoding = "latin-1", index = False)
    

    print(f"Arquivos salvos na pasta {caminho}")

    DF2[[col_codigo_pac_uni,
         "nome",
         "nome_mae",
         "ibge6_res",
         "data_nascimento",
         "Sexo_cat",
         "data_min",
         "data_obito",
         "data_registro_min",
         "PVHA"]].to_csv(f"B:/Bancos Atuais HIV/Nominal/PVHA_nominal.csv",
                            sep = ";", encoding = "latin-1", index = False)
    

    print(f"Arquivo nominal na pasta B:/Bancos Atuais HIV/Nominal/")

    
    # Limpeza da memória

    del Pac, Disp, CV, CD4, obitos_siclom, Sim, Sim_total, Sinan_A, Sinan_G, Sinan_cr, Cod_uni, \
    Cod_tab, SIM, Sinan_A_org, Sinan_G_org, Sinan_cr_org, DF_var, Disp_Prim, CV_Prim, CD4_Prim, \
    Sim_Prim, Sinan_Prim, Sinan_G_Prim, Sinan_cr_Prim, obitos_siclom_Prim, DF, DF1, DF2, \
    colunas_datamin, datas_siclom, datas_sim, datas_sinan_min, \
    Rac, Rac_escolha, Rac_cond, Sex, Sex_escolha, Sem_cod_prob, start_value, var_padrao, \
    var_list, col, linhas_total, colunas_total, mask_cd4, mask_cv, mask_obito, mask, \
    Sinan_cr_drop, Sinan_G_drop, Sinan_A_drop, Sim_drop, Sim_total_drop

    gc.collect()



def gerar_bases_PVHA(hoje: datetime.date,
                     ano_corte:int,
                     col_codigo_pac_uni: str = "Cod_unificado",
                     col_ano = "ano",
                     col_ano_disp= "ano_disp",
                     col_ano_cv = "ano_coleta_cv",
                     col_ano_cd4 = "ano_coleta_cd4",
                     cols_datas = ["data_dispensa","data_hora_coleta_cv", "data_hora_coleta_cd4"],
                     col_data_nascimento = "data_nascimento",
                     colunas_data_min_exame = ["data_coleta_cd4_Prim", "data_coleta_cv_Prim","data_sol_cd4_Prim","data_sol_cv_Prim"],
                     col_copias = "copias",
                     col_comentario_copias = "comentario_copias",
                     col_contagem_cd4 = "contagem_cd4",
                     obito_admin = True,
                     anos_obito_adm = 10,
                     data_disp_prim = "data_dispensa_Prim"):
    
    ano = hoje.year

    # Carregamento dos Bancos
    PVHA, Pac, Disp, CV, CD4, Sinan_A, Sinan_cong, Sinan_G, Sinan_cr, Sim = carregar_bases_PVHA(
        hoje=hoje,
        PVHA = True,
        Cadastro = True,
        Disp = True,
        CV = True,
        CD4=True,
        Sinan=True,
        Sinan_G=True,
        Sinan_cr=True,
        Sim=True
    )

    # Organização dos Bancos

    ### Siclom\SISCEL

    # Organiza o cadastro
    PVHA = organizacao_cadastro(
        PVHA,
        Pac,
        hoje
    )

    # Organiza a dispensação
    Disp = organizacao_disp(
        hoje, Disp,
        esquemas = True,
        Prim = True
    )

    CV = organizacao_cv(
        hoje, CV,
        categorizar = True,
        Prim = True
    )

    CD4 = organizacao_cd4(
        hoje, CD4, 
        categorizar = True,
        Prim = True
    )

    ### SINAN

    Sinan_A  = organizacao_sinan(Sinan_A, tipo = "Ad")

    Sinan_G  = organizacao_sinan(Sinan_G, tipo = "gest")

    Sinan_cr = organizacao_sinan(Sinan_cr, tipo = "cr")

    Sinan_cong["CS_ESCOL_N"] = np.nan
    Sinan_cong = organizacao_sinan(Sinan_cong, tipo = "cong")


    ########################### UNIÃO DOS BANCOS ###############################################

    Disp_Prim, CV_Prim, CD4_Prim = bases_prim(DF_lista = [Disp, CV, CD4], DF_nome = ["Disp", "CV", "CD4"], col_data = ["data_dispensa", "data_hora_coleta_cv", "data_hora_coleta_cd4"]).values()

    # Pega apenas a entrada mais antiga
    Sinan_A = Sinan_A.sort_values(["DT_NOTIFIC_Ad"], ascending = True)
    Sinan_Prim = Sinan_A.drop_duplicates(subset = [col_codigo_pac_uni], keep="first")
    Sinan_Prim = Sinan_Prim.reset_index(drop = True)

    Sinan_G = Sinan_G.sort_values(["DT_NOTIFIC_gest"], ascending = True)
    Sinan_G_Prim = Sinan_G.drop_duplicates(subset = [col_codigo_pac_uni], keep="first")
    Sinan_G_Prim = Sinan_G_Prim.reset_index(drop = True)

    Sinan_cr = Sinan_cr.sort_values(["DT_NOTIFIC_cr"], ascending = True)
    Sinan_cr_Prim = Sinan_cr.drop_duplicates(subset = [col_codigo_pac_uni], keep="first")
    Sinan_cr_Prim = Sinan_cr_Prim.reset_index(drop = True)

    Sim = Sim.drop_duplicates(subset = [col_codigo_pac_uni], keep="first")
    Sim = Sim.reset_index(drop = True)


    DF = pd.merge(PVHA, Disp_Prim, on = col_codigo_pac_uni, how = "left")
    DF = pd.merge(DF, CD4_Prim, on = col_codigo_pac_uni, how = "left")
    DF = pd.merge(DF, CV_Prim, on = col_codigo_pac_uni, how = "left")
    DF = pd.merge(DF, Sinan_Prim, on = col_codigo_pac_uni, how = "left")
    DF = pd.merge(DF, Sinan_cong, on = col_codigo_pac_uni, how = "left")
    DF = pd.merge(DF, Sinan_G_Prim, on = col_codigo_pac_uni, how = "left")
    DF = pd.merge(DF, Sinan_cr_Prim, on = col_codigo_pac_uni, how = "left")
    DF = pd.merge(DF, Sim, on = col_codigo_pac_uni, how = "left")


    del Disp_Prim, CV_Prim, CD4_Prim, Sinan_A, Sinan_G, Sinan_cr
    gc.collect()

    print(DF.shape)
    print()
    display(DF["ano_min"].value_counts().sort_index(ascending = False))
    print()
    display(DF["ano_obito"].value_counts().sort_index(ascending = False))
    print()

    Disp.sort_values(["data_dispensa"], ascending = [False], inplace = True)
    Disp.reset_index(drop = True, inplace = True)

    CV.sort_values(["data_hora_coleta_cv"], ascending = [False], inplace = True)
    CV.reset_index(drop = True, inplace = True)

    CD4.sort_values(["data_hora_coleta_cd4"], ascending = [False], inplace = True)
    CD4.reset_index(drop = True, inplace = True)

    # Organizar as bases pelo ultimo evento de cada ano
    PVHA_ano, Disp_ano, CV_ano, CD4_ano = bases_ult_ano(hoje, DF_lista=[PVHA, Disp, CV, CD4], DF_nome=["PVHA","Disp", "CV", "CD4"], col_data=["ano_min","ano_disp", "ano_coleta_cv", "ano_coleta_cd4"]).values()

    print("Base PVHA_ano")
    display(PVHA_ano[col_ano].value_counts().sort_index(ascending = False))
    print()

    DF_ult_ano = pd.merge(PVHA_ano,Disp_ano, on = [col_ano, col_codigo_pac_uni], how = "left")
    DF_ult_ano = pd.merge(DF_ult_ano,CD4_ano, on = [col_ano,col_codigo_pac_uni], how = "left")
    DF_ult_ano = pd.merge(DF_ult_ano,CV_ano, on = [col_ano,col_codigo_pac_uni], how = "left")

    DF_ult_ano = pd.merge(DF_ult_ano, Sinan_Prim[[
        col_codigo_pac_uni,'Escol_num_Ad',"ID_MN_RESI_Ad",
        "ID_MUNICIP_Ad","DT_NOTIFIC_Ad"]], on = col_codigo_pac_uni, how = "left")
    
    DF_ult_ano = pd.merge(DF_ult_ano, Sinan_G_Prim[[
        col_codigo_pac_uni,'Escol_num_gest',"ID_MN_RESI_gest",
        "ID_MUNICIP_gest","DT_NOTIFIC_gest"]], on = col_codigo_pac_uni, how = "left")
    
    DF_ult_ano = pd.merge(DF_ult_ano, Sinan_cr_Prim[[
        col_codigo_pac_uni,'Escol_num_cr',"ID_MN_RESI_cr",
        "ID_MUNICIP_cr","DT_NOTIFIC_cr"]], on = col_codigo_pac_uni, how = "left")

    DF_ult_ano = pd.merge(DF_ult_ano, Sinan_cong[[
        col_codigo_pac_uni,'Escol_num_cong',"ID_MN_RESI_cong",
        "ID_MUNICIP_cong","DT_NOTIFIC_cong"]], on = col_codigo_pac_uni, how = "left")

    print("Base DF_ult_ano após o Merge")
    display(DF_ult_ano[col_ano].value_counts().sort_index(ascending = False))
    print()

    del Disp_ano, CV_ano, CD4_ano, PVHA_ano, PVHA, Disp, CV, CD4, Sinan_Prim, Sinan_G_Prim, Sinan_cr_Prim, Sinan_cong
    gc.collect()


    ######################## PVHA_PRIM ###########################################

    ### Organização de Variáveis
    DF['Escol_num'] = DF[['Escol_num','Escol_num_Ad','Escol_num_gest','Escol_num_cr',"Escol_num_cong"]].max(axis=1)

    Escol = [
        (DF["Escol_num"] == 1),
        (DF["Escol_num"] == 2),
        (DF["Escol_num"] == 3)
    ]

    Escol_escolha = ["0 a 7 anos",
                    "8 a 11 anos",
                    "12 e mais anos"]

    DF["Escol_cat"] = np.select(Escol, Escol_escolha, default="Não Informado")
    display(DF["Escol_cat"].value_counts(dropna = False))
    print()

    # Unifica o codigo_ibge_resid
    DF = fg.ibge_resid(DF, cols_ibge = [
                        "ID_MN_RESI_Ad","ID_MN_RESI_gest","ID_MN_RESI_cr","ID_MN_RESI_cong",
                        "cod_ibge_solicitante_cv","cod_ibge_solicitante_cd4","cod_ibge_udm",
                        "ID_MUNICIP_Ad","ID_MUNICIP_gest", "ID_MUNICIP_cr","ID_MUNICIP_cong",
                        "CODMUNRES","CODMUNOCOR","CODMUNCART"
                        ])

    # Organiza o IBGE das instituições Solicitantes
    DF = fg.ibge_inst_sol_exames(DF)

    colunas_datamin = [
        "data_hora_coleta_cd4",
        "data_hora_coleta_cv",
        "data_solicitacao_cd4",
        "data_solicitacao_cv"
    ]

    # Encontra a menor data em um vetor de colunas
    DF = fg.encontrar_menor_data_vetorizado(DF, 
                                            colunas_datas = colunas_datamin,
                                            nome_dt_min = "data_min_exame")
    DF["ano_min_exame"] = DF["data_min_exame"].dt.year

    ### Criação de variáveis para o PIMC

    # Cria a coluna com a quantidade de dias entre o diagnóstico e a primeira dispensação
    DF["Dias_diag_TARV"] = (DF["data_dispensa"] - DF["data_min"]).dt.days

    cond = [
        DF["data_dispensa"].isnull() == True,
        DF["Dias_diag_TARV"] < 1,
        DF["Dias_diag_TARV"] < 8,
        DF["Dias_diag_TARV"] < 30,
        DF["Dias_diag_TARV"] < 180,
        DF["Dias_diag_TARV"] >= 180
    ]

    retorno = [
        "Não iniciou TARV",
        "No mesmo dia",
        "1-7 dias",
        "8-30 dias",
        "1-6 meses",
        "Mais de 6 meses"  
    ]

    DF["Dias_diag_TARV_cat"] = np.select(cond,retorno, default = "Erro")
    display("Tempo entre diagnóstico e TARV ",DF[["ano_min","Dias_diag_TARV_cat"]].value_counts(dropna = False).sort_index(ascending = False))


    # Cria a coluna com a quantidade de dias entre o exame e a primeira dispensação
    DF["Dias_exame_TARV"] = (DF["data_dispensa"] - DF["data_min_exame"]).dt.days

    cond = [
        DF["data_dispensa"].isnull() == True,
        DF["data_min_exame"].isnull() == True,
        DF["Dias_exame_TARV"] < -15,
        DF["Dias_exame_TARV"] < 1,
        DF["Dias_exame_TARV"] < 8,
        DF["Dias_exame_TARV"] < 30,
        DF["Dias_exame_TARV"] < 180,
        DF["Dias_exame_TARV"] >= 180
    ]

    retorno = [
        "Não iniciou TARV",
        "Não tem exames disponíveis",
        "Iniciou TARV antes de exames no SUS",
        "No mesmo dia",
        "1-7 dias",
        "8-30 dias",
        "1-6 meses",
        "Mais de 6 meses"  
    ]

    DF["Dias_exame_TARV_cat"] = np.select(cond,retorno, default = "Erro")
    display("Tempo entre 1o exame e TARV ",DF[["ano_min","Dias_exame_TARV_cat"]].value_counts(dropna = False).sort_index(ascending = False))

    mask = DF["Nome_mun_exame"].isnull()
    DF.loc[mask, "ibge_inst_exames"] = DF.loc[mask, "cod_ibge6_exame"].apply(
        lambda x: str(x)[:2] + "99999" if pd.notna(x) else "9999999")
    

    # Classifica os Municípios com menos de 50 pessoas por ano em uma categoria unificada

    # Define os 10 anos mais recentes
    anos_recentes = sorted(DF["ano_min"].dropna().unique())[-10:-1]

    # Conta quantas pessoas por (município, ano)
    contagem = (
        DF[DF["ano_min"].isin(anos_recentes)]
        .groupby(["ibge_inst_exames", "ano_min"])[col_codigo_pac_uni]
        .count()
        .reset_index(name="n_pessoas")
    )

    # Todos os municípios (inclusive nulos) nos dados originais
    todos_municipios = DF["ibge_inst_exames"].unique()

    # Municípios que aparecem na contagem (com pelo menos 1 pessoa em algum ano)
    municipios_contados = contagem["ibge_inst_exames"].unique()

    # Municípios completamente ausentes da contagem (possivelmente 0 pessoas ou dados ausentes)
    municipios_ausentes = set(todos_municipios) - set(municipios_contados)

    # Seleciona os municípios com < 30 pessoas em ao menos um desses anos
    municipios_menos_30 = np.concatenate([
        contagem[contagem["n_pessoas"] < 30]["ibge_inst_exames"].unique(),
        list(municipios_ausentes)
    ])

    # Cria a condição booleana
    condicao_mun = DF["ibge_inst_exames"].astype(int).isin(pd.to_numeric(municipios_menos_30).astype(int))

    DF["Nome_mun_exame_30+"] = np.where(condicao_mun,
                                            "Município com menos de 30 pessoas identificadas por ano",
                                            DF["Nome_mun_exame"])

    DF["Cod_ibge_insti_sol_30+"] = np.where(condicao_mun,
                                        DF["ibge_inst_exames"].apply(lambda x: str(x)[:2] + "99999"),
                                        DF["ibge_inst_exames"])

    DF["Cod_ibge_insti_sol_30+"] = pd.to_numeric(DF["Cod_ibge_insti_sol_30+"], errors = "coerce").astype("Int64")


    # Classifica os Serviços com menos de 30 pessoas por ano em uma categoria unificada

    # Conta quantas pessoas por (instituicao, município, ano)
    contagem = (
        DF[DF["ano_min"].isin(anos_recentes)]
        .groupby(["nome_inst_exames","ibge_inst_exames", "ano_min"])[col_codigo_pac_uni]
        .count()
        .reset_index(name="n_pessoas_inst")
    )

    # Todas as combinações presentes no DF, inclusive com nulos
    todas_insts = DF[["nome_inst_exames", "ibge_inst_exames"]].drop_duplicates()

    # Combinações que apareceram na contagem (têm pelo menos uma pessoa em algum ano)
    insts_com_contagem = contagem[["nome_inst_exames", "ibge_inst_exames"]].drop_duplicates()

    # Força os tipos das colunas para string
    todas_insts["nome_inst_exames"] = todas_insts["nome_inst_exames"].astype(str)
    todas_insts["ibge_inst_exames"] = todas_insts["ibge_inst_exames"].astype(int)

    insts_com_contagem["nome_inst_exames"] = insts_com_contagem["nome_inst_exames"].astype(str)
    insts_com_contagem["ibge_inst_exames"] = insts_com_contagem["ibge_inst_exames"].astype(int)

    # Usa merge anti-join para pegar as que estão ausentes da contagem
    insts_ausentes = todas_insts.merge(insts_com_contagem, 
                                    on=["nome_inst_exames", "ibge_inst_exames"],
                                    how="left", indicator=True)
    insts_ausentes = insts_ausentes[insts_ausentes["_merge"] == "left_only"].drop(columns=["_merge"])

    insts_menos_30_ou_ausentes = pd.concat([
        contagem[contagem["n_pessoas_inst"] < 30][["nome_inst_exames", "ibge_inst_exames"]].drop_duplicates(),
        insts_ausentes
    ]).drop_duplicates()

    # Transforma em tuplas
    insts_menos_30_ou_ausentes_tuplas = list(insts_menos_30_ou_ausentes.itertuples(index=False, name=None))

    # Cria a chave de combinação nos dois DataFrames
    DF["chave_inst"] = list(zip(DF["nome_inst_exames"], DF["ibge_inst_exames"]))

    # Cria a condição booleana com base nas tuplas
    condicao_inst = DF["chave_inst"].isin(insts_menos_30_ou_ausentes_tuplas)

    DF["nome_inst_exames_30+"] = np.where(condicao_inst,
                                            "Serviços com menos de 30 pessoas identificadas por ano",
                                            DF["nome_inst_exames"])
    
    DF.drop(columns=["chave_inst"], inplace=True)


    ######################## PVHA_ULT_ANO ###########################################

    # Preenche as colunas do ano que estavam vazias.
    DF_ult_ano["ano_min"] = DF_ult_ano["ano_min"].fillna(np.nan)
    DF_ult_ano["ano_min"] = DF_ult_ano["ano_min"].astype("Int64")
    DF_ult_ano["ano_obito"] = DF_ult_ano["ano_obito"].fillna(np.nan)
    DF_ult_ano["ano_obito"] = DF_ult_ano["ano_obito"].astype("Int64")
    DF_ult_ano[col_ano_disp] = DF_ult_ano[col_ano_disp].fillna(np.nan)
    DF_ult_ano[col_ano_disp] = DF_ult_ano[col_ano_disp].astype("Int64")
    DF_ult_ano[col_ano_cv] = DF_ult_ano[col_ano_cv].fillna(np.nan)
    DF_ult_ano[col_ano_cv] = DF_ult_ano[col_ano_cv].astype("Int64")
    DF_ult_ano[col_ano_cd4] = DF_ult_ano[col_ano_cd4].fillna(np.nan)
    DF_ult_ano[col_ano_cd4] = DF_ult_ano[col_ano_cd4].astype("Int64")

    # Filtra excluindo: linhas de anos a partir do ano do óbito.
    DF_ult_ano = DF_ult_ano[((DF_ult_ano["ano_obito"].isnull()) |
                            (DF_ult_ano[col_ano] < DF_ult_ano["ano_obito"]))].copy()
    
    print("Base DF_ult_ano após o excluir os óbitos")
    display(DF_ult_ano[col_ano].value_counts().sort_index(ascending = False))
    print()

    DF_ult_ano['Escol_num'] = DF_ult_ano[['Escol_num','Escol_num_Ad','Escol_num_gest','Escol_num_cr','Escol_num_cong']].max(axis=1)

    Escol = [
        (DF_ult_ano["Escol_num"] == 1),
        (DF_ult_ano["Escol_num"] == 2),
        (DF_ult_ano["Escol_num"] == 3)
    ]

    Escol_escolha = ["0 a 7 anos",
                    "8 a 11 anos",
                    "12 e mais anos"]

    DF_ult_ano["Escol_cat"] = np.select(Escol, Escol_escolha, default="Não Informado")
    display(DF_ult_ano["Escol_cat"].value_counts(dropna = False))
    print()

    # Calcula a idade em todas as datas passadas no cols_idade.
    for cols in cols_datas:
        DF_ult_ano[f"Idade_{cols}_cat"] = fg.idade_cat(DF = DF_ult_ano, data_ref=cols, data_nasc=col_data_nascimento)
        display(DF_ult_ano[f"Idade_{cols}_cat"].value_counts().sort_index())
        print()


    DF_ult_ano = fg.encontrar_menor_data_vetorizado(df=DF_ult_ano, colunas_datas=colunas_data_min_exame, nome_dt_min = "data_min_exame")
    DF_ult_ano["data_min_exame"] = pd.to_datetime(DF_ult_ano["data_min_exame"],errors="coerce").dt.normalize()
    DF_ult_ano["ano_min_exame"]= DF_ult_ano["data_min_exame"].dt.year
    display(DF_ult_ano["ano_min_exame"].value_counts(dropna = False).sort_index())
    print()

    # Crie uma série de datas com o último dia do ano
    DF_ult_ano["Ult_dia_ano"] = pd.to_datetime(DF_ult_ano[col_ano].astype(str) + "-12-31")

    # Ajusta o último dia do ano para o dia de fechamento da base no ano atual
    DF_ult_ano.loc[DF_ult_ano[col_ano] == ano, "Ult_dia_ano"] = hoje

    # Normalize as datas para garantir que apenas a data seja mantida (sem informações de hora)
    DF_ult_ano["Ult_dia_ano"] = DF_ult_ano["Ult_dia_ano"].dt.normalize()

    # Calcula a idade no último dia de cada ano (importante para Cascata e BI)
    DF_ult_ano["Idade_cascata"] = fg.idade_cat(DF = DF_ult_ano, data_ref="Ult_dia_ano", data_nasc=col_data_nascimento)
    display(DF_ult_ano["Idade_cascata"].value_counts().sort_index())
    print()

    # Calcula a idade no último dia de cada ano nas categorias do spectrum
    DF_ult_ano["Idade_Spec_crianca"] = fg.idade_cat(DF = DF_ult_ano, data_ref="Ult_dia_ano", data_nasc=col_data_nascimento, faixas_etarias="spectrum_crianca")
    display(DF_ult_ano["Idade_Spec_crianca"].value_counts().sort_index())
    print()   
    DF_ult_ano["Idade_Spec1"] = fg.idade_cat(DF = DF_ult_ano, data_ref="Ult_dia_ano", data_nasc=col_data_nascimento, faixas_etarias="spectrum1")
    display(DF_ult_ano["Idade_Spec1"].value_counts().sort_index())
    print()   
    DF_ult_ano["Idade_Spec2"] = fg.idade_cat(DF = DF_ult_ano, data_ref="Ult_dia_ano", data_nasc=col_data_nascimento, faixas_etarias="spectrum2")
    display(DF_ult_ano["Idade_Spec2"].value_counts().sort_index())
    print()   

    data_disp = next((col for col in cols_datas if "disp" in col.lower()), None)

    data_cv = next((col for col in cols_datas if "cv" in col.lower() or "carga" in col.lower()), None)

    data_cd4 = next((col for col in cols_datas if "cd4" in col.lower()), None)

        
    # Organiza o IBGE de residência
    DF_ult_ano = fg.ibge_resid(DF=DF_ult_ano, cols_ibge = [
                        "ID_MN_RESI_Ad","ID_MN_RESI_gest","ID_MN_RESI_cr", "ID_MN_RESI_cong",
                        "cod_ibge_solicitante_cv","cod_ibge_solicitante_cd4","cod_ibge_udm",
                        "ID_MUNICIP_Ad","ID_MUNICIP_gest", "ID_MUNICIP_cr", "ID_MUNICIP_cong"
                        ])

    # Organiza o IBGE das instituições Solicitantes
    DF_ult_ano = fg.ibge_inst_sol_exames(DF_ult_ano)

    ## A Coluna "Status_ano" mostra a situação do paciente em relação ao último dia do ano em que a dispensação foi feita, ou em relação a data atual
    # Calcula o atraso em relação ao último dia do ano da dispensação
    DF_ult_ano["Atraso_ano"] = (DF_ult_ano["Ult_dia_ano"] - DF_ult_ano["dt_lim_prox_disp"]).dt.days

    cond_status = [
        (DF_ult_ano["Atraso_ano"] <= 60),
        (DF_ult_ano["Atraso_ano"] > 60),
        (DF_ult_ano["data_min_exame"].isnull() == False)
        ]

    escolha_status = ["Tarv", "Interrupção de Tarv", "Gap de tratamento"]

    DF_ult_ano["Status_ano"] = np.select(cond_status, escolha_status, default= None)

    DF_ult_ano["Status_ano"] = DF_ult_ano["Status_ano"].fillna("Gap de vinculação")

    display("Status no último dia do ano, para o ano vigente, considera-se a data de fechamento do banco",
            pd.DataFrame(DF_ult_ano[[col_ano,"Status_ano"]].value_counts(dropna = False).sort_index(ascending = False)))
    print()

    # Cria uma variável com o vínculo no ano
    Cond_vin = [
        (DF_ult_ano[col_ano] == ano) &
        (((DF_ult_ano["Hoje"] - DF_ult_ano[data_cv]).dt.days <= 365) |
        ((DF_ult_ano["Hoje"] - DF_ult_ano[data_cd4]).dt.days <= 365) |
        ((DF_ult_ano["Hoje"] - DF_ult_ano[data_disp]).dt.days <= 365)),
        
        ((DF_ult_ano[col_ano] == DF_ult_ano[data_cv].dt.year) |
        (DF_ult_ano[col_ano] == DF_ult_ano[data_cd4].dt.year) |
        (DF_ult_ano[col_ano] == DF_ult_ano[data_disp].dt.year)),
    ]

    escolha_vin = [
        "Vinculado", "Vinculado"
    ]

    DF_ult_ano["Vinculado_ano"]  = np.select(Cond_vin,escolha_vin, default = "Não vinculado")
    display("Vinculados por ano, considera os últimos 12 meses para o ano vigente",
            pd.DataFrame(DF_ult_ano[["ano","Vinculado_ano"]].value_counts(dropna = False).sort_index(ascending = False)))
    print()

    # Conta o número de pessoas que fazem acompanhamento (exame CV) em cada instituição.
    DF_ult_ano["N_PVHA_CV_ano_Instituicao"] = DF_ult_ano[DF_ult_ano["nome_inst_exames"] != "SEM EXAMES NO SUS"].groupby(["ano","nome_inst_exames"])[col_codigo_pac_uni].transform("count")

    cond_inst = [
            (DF_ult_ano["N_PVHA_CV_ano_Instituicao"] >= 1000),
            (DF_ult_ano["N_PVHA_CV_ano_Instituicao"] >= 500),
            (DF_ult_ano["N_PVHA_CV_ano_Instituicao"] >= 1)
    ]

    escolhas_inst = [
            "Instituições acompanhando mais de 1000 pessoas",
            "Instituições acompanhando entre 500 e 1000 pessoas",
            "Instituições acompanhando menos de 500 pessoas"
    ]

    DF_ult_ano["Grupo_instituicao"] = np.select(cond_inst,escolhas_inst, default = "Sem pedidos de exames")
    display(DF_ult_ano["Grupo_instituicao"].value_counts())
    print()

    # Categoriza as faixas de Carga Viral
    DF_ult_ano[col_copias] = pd.to_numeric(DF_ult_ano[col_copias], errors='coerce')
    DF_ult_ano[col_comentario_copias] = pd.to_numeric(DF_ult_ano[col_comentario_copias], errors='coerce')

    cond_cv = [
            (((DF_ult_ano[col_ano_cv] == DF_ult_ano[col_ano]) |
            ((DF_ult_ano[col_ano] == ano) & ((pd.to_datetime(hoje) - DF_ult_ano[data_cv]).dt.days <= 365))) &
            ((DF_ult_ano[col_copias].isnull() == True) & (DF_ult_ano[col_comentario_copias] == 0)) |
            ((DF_ult_ano[col_copias].isnull() == True) & (DF_ult_ano[col_comentario_copias] == 1)) |
            (DF_ult_ano[col_copias] < 50)
            ),
            (
            ((DF_ult_ano[col_ano_cv] == DF_ult_ano[col_ano]) |
            ((DF_ult_ano[col_ano] == ano) & ((pd.to_datetime(hoje) - DF_ult_ano[data_cv]).dt.days <= 365))) &
            (((DF_ult_ano[col_copias].isnull() == True) & (DF_ult_ano[col_comentario_copias] == 2)) |
            (DF_ult_ano[col_copias] >= 10000))
            ),
            (
            ((DF_ult_ano[col_ano_cv] == DF_ult_ano[col_ano]) |
            ((DF_ult_ano[col_ano] == ano) & ((pd.to_datetime(hoje) - DF_ult_ano[data_cv]).dt.days <= 365))) &
            (DF_ult_ano[col_copias] >= 50) & (DF_ult_ano[col_copias] < 200)
            ),
            (
            ((DF_ult_ano[col_ano_cv] == DF_ult_ano[col_ano]) |
            ((DF_ult_ano[col_ano] == ano) & ((pd.to_datetime(hoje) - DF_ult_ano[data_cv]).dt.days <= 365))) &
            (DF_ult_ano[col_copias] >= 200) & (DF_ult_ano[col_copias] < 1000)
            ),
            (
            ((DF_ult_ano[col_ano_cv] == DF_ult_ano[col_ano]) |
            ((DF_ult_ano[col_ano] == ano) & ((pd.to_datetime(hoje) - DF_ult_ano[data_cv]).dt.days <= 365))) &        
            (DF_ult_ano[col_copias] >= 1000) & (DF_ult_ano[col_copias] < 10000)
            )
    ]

    escolhas_cv = [
            "<50","10.000+","50-199","200-999","1.000-9.999"
    ]

    DF_ult_ano["CV_cat"] = np.select(cond_cv,escolhas_cv, default = "Não disponível")
    display(DF_ult_ano["CV_cat"].value_counts())
    print()

    # Cria a categoria com 1000 copias de CV   
    cond_cv = [
            (DF_ult_ano["CV_cat"] == "1.000-9.999") | (DF_ult_ano["CV_cat"] == "10.000+"),
            (DF_ult_ano["CV_cat"] == "<50") | (DF_ult_ano["CV_cat"] == "50-199") | (DF_ult_ano["CV_cat"] == "200-999")   
    ]

    escolhas_cv = [
            ">=1.000",
            "<1.000"
    ]

    DF_ult_ano["CV_cat1000"] = np.select(cond_cv,escolhas_cv, default = "Não disponível")
    display(DF_ult_ano["CV_cat1000"].value_counts())
    print()

    # Cria a categoria com 200 copias de CV
    cond_cv = [
            (DF_ult_ano["CV_cat"] == "<50") | (DF_ult_ano["CV_cat"] == "50-199"),
            (DF_ult_ano["CV_cat"] == "1.000-9.999") | (DF_ult_ano["CV_cat"] == "10.000+") |
            (DF_ult_ano["CV_cat"] == "200-999")   
    ]

    escolhas_cv = [
            "<200",
            ">=200"
    ]

    DF_ult_ano["CV_cat200"] = np.select(cond_cv,escolhas_cv, default = "Não disponível")
    display(DF_ult_ano["CV_cat200"].value_counts())
    print()

    # Cria a categoria com 50 copias de CV   
    cond_cv = [
            (DF_ult_ano["CV_cat"] == "<50"),
            (DF_ult_ano["CV_cat"] == "1.000-9.999") | (DF_ult_ano["CV_cat"] == "10.000+") |
            (DF_ult_ano["CV_cat"] == "50-199") | (DF_ult_ano["CV_cat"] == "200-999")   
    ]

    escolhas_cv = [
            "<50",
            ">=50"
    ]

    DF_ult_ano["CV_cat50"] = np.select(cond_cv,escolhas_cv, default = "Não disponível")
    display(DF_ult_ano["CV_cat50"].value_counts())
    print()

    # Cria uma variável com a indetecção viral em 50 cópias
    Cond_sup = [
        (
        (DF_ult_ano["Status_ano"] == "Tarv") & 
        ((DF_ult_ano[data_cv] - DF_ult_ano[data_disp_prim]).dt.days > 180) &
        (DF_ult_ano["CV_cat50"] == "<50")
        ),
        (
        (DF_ult_ano["Status_ano"] == "Tarv") & 
        ((DF_ult_ano[data_cv] - DF_ult_ano[data_disp_prim]).dt.days > 180) &
        (DF_ult_ano["CV_cat50"] == ">=50")
        )
    ]

    escolha_sup = [
    "CV indetectável","CV detectável"
    ]

    DF_ult_ano["Deteccao_CV50"]  = np.select(Cond_sup,escolha_sup, default = "Sem dados de CV")
    display("Indetecção Viral pós TARV na última CV do ano (<50 cp)",
            DF_ult_ano[[col_ano,"Deteccao_CV50"]].value_counts(dropna = False).sort_index())

    # Cria uma variável com a supressão viral em 200 cópias
    Cond_sup = [
        (
        (DF_ult_ano["Status_ano"] == "Tarv") & 
        ((DF_ult_ano[data_cv] - DF_ult_ano[data_disp_prim]).dt.days > 180) &
        (DF_ult_ano["CV_cat200"] == "<200")
        ),
        (
        (DF_ult_ano["Status_ano"] == "Tarv") & 
        ((DF_ult_ano[data_cv] - DF_ult_ano[data_disp_prim]).dt.days > 180) &
        (DF_ult_ano["CV_cat200"] == ">=200")
        )
    ]

    escolha_sup = [
        "CV suprimida","CV não suprimida"
    ]

    DF_ult_ano["Deteccao_CV200"]  = np.select(Cond_sup,escolha_sup, default = "Sem dados de CV")
    display("Supressão Viral pós TARV na última CV do ano (<200 cp)",
            DF_ult_ano[[col_ano,"Deteccao_CV200"]].value_counts(dropna = False).sort_index())


    # Cria uma variável com a supressão viral em 1000 cópias
    Cond_sup = [
        (
        (DF_ult_ano["Status_ano"] == "Tarv") & 
        ((DF_ult_ano[data_cv] - DF_ult_ano[data_disp_prim]).dt.days > 180) &
        (DF_ult_ano["CV_cat1000"] == "<1.000")
        ),
        (
        (DF_ult_ano["Status_ano"] == "Tarv") & 
        ((DF_ult_ano[data_cv] - DF_ult_ano[data_disp_prim]).dt.days > 180) &
        (DF_ult_ano["CV_cat1000"] == ">=1.000")
        )
    ]

    escolha_sup = [
        "CV suprimida","CV não suprimida"
    ]

    DF_ult_ano["Deteccao_CV1000"]  = np.select(Cond_sup,escolha_sup, default = "Sem dados de CV")
    display("Supressão Viral pós TARV na última CV do ano (<1000 cp)",
            DF_ult_ano[[col_ano,"Deteccao_CV1000"]].value_counts(dropna = False).sort_index())

    DF_ult_ano[col_contagem_cd4] = pd.to_numeric(DF_ult_ano[col_contagem_cd4], errors="coerce")

    # criando variável de categorias do valor do CD4
    Cond_CD4cat = [
        (DF_ult_ano[col_contagem_cd4] > 0) & (DF_ult_ano[col_contagem_cd4] < 200),
        (DF_ult_ano[col_contagem_cd4] >= 200) & (DF_ult_ano[col_contagem_cd4] < 350),
        (DF_ult_ano[col_contagem_cd4] >= 350) & (DF_ult_ano[col_contagem_cd4] < 500),
        (DF_ult_ano[col_contagem_cd4] >= 500) & (DF_ult_ano[col_contagem_cd4] < 2500)
        ]

    Escolha_CD4cat = ["0-199", "200-349", "350-499", "500+"]

    DF_ult_ano["CD4_cat"] = np.select(Cond_CD4cat, Escolha_CD4cat, default="Sem resultado de CD4")

    display(DF_ult_ano["CD4_cat"].value_counts(dropna = False).sort_index())
    print()

    # criando variável de categorias do valor do CD4
    Cond_CD4cat2 = [
        (DF_ult_ano[col_contagem_cd4] > 0) & (DF_ult_ano[col_contagem_cd4] < 200),
        (DF_ult_ano[col_contagem_cd4] >= 200) & (DF_ult_ano[col_contagem_cd4] < 2500)
        ]

    Escolha_CD4cat2 = ["<200 células / mm³", "≥200 células / mm³"]

    DF_ult_ano["CD4_cat200"] = np.select(Cond_CD4cat2, Escolha_CD4cat2, default="Sem resultado de CD4")

    display(DF_ult_ano["CD4_cat200"].value_counts(dropna = False).sort_index())
    print()

    # Calcular o último ano de registro de evento de cada pessoa
    ultimo_evento = DF_ult_ano.groupby(col_codigo_pac_uni).agg({
        "ano_min":"max",
        col_ano_disp: 'max',
        col_ano_cv: 'max',
        col_ano_cd4: 'max'
    }).max(axis=1).reset_index()
    ultimo_evento.columns = [col_codigo_pac_uni, 'ultimo_evento']

    ultimo_vinculo = DF_ult_ano.groupby(col_codigo_pac_uni).agg({
        col_ano_disp: 'max',
        col_ano_cv: 'max',
        col_ano_cd4: 'max'
    }).max(axis=1).reset_index()
    ultimo_vinculo.columns = [col_codigo_pac_uni, 'ultimo_vinculo']

    # Mesclar essa informação de último evento com o DataFrame original
    DF_ult_ano = DF_ult_ano.merge(ultimo_evento, on=col_codigo_pac_uni)
    DF_ult_ano = DF_ult_ano.merge(ultimo_vinculo, on=col_codigo_pac_uni)

    # Retira os óbitos presumidos (pessoas há mais de 10 anos sem aparecer na base)
    if obito_admin is True:

        # Filtrar os dados onde a diferença entre o ano atual e o último evento é menor ou igual a 10 anos
        DF_ult_ano = DF_ult_ano[DF_ult_ano[col_ano] - DF_ult_ano['ultimo_evento'] <= anos_obito_adm]

    print("Base DF_ult_ano após o excluir os óbitos ADM")
    display(DF_ult_ano[col_ano].value_counts().sort_index(ascending = False))
    print()

    # Adiciona nova coluna para os vinculados nos últimos 5 anos
    DF_ult_ano['Vinculado_5anos'] = np.where((DF_ult_ano['ultimo_vinculo'].notna()) &
                                            (DF_ult_ano["ano"] - DF_ult_ano['ultimo_vinculo'] <= 5), "Vinculado", "Não Vinculado")

    # Adiciona nova coluna para os vinculados nos últimos 5 anos
    DF_ult_ano['Vinculado_10anos'] = np.where((DF_ult_ano['ultimo_vinculo'].notna()) &
                                            (DF_ult_ano["ano"] - DF_ult_ano['ultimo_vinculo'] <= 10), "Vinculado", "Não Vinculado")

    # Filtra o ano de corte 
    DF_ult_ano = DF_ult_ano[DF_ult_ano[col_ano] >= ano_corte]

    print("Base DF_ult_ano após o excluir os anos de corte")
    display(DF_ult_ano[col_ano].value_counts().sort_index(ascending = False))
    print()

    # Filtra os cadastros inválidos do Siclom
    DF_ult_ano = DF_ult_ano[
        ~((DF_ult_ano["st_num"].isin([2,3])) &
          (DF_ult_ano["Status_ano"].isin(["Interrupção de Tarv","Gap de tratamento"])) &
          (DF_ult_ano["ano"] - DF_ult_ano['ultimo_vinculo'] >= 2))
        ]

    print("Base DF_ult_ano após o excluir os cadastros inválidos do SICLOM")
    display(DF_ult_ano[col_ano].value_counts().sort_index(ascending = False))
    print()

    # Classifica os Municípios com menos de 50 pessoas por ano em uma categoria unificada
    # Conta quantas pessoas por (município, ano)
    contagem = (
        DF_ult_ano
        .groupby(["ibge_inst_exames", "ano"])[col_codigo_pac_uni]
        .count()
        .reset_index(name="n_pessoas")
    )

    # Seleciona os municípios com < 50 pessoas em ao menos um desses anos
    municipios_menos_50 = contagem[contagem["n_pessoas"] < 50]["ibge_inst_exames"].unique()

    # Cria a condição booleana
    condicao_mun = DF_ult_ano["ibge_inst_exames"].isin(municipios_menos_50)


    DF_ult_ano["Nome_mun_exame_50+"] = np.where(condicao_mun,
                                            "Município com menos de 50 pessoas no cuidado",
                                            DF_ult_ano["Nome_mun_exame"])

    DF_ult_ano["Cod_ibge_insti_sol_50+"] = np.where(condicao_mun,
                                        (DF_ult_ano["ibge_inst_exames"].apply(lambda x: str(x)[:2] + "99999")),
                                        DF_ult_ano["ibge_inst_exames"])

    DF_ult_ano["Cod_ibge_insti_sol_50+"] = pd.to_numeric(DF_ult_ano["Cod_ibge_insti_sol_50+"], errors = "coerce").astype("Int64")

    # Classifica os Serviços com menos de 50 pessoas por ano em uma categoria unificada
    # Conta quantas pessoas por (instituicao, município, ano)
    contagem = (
        DF_ult_ano
        .groupby(["nome_inst_exames","ibge_inst_exames", "ano"])[col_codigo_pac_uni]
        .count()
        .reset_index(name="n_pessoas_inst")
    )

    # Seleciona os municípios com < 50 pessoas em ao menos um desses anos
    instituicao_menos_50_tuplas = list(contagem[contagem["n_pessoas_inst"] < 50][["nome_inst_exames","ibge_inst_exames"]]
                            .drop_duplicates()
                            .itertuples(index=False, name=None))

    # Cria a chave de combinação nos dois DataFrames
    DF_ult_ano["chave_inst"] = list(zip(DF_ult_ano["nome_inst_exames"], DF_ult_ano["ibge_inst_exames"]))

    # Cria a condição booleana com base nas tuplas
    condicao_inst = DF_ult_ano["chave_inst"].isin(instituicao_menos_50_tuplas)
    

    DF_ult_ano["nome_inst_exames_50+"] = np.where(condicao_inst,
                                            "Serviços com menos de 50 pessoas no cuidado",
                                            DF_ult_ano["nome_inst_exames"])

    DF_ult_ano.drop(columns=["chave_inst"], inplace=True)



    ######################### GERA O BANCO PVHA_PRIM_ULT ##########################


    cols_drop_prim = [
        "escolaridade",
        "data_ult_atu",
        'dt_lim_prox_disp',
        'dt_real_disp_seguinte',
        'Atraso',
        'Status_cadaDisp',
        'Abandono',
        'Abandono_sum',
        'N_dispensas',
        'Atraso_sum',
        'Abandono_raz',
        'Atraso_raz',
        'Atraso_cat',
        'Abandono_ano',
        'esquemaAMA_Prim',
        'data_dispensa_Prim',
        'ano_disp_Prim',
        'N_exames_ano_cd4',
        'data_coleta_cd4_Prim',
        'contagem_cd4_Prim',
        'data_sol_cd4_Prim',
        'ano_coleta_cd4_Prim',
        'CD4_cat_Prim',
        'CD4_cat200_Prim',
        'data_coleta_cv_Prim',
        'copia_cv_Prim',
        'comentario_copias_Prim',
        'data_sol_cv_Prim',
        'ano_coleta_cv_Prim',
        'CV_cat_Prim'
    ]

    cols_rename_prim = [
        'codigo_udm',
        'nm_udm',
        'tp_servico_atendimento',
        'cod_ibge_udm',
        'st_pub_priv',
        "esquema",
        "esquema_forma",
        "duracao",
        "esquema",
        "mesN_disp",
        'duracao_sum',
        'Duracao_cat',
        'esquema_AMA',
        'Classe_Med',
        'Classe_Med_Count',
        'esquema_cat',
        "data_dispensa",
        'mes_disp',
        'mes_ano',
        "ano_disp",
        'cod_ibge_solicitante_cd4',
        'cd_inst_sol_cd4',
        'nm_inst_sol_cd4',
        'tipo_inst_sol_cd4',
        'data_solicitacao_cd4',
        'data_hora_coleta_cd4',
        'contagem_cd4',
        'ano_coleta_cd4',
        'mesN_coleta_cd4',
        'mes_coleta_cd4',
        'CD4_cat',
        'CD4_cat200',
        'CD4_cat350',
        'cod_ibge_solicitante_cv',
        'cd_inst_sol_cv',
        'nm_inst_sol_cv',
        'tipo_inst_sol_cv',
        'data_solicitacao_cv',
        'data_hora_coleta_cv',
        'copias',
        'comentario_copias',
        'ano_coleta_cv',
        'mesN_coleta_cv',
        'mes_coleta_cv',
        'N_exames_ano_cv',
        'CV_detec_count',
        'CV_5000_count',
        'CV_cat',
        'CV_cat1000',
        'CV_cat50',
        'CV_cat200',
        'CV_cat500',
        'nome_inst_exames',
        'uf_inst',
        'uf_inst_completo',
        'reg_inst',
        'cod_ibge6_exame',
        'ibge_inst_exames',
        'Nome_mun_exame',
        'Populacao_exame',
        'Nome_mun_exame_30+',
        'Cod_ibge_insti_sol_30+',
        'nome_inst_exames_30+',
        "cd_crm",
        "uf_crm"
    ]

    DF = DF.drop(columns = cols_drop_prim)

    for cols in cols_rename_prim:
        DF = DF.rename(columns = {cols:f"{cols}_prim"})

    PVHA_ult = DF_ult_ano[DF_ult_ano["ano"] == ano].copy()

    cols_drop_ult = [
        'ano',
        'data_min',
        'data_obito',
        'data_nascimento',
        'Sexo_cat',
        'Raca_cat',
        'PVHA',
        'Hoje',
        'Raca_cat2',
        'ano_min',
        'ano_obito',
        'Idade_data_min',
        'Idade_vinc_cat',
        'Idade_vinc_cat_anterior',
        'Idade_Hoje',
        'Idade_atual',
        'codigo_paciente',
        'uf_residencia',
        'acomp_medico',
        'escolaridade',
        'data_ult_atu',
        'cd_pais',
        'co_orientacao_sexual',
        'co_genero',
        'Escol_num',
        'Genero_cat',
        'Orientacao_cat',
        'Pop_genero',
        'pais',
        'Status_cadaDisp',
        'Abandono',
        'Abandono_raz',
        'Atraso_raz',
        'Abandono_ano',
        'esquemaAMA_Prim',
        'data_dispensa_Prim',
        'ano_disp_Prim',
        'data_coleta_cd4_Prim',
        'contagem_cd4_Prim',
        'data_sol_cd4_Prim',
        'ano_coleta_cd4_Prim',
        'CD4_cat_Prim',
        'CD4_cat200_Prim',
        'data_coleta_cv_Prim',
        'copia_cv_Prim',
        'comentario_copias_Prim',
        'data_sol_cv_Prim',
        'ano_coleta_cv_Prim',
        'CV_cat_Prim',
        'data_min_exame',
        'origem_data_min_exame',
        'ano_min_exame',
        'Ult_dia_ano',
        'Idade_Ult_dia_ano',
        'cod_ibge6_res',
        'capital_res',
        'uf_res',
        'reg_res',
        'Nome_mun_resid',
        'Populacao_resid',
        'codigo_ibge_resid',
        "Escol_num_Ad",
        'ID_MN_RESI_Ad',
        'ID_MUNICIP_Ad',
        'DT_NOTIFIC_Ad',
        'Escol_num_gest',
        'ID_MN_RESI_gest',
        'ID_MUNICIP_gest',
        'DT_NOTIFIC_gest',
        'Escol_num_cr',
        'ID_MN_RESI_cr',
        'ID_MUNICIP_cr',
        'DT_NOTIFIC_cr',
        "Escol_num_cong",
        'ID_MN_RESI_cong',
        'ID_MUNICIP_cong',
        'DT_NOTIFIC_cong',
        'Escol_cat'        

    ]

    cols_rename_ult = [
        'codigo_udm',
        'nm_udm',
        'tp_servico_atendimento',
        'cod_ibge_udm',
        'st_pub_priv',
        'data_dispensa',
        'esquema',
        'esquema_forma',
        'duracao',
        'ano_disp',
        'mesN_disp',
        'mes_disp',
        'mes_ano',
        'duracao_sum',
        'dt_lim_prox_disp',
        'dt_real_disp_seguinte',
        'Atraso',
        'Abandono_sum',
        'N_dispensas',
        'Atraso_sum',
        'Atraso_cat',
        'Duracao_cat',
        'esquema_AMA',
        'Classe_Med',
        'Classe_Med_Count',
        'esquema_cat',
        'cod_ibge_solicitante_cd4',
        'cd_inst_sol_cd4',
        'nm_inst_sol_cd4',
        'tipo_inst_sol_cd4',
        'data_solicitacao_cd4',
        'data_hora_coleta_cd4',
        'contagem_cd4',
        'ano_coleta_cd4',
        'mesN_coleta_cd4',
        'mes_coleta_cd4',
        'N_exames_ano_cd4',
        'CD4_cat',
        'CD4_cat200',
        'CD4_cat350',
        'cod_ibge_solicitante_cv',
        'cd_inst_sol_cv',
        'nm_inst_sol_cv',
        'tipo_inst_sol_cv',
        'data_solicitacao_cv',
        'data_hora_coleta_cv',
        'copias',
        'comentario_copias',
        'ano_coleta_cv',
        'mesN_coleta_cv',
        'mes_coleta_cv',
        'N_exames_ano_cv',
        'CV_detec_count',
        'CV_5000_count',
        'CV_cat',
        'CV_cat1000',
        'CV_cat50',
        'CV_cat200',
        'CV_cat500',
        'Idade_data_dispensa',
        'Idade_data_dispensa_cat',
        'Idade_data_hora_coleta_cv',
        'Idade_data_hora_coleta_cv_cat',
        'Idade_data_hora_coleta_cd4',
        'Idade_data_hora_coleta_cd4_cat',   
        'Idade_cascata',
        'Idade_Spec_crianca',
        'Idade_Spec1',
        'Idade_Spec2',
        'nome_inst_exames',
        'uf_inst',
        'uf_inst_completo',
        'reg_inst',
        'cod_ibge6_exame',
        'ibge_inst_exames',
        'Nome_mun_exame',
        'Populacao_exame',
        'Atraso_ano',
        'Status_ano',
        'Vinculado_ano',
        'N_PVHA_CV_ano_Instituicao',
        'Grupo_instituicao',
        'Deteccao_CV50',
        'Deteccao_CV200',
        'Deteccao_CV1000',
        'Vinculado_5anos',
        'Vinculado_10anos',
        'Nome_mun_exame_50+',
        'Cod_ibge_insti_sol_50+',
        'nome_inst_exames_50+',
        "cd_crm",
        "uf_crm"
    ]

    PVHA_ult = PVHA_ult.drop(columns = cols_drop_ult)

    for cols in cols_rename_ult:
        PVHA_ult = PVHA_ult.rename(columns = {cols:f"{cols}_ult"})

    PVHA_prim_ult = pd.merge(DF, PVHA_ult, on = "Cod_unificado", how = "left")

    ### Atualiza o Status atual de todas as colunas novas

    display(PVHA_prim_ult["Status_ano_ult"].value_counts(dropna = False))
    print()

    cond = [
        PVHA_prim_ult["ano_obito"].isnull() == False,
        (PVHA_prim_ult["Status_ano_ult"].isnull()) & (PVHA_prim_ult["ano_disp_prim"].isnull()) & (PVHA_prim_ult["ano_min_exame"].isnull()),
        PVHA_prim_ult["Status_ano_ult"].isnull()
    ]

    escolha = [
        "Óbito",
        "Óbito_adm_sinan",
        "Óbito_adm"
    ]


    PVHA_prim_ult["Status_atual_completo"] = np.select(cond,escolha, default = PVHA_prim_ult["Status_ano_ult"])
    display(PVHA_prim_ult["Status_atual_completo"].value_counts(dropna = False))
    print()

    display(PVHA_prim_ult["Vinculado_10anos_ult"].value_counts(dropna = False))
    print()

    cond = (
        (PVHA_prim_ult["Vinculado_10anos_ult"] == "Vinculado") | (PVHA_prim_ult["Status_atual_completo"] == "Gap de vinculação")
    )

    PVHA_prim_ult["Diagnosticados_vivos"] = np.where(cond, "Diagnosticados", "Óbito")
    display(PVHA_prim_ult["Diagnosticados_vivos"].value_counts(dropna = False))
    print()


    escolhas_cd4 = (
        ((PVHA_prim_ult["data_hora_coleta_cd4_prim"] - PVHA_prim_ult["data_min"]).dt.days <=180) &
        (PVHA_prim_ult["CD4_cat_prim"] != "Sem resultado de CD4")
    )
    PVHA_prim_ult["CD4_diag_prim_cat"] = np.where(escolhas_cd4, PVHA_prim_ult["CD4_cat_prim"], "Sem CD4 próximo ao diagnóstico")
    display(PVHA_prim_ult["CD4_diag_prim_cat"].value_counts(dropna = False))
    print()
    PVHA_prim_ult["CD4_diag_prim_cat200"] = np.where(escolhas_cd4, PVHA_prim_ult["CD4_cat200_prim"], "Sem CD4 próximo ao diagnóstico")
    display(PVHA_prim_ult["CD4_diag_prim_cat200"].value_counts(dropna = False))
    print()
    PVHA_prim_ult["CD4_diag_prim_cat350"] = np.where(escolhas_cd4, PVHA_prim_ult["CD4_cat350_prim"], "Sem CD4 próximo ao diagnóstico")
    display(PVHA_prim_ult["CD4_diag_prim_cat350"].value_counts(dropna = False))


    PVHA_prim_ult["Raca_cat2"] = PVHA_prim_ult["Raca_cat2"].fillna("Não Informado")
    PVHA_prim_ult["Sexo_cat"] = PVHA_prim_ult["Sexo_cat"].fillna("Não Informado")
    DF_ult_ano["Raca_cat2"] = DF_ult_ano["Raca_cat2"].fillna("Não Informado")
    DF_ult_ano["Sexo_cat"] = DF_ult_ano["Sexo_cat"].fillna("Não Informado")

    return PVHA_prim_ult, DF_ult_ano



def gerar_base_BI(DF_ult_ano:pd.DataFrame = None,
                  DF_Prim:pd.DataFrame = None,
                  salvar_local:str = "//SAP109/Bancos AMA/Arquivos Atuais/PowerBI/HIV/"):
    """
    Carrega, organiza e filtra os dados de dispensação de TARV, exames de CD4 e carga viral, 
    salvando-os em um arquivo CSV formatado para a atualização do painel BI do Monitoramento Clínico do HIV.

    Se o DataFrame não for fornecido, ele será gerado utilizando a função `gerar_base_completa`.
    O arquivo CSV gerado contém uma linha por ano por paciente, com informações do último evento 
    de cada ano e colunas com informações do primeiro evento na vida do paciente.

    Parameters
    ----------
    hoje: datetime.date
        Data de fechamento dos bancos de análise no formato datetime.date(aaaa, mm, dd).
    
    ano_corte: int
        Ano de corte para filtrar os dados das bases de dispensação, carga viral e contagem de CD4.
        Default = ano - 10
    
    DF_ult_ano: pd.DataFrame, optional
        DataFrame contendo os dados organizados e filtrados dos pacientes. Se não for fornecido, 
        será gerado usando a função `gerar_base_completa`.
    
    salvar_local: str, optional
        Diretório onde o arquivo CSV será salvo. O caminho completo deve ser especificado. Caso não seja,
        o arquivo será salvo na pasta local.

    Returns
    -------
    None

    Example
    -------
    >>> hoje = pd.to_datetime(2024, 1, 31)
    >>> ano = hoje.year
    >>> ano_corte = ano - 10
    >>> gerar_base_BI(hoje, ano_corte, salvar_local='/caminho/para/salvar')
    """

    colunas_ult = [
        "ano",
        "Cod_unificado",
        "PVHA",
        "Vinculado_ano",
        "Status_ano",
        "Deteccao_CV50",
        "Deteccao_CV200",
        "Deteccao_CV1000",
        "data_dispensa",
        "Idade_cascata",
        "Raca_cat2",
        "Escol_cat",
        "Sexo_cat",
        "CD4_cat200",
        "ano_coleta_cd4",
        "Cod_ibge_insti_sol_50+",
        "Nome_mun_exame_50+",
        "nome_inst_exames_50+",
        "uf_inst_completo",
        "esquema_cat",
        "ibge_inst_exames",
        "Nome_mun_exame",
        "nome_inst_exames",
        'codigo_ibge_resid',
        'uf_res',
        'Nome_mun_resid'
        ]
    
    colunas_prim =[
        "Cod_unificado",
        "Cod_ibge_insti_sol_30+_prim",
        "Nome_mun_exame_30+_prim",
        "nome_inst_exames_30+_prim",
        "uf_inst_completo_prim",
        "Raca_cat2",
        "Escol_cat",
        "Sexo_cat",
        "Idade_vinc_cat",
        "data_min",
        "ano_min",
        "data_dispensa_prim",
        "CD4_diag_prim_cat200",
        "CD4_diag_prim_cat350",
        "CD4_diag_prim_cat",
        "Dias_diag_TARV_cat",
        'ibge_inst_exames_prim',
        'Nome_mun_exame_prim',
        'nome_inst_exames_prim',
        'codigo_ibge_resid',
        'uf_res',
        'Nome_mun_resid'
        ]

    DF = DF_Prim[colunas_prim].replace('>=','≥', regex=True)
    DF_ult = DF_ult_ano[colunas_ult].replace('>=','≥', regex=True)

    DF.to_csv(f"{salvar_local}Planilha_BI_HIV_prim.csv", index = False, sep = ";", decimal = ".")
    DF_ult.to_csv(f"{salvar_local}Planilha_BI_HIV_ult.csv", index = False, sep = ";", decimal = ".")

    print(f"Arquivos CSV salvos com sucesso na pasta {salvar_local}.")



def organizacao_geno(hoje: datetime.date, DF:pd.DataFrame, Cod_tab:pd.DataFrame, col_codigo_paciente:str = "cod_pac",
                    col_codigo_pac_uni:str = "Cod_unificado",
                    colunas_copias:list = ["copias_ultima_cv","copias_penultima_cv"], col_data_ultima_cv:str = "data_ultima_cv",
                    col_data_penultima_cv:str = "data_penultima_cv"):
    

    DF["codigo_paciente"] = pd.to_numeric(DF[col_codigo_paciente], errors="coerce")
    # Une o cadasto ao código Unificado
    DF = pd.merge(DF,Cod_tab, on = "codigo_paciente", how = "left")
     
    DF[col_codigo_pac_uni] = pd.to_numeric(DF[col_codigo_pac_uni], errors='coerce')
    DF = fg.ajusta_data_linha_vetorizado(DF, coluna_data = "data_coleta", coluna_retorno = "dt_geno")

    Geno = DF[(DF["dt_geno"] <= pd.to_datetime(hoje))].copy()

    Geno.drop(columns=["codigo_paciente",col_codigo_paciente, "cod_pac_final"], inplace=True)

    Geno["ano_geno"] = Geno["dt_geno"].dt.year

    # Ajuste das datas das últimas CVs do formulário
    Geno["data_ultima_cv"] = pd.to_datetime(Geno[col_data_ultima_cv], errors="coerce")
    Geno["data_penultima_cv"] = pd.to_datetime(Geno[col_data_penultima_cv], errors="coerce")

    for col in colunas_copias:
        Geno[f"{col}_limpo"] = Geno[col].apply(lambda x: str(x).strip())
        
        lista_indetec = [
            "< L. Min.","Não detec","Não detect","N?o detec","<50","<400",'0,0',
            "men 50","<L MIN",'<LIM','< 50','<40','< 40', '< min'
        ]
        for var in lista_indetec:
            Geno[f"{col}_limpo"] = Geno[f"{col}_limpo"].replace(var,0)
        
        
        lista_detec = [
            "> L. Max.",">500000",">750000",">ldl",">50","> 50000","máx",'> l. MA',
            ">L MAX", '> l max','>40000', '>2000', '>100000', '6079+9','>lim mx',
            '1.528.2', '.1.569', '>LMAX', '> L.max', '>LIM.MA', '>LIM', '>VAL.DE',
            '>500', '>limit', '>l.max', '>l. max', '>limite', '>maximo','>500.00',
            '>limmax', '>lim', '> L max', '>LIMITE', '>L.MAX.', 'lmax','33.209'
        ]
        for var in lista_detec:
            Geno[f"{col}_limpo"] = Geno[f"{col}_limpo"].replace(var,500000)
        
        
        lista_seminfo = [
            "",'n inf', 'ninf', 'IGNORAD', 'ninform', 'NINFORM', 'NINFOR', 'ntem',
            '-','ANDAM', 'não tem','ind','<MAXIMO', '<maximo', '<MAX'
        ]
        for var in lista_seminfo:
            Geno[f"{col}_limpo"] = Geno[f"{col}_limpo"].replace(var,np.nan)


    # Cria uma lista vazia para armazenar os valores problemáticos
    valores_problema = []

    # Itera sobre cada valor único na coluna
    for col in colunas_copias:
        for val in Geno[f"{col}_limpo"].unique():
            try:
                # Tenta converter o valor para número (float), mas não como int ainda
                num_val = pd.to_numeric(val, errors='raise')
            except (ValueError, TypeError):
                # Se houver um erro, adiciona o valor à lista de problemas
                valores_problema.append(val)

    # Exibindo a lista de valores problemáticos
    print("valores problema")
    print(valores_problema)

    # Remover pontos e converter para int
    Geno["copias_ultima_cv_limpo"] = Geno["copias_ultima_cv_limpo"].str.replace(".", "", regex=False)
    Geno["copias_ultima_cv_limpo"] = pd.to_numeric(Geno["copias_ultima_cv_limpo"], errors='coerce')

    # Remover pontos e converter para int
    Geno["copias_penultima_cv_limpo"] = Geno["copias_penultima_cv_limpo"].str.replace(".", "", regex=False)
    Geno["copias_penultima_cv_limpo"] = pd.to_numeric(Geno["copias_penultima_cv_limpo"], errors='coerce')

    cond_cv = [
        (Geno["copias_ultima_cv_limpo"] < 500),
        (Geno["copias_ultima_cv_limpo"] >= 500)
    ]

    escolhas_cv = [
            "<500",">=500"
    ]

    Geno["CV_ult_formgeno_cat500"] = np.select(cond_cv,escolhas_cv, default = "Não disponível")

    cond_cv = [
        (Geno["copias_penultima_cv_limpo"] < 200),
        (Geno["copias_penultima_cv_limpo"] >= 200)
    ]

    escolhas_cv = [
            "<200",">=200"
    ]

    Geno["CV_penult_formgeno_cat200"] = np.select(cond_cv,escolhas_cv, default = "Não disponível")

    return Geno



def gerar_base_dispGeno(hoje: datetime.date, Disp:pd.DataFrame, CV:pd.DataFrame, PVHA:pd.DataFrame,
                        Geno:pd.DataFrame, ano:int, tempo_max_cvs:int = 180, tempo_inicioTARV:int = 152,
                        lista_cols_linhas:list = ["Resultado_convencional","Resultado_integrase",
                                "SequenceName", "SequenceName_integrase",
                                'ATVrScore', 'DRVrScore', 'FPVrScore','IDVrScore',
                                'LPVrScore', 'NFVScore', 'SQVrScore', 'TPVrScore',
                                'ABCScore', 'AZTScore', 'D4TScore', 'DDIScore',
                                'FTCScore', '3TCScore', 'TDFScore', 'EFVScore',
                                'ETRScore', 'NVPScore', 'RPVScore', 'BICScore',
                                'DTGScore', 'EVGScore', 'RALScore', 'assignment',
                                "data_ultima_cv", "CV_ult_formgeno_cat500","data_penultima_cv",
                                "CV_penult_formgeno_cat200"],
                        col_codigo_pac_uni:str = "Cod_unificado"):

    Geno1 = Geno[(Geno["dt_geno"] >= f"{ano}-01-01") & (Geno["dt_geno"] <= f"{ano + 1}-12-31")].copy()

    count_geno = Geno1[col_codigo_pac_uni].value_counts().iloc[0]
    
    Geno_linha = fg.Na_linha(DF = Geno1,
                        lista_colunas = lista_cols_linhas,
                        col_data = "dt_geno",
                        cod_indent = col_codigo_pac_uni,
                        n = count_geno)

    Geno_linha.replace("",np.nan, inplace = True)

    Disp1 = Disp[(Disp["data_dispensa"] >= f"{ano - 1}-01-01") & (Disp["data_dispensa"] <= f"{ano + 1}-12-31")].copy()

    count_disp = Disp1[col_codigo_pac_uni].value_counts().iloc[0]

    Disp1["data_dispensa_Prim_vida"] = Disp1["data_dispensa_Prim"].copy()

    # Cria a variárivel em TARV no último dia do ano de análise
    Disp1["data_dispensa_ult_ano"] = Disp1[Disp1["ano_disp"] == ano].groupby([col_codigo_pac_uni])["data_dispensa"].transform("first")
    Disp1["data_dispensa_ult_ano"] = Disp1.groupby(col_codigo_pac_uni)["data_dispensa_ult_ano"].transform("max")
    Disp1["duracao_ult_ano"] = Disp1[Disp1["ano_disp"] == ano].groupby([col_codigo_pac_uni])["duracao_sum"].transform("first")
    Disp1["duracao_ult_ano"] = Disp1.groupby(col_codigo_pac_uni)["duracao_ult_ano"].transform("max")

    # Marca o status de TARV na data da geno
    Disp1["fim_ano"] = pd.to_datetime(f"{ano}-12-31")
    Disp1["hoje"] = pd.to_datetime(hoje)
    Disp1["fim_ano"] = np.where((Disp1["ano_disp"] == ano), Disp1["hoje"], Disp1["fim_ano"])
    Disp1["TARV_fim_ano"] = fg.TARV(DF = Disp1, col_data_ref = "fim_ano", col_dt_disp = "data_dispensa_ult_ano")

    Disp_linha = fg.Na_linha(DF = Disp1,
                        lista_colunas = ["esquema_AMA","duracao_sum"],
                        col_data = "data_dispensa",
                        cod_indent = col_codigo_pac_uni,
                        n = count_disp)

    Disp_geno = pd.merge(Disp_linha,Geno_linha, on = col_codigo_pac_uni, how = "left")

    # Marca a última dispensação realizada antes da data da Geno
    for index, row in Disp_geno.iterrows():
        found = False
        for i in range(count_disp):
            if (row["dt_geno"] - row[f"data_dispensa_{i}"]).days > 0:
                Disp_geno.at[index, "dt_disp_ult_geno"] = row[f"data_dispensa_{i}"]
                Disp_geno.at[index, "esquema_ult_geno"] = row[f"esquema_AMA_{i}"]
                Disp_geno.at[index, "duracao_ult_geno"] = row[f"duracao_sum_{i}"]
                found = True
                break  # Interrompe o loop após a primeira ocorrência
        
        if not found:
            # Se nenhum teste foi bem-sucedido, você pode definir valores padrão ou deixar em branco
            Disp_geno.at[index, "dt_disp_ult_geno"] = pd.NaT
            Disp_geno.at[index, "esquema_ult_geno"] = pd.NaT
            Disp_geno.at[index, "duracao_ult_geno"] = pd.NaT

    # Marca o status de TARV na data da geno
    Disp_geno["Status_geno"] = fg.TARV(DF = Disp_geno, col_data_ref = "dt_geno", col_dt_disp = "dt_disp_ult_geno")

    # União com o banco de cadastro para avaliar a residência
    Disp_geno_pac = pd.merge(PVHA, Disp_geno, on = col_codigo_pac_uni, how = "left")
    Disp_geno_pac = fg.ibge_resid(DF = Disp_geno_pac, cols_ibge = ["cod_ibge_udm"])

    # Filtro para idade maior de 18 anos no início do ano de análise
    Disp_geno_pac["inicio_ano"] = pd.to_datetime(f"{ano}-01-01")
    Disp_geno_pac["Idade_18"] = fg.idade_cat(DF = Disp_geno_pac, data_ref = "inicio_ano", faixas_etarias= [(0, 12),(12, 18),(18,99)])
    Disp_geno_pac = Disp_geno_pac[(Disp_geno_pac["Idade_18"] == "Mais de 18 anos")]

    CV1 = CV[(CV["data_hora_coleta_cv"] >= f"{ano - 1}-01-01") & (CV["data_hora_coleta_cv"] <= f"{ano + 1}-12-31")].copy()

    CV1["CV_count_ano"] = CV1[CV1["ano_coleta_cv"] == ano].groupby([col_codigo_pac_uni])["ano_coleta_cv"].transform("count")

    count_cv = CV1[col_codigo_pac_uni].value_counts().iloc[0]
    print(count_cv)

    CV_linha = fg.Na_linha(DF = CV1,
                        lista_colunas = ["CV_cat200","CV_cat500"],
                        col_data = "data_hora_coleta_cv",
                        cod_indent = col_codigo_pac_uni,
                        n = count_cv)
    
    print(CV_linha.shape)

    DF = pd.merge(Disp_geno_pac, CV_linha, on = col_codigo_pac_uni, how = "left")

    # Número máximo de dias entre os dois exames seguidos e detectáveis de CV
    def falhas(row, tempo_inicioTARV = tempo_inicioTARV):
        for i in range(count_cv-1, 0, -1):
            if pd.notnull(row[f"data_hora_coleta_cv_{i}"]) and (row[f"data_hora_coleta_cv_{i}"] - row["data_dispensa_Prim_vida"]).days > tempo_inicioTARV and row[f"CV_cat200_{i}"] == ">=200": # 152 dias do início da TARV para que a segunda CV após 28 dias tenha no mínimo 180 dias da TARV
                
                if i - 1 > 0 and pd.notnull(row[f"data_hora_coleta_cv_{i-1}"]) and (row[f"data_hora_coleta_cv_{i-1}"] - row[f"data_hora_coleta_cv_{i}"]).days >= 28 and (row[f"data_hora_coleta_cv_{i-1}"] - row[f"data_hora_coleta_cv_{i}"]).days <= tempo_max_cvs and row[f"CV_cat500_{i-1}"] == ">=500":
                    for r in range(i-2, 0, -1):
                        if pd.notnull(row[f"data_hora_coleta_cv_{r}"]) and (row[f"data_hora_coleta_cv_{r}"] - row[f"data_hora_coleta_cv_{i-1}"]).days >= 180 and row[f"CV_cat200_{r}"] == ">=200":
                            if r - 1 > 0 and pd.notnull(row[f"data_hora_coleta_cv_{r-1}"]) and (row[f"data_hora_coleta_cv_{r-1}"] - row[f"data_hora_coleta_cv_{r}"]).days >= 28 and (row[f"data_hora_coleta_cv_{r-1}"] - row[f"data_hora_coleta_cv_{r}"]).days <= tempo_max_cvs and row[f"CV_cat500_{r-1}"] == ">=500":
                                return row[f"data_hora_coleta_cv_{i-1}"], row[f"data_hora_coleta_cv_{r-1}"]
                            elif r - 2 > 0 and pd.notnull(row[f"data_hora_coleta_cv_{r-2}"]) and (row[f"data_hora_coleta_cv_{r-2}"] - row[f"data_hora_coleta_cv_{r}"]).days >= 28 and (row[f"data_hora_coleta_cv_{r-2}"] - row[f"data_hora_coleta_cv_{r}"]).days <= tempo_max_cvs and row[f"CV_cat500_{r-2}"] == ">=500":
                                return row[f"data_hora_coleta_cv_{i-1}"], row[f"data_hora_coleta_cv_{r-2}"]
                            elif r - 3 > 0 and pd.notnull(row[f"data_hora_coleta_cv_{r-3}"]) and (row[f"data_hora_coleta_cv_{r-3}"] - row[f"data_hora_coleta_cv_{r}"]).days >= 28 and (row[f"data_hora_coleta_cv_{r-3}"] - row[f"data_hora_coleta_cv_{r}"]).days <= tempo_max_cvs and row[f"CV_cat500_{r-3}"] == ">=500":
                                return row[f"data_hora_coleta_cv_{i-1}"], row[f"data_hora_coleta_cv_{r-3}"]
                    return row[f"data_hora_coleta_cv_{i-1}"], None
                
                elif i - 2 > 0 and pd.notnull(row[f"data_hora_coleta_cv_{i-2}"]) and (row[f"data_hora_coleta_cv_{i-2}"] - row[f"data_hora_coleta_cv_{i}"]).days >= 28 and (row[f"data_hora_coleta_cv_{i-2}"] - row[f"data_hora_coleta_cv_{i}"]).days <= tempo_max_cvs and row[f"CV_cat500_{i-2}"] == ">=500":
                    for r in range(i-3, 0, -1):
                        if pd.notnull(row[f"data_hora_coleta_cv_{r}"]) and (row[f"data_hora_coleta_cv_{r}"] - row[f"data_hora_coleta_cv_{i-2}"]).days >= 180 and row[f"CV_cat200_{r}"] == ">=200":
                            if r - 1 > 0 and pd.notnull(row[f"data_hora_coleta_cv_{r-1}"]) and (row[f"data_hora_coleta_cv_{r-1}"] - row[f"data_hora_coleta_cv_{r}"]).days >= 28 and (row[f"data_hora_coleta_cv_{r-1}"] - row[f"data_hora_coleta_cv_{r}"]).days <= tempo_max_cvs and row[f"CV_cat500_{r-1}"] == ">=500":
                                return row[f"data_hora_coleta_cv_{i-2}"], row[f"data_hora_coleta_cv_{r-1}"]
                            elif r - 2 > 0 and pd.notnull(row[f"data_hora_coleta_cv_{r-2}"]) and (row[f"data_hora_coleta_cv_{r-2}"] - row[f"data_hora_coleta_cv_{r}"]).days >= 28 and (row[f"data_hora_coleta_cv_{r-2}"] - row[f"data_hora_coleta_cv_{r}"]).days <= tempo_max_cvs and row[f"CV_cat500_{r-2}"] == ">=500":
                                return row[f"data_hora_coleta_cv_{i-2}"], row[f"data_hora_coleta_cv_{r-2}"]
                            elif r - 3 > 0 and pd.notnull(row[f"data_hora_coleta_cv_{r-3}"]) and (row[f"data_hora_coleta_cv_{r-3}"] - row[f"data_hora_coleta_cv_{r}"]).days >= 28 and (row[f"data_hora_coleta_cv_{r-3}"] - row[f"data_hora_coleta_cv_{r}"]).days <= tempo_max_cvs and row[f"CV_cat500_{r-3}"] == ">=500":
                                return row[f"data_hora_coleta_cv_{i-2}"], row[f"data_hora_coleta_cv_{r-3}"]
                    return row[f"data_hora_coleta_cv_{i-2}"], None
                
                elif i - 3 > 0 and pd.notnull(row[f"data_hora_coleta_cv_{i-3}"]) and (row[f"data_hora_coleta_cv_{i-3}"] - row[f"data_hora_coleta_cv_{i}"]).days >= 28 and (row[f"data_hora_coleta_cv_{i-3}"] - row[f"data_hora_coleta_cv_{i}"]).days <= tempo_max_cvs and row[f"CV_cat500_{i-3}"] == ">=500":
                    for r in range(i-4, 0, -1):
                        if pd.notnull(row[f"data_hora_coleta_cv_{r}"]) and (row[f"data_hora_coleta_cv_{r}"] - row[f"data_hora_coleta_cv_{i-3}"]).days >= 180 and row[f"CV_cat200_{r}"] == ">=200":
                            if r - 1 > 0 and pd.notnull(row[f"data_hora_coleta_cv_{r-1}"]) and (row[f"data_hora_coleta_cv_{r-1}"] - row[f"data_hora_coleta_cv_{r}"]).days >= 28 and (row[f"data_hora_coleta_cv_{r-1}"] - row[f"data_hora_coleta_cv_{r}"]).days <= tempo_max_cvs and row[f"CV_cat500_{r-1}"] == ">=500":
                                return row[f"data_hora_coleta_cv_{i-3}"], row[f"data_hora_coleta_cv_{r-1}"]
                            elif r - 2 > 0 and pd.notnull(row[f"data_hora_coleta_cv_{r-2}"]) and (row[f"data_hora_coleta_cv_{r-2}"] - row[f"data_hora_coleta_cv_{r}"]).days >= 28 and (row[f"data_hora_coleta_cv_{r-2}"] - row[f"data_hora_coleta_cv_{r}"]).days <= tempo_max_cvs and row[f"CV_cat500_{r-2}"] == ">=500":
                                return row[f"data_hora_coleta_cv_{i-3}"], row[f"data_hora_coleta_cv_{r-2}"]
                            elif r - 3 > 0 and pd.notnull(row[f"data_hora_coleta_cv_{r-3}"]) and (row[f"data_hora_coleta_cv_{r-3}"] - row[f"data_hora_coleta_cv_{r}"]).days >= 28 and (row[f"data_hora_coleta_cv_{r-3}"] - row[f"data_hora_coleta_cv_{r}"]).days <= tempo_max_cvs and row[f"CV_cat500_{r-3}"] == ">=500":
                                return row[f"data_hora_coleta_cv_{i-3}"], row[f"data_hora_coleta_cv_{r-3}"]
                    return row[f"data_hora_coleta_cv_{i-3}"], None
                
        return None, None

    # Aplicar a função falha ao DataFrame e dividir a saída em duas colunas
    # Salva a data da segunda carga viral que compõe a falha
    DF[['Data_falha1', 'Data_falha2']] = DF.apply(falhas, axis=1, result_type="expand")

    # Cria a Falha pelo resultado do formulário
    cond = (
        ((DF["data_penultima_cv"] - DF["data_dispensa_Prim_vida"]).dt.days > 28) &
        (DF["CV_penult_formgeno_cat200"] == ">=200") &
        (DF["CV_ult_formgeno_cat500"] == ">=500") &
        (DF["data_ultima_cv"].dt.year == ano)
    )

    DF["Falha_form"] = np.where(cond, "Falha", "Sem falha")

    DF["tempo_inicio_TARV_form"] = (DF["data_penultima_cv"] - DF["data_dispensa_Prim_vida"]).dt.days
    
    # Unifica as falhas em uma única coluna com a primeira falha do ano complementada pelo formulário
    DF['Data_falha1'] = pd.to_datetime(DF['Data_falha1'], errors="coerce")
    DF['Data_falha2'] = pd.to_datetime(DF['Data_falha2'], errors="coerce")
    condicao = [
        (DF['Data_falha1'].dt.year == ano) & DF['Data_falha1'].notna(),
        (DF['Data_falha2'].dt.year == ano) & DF['Data_falha2'].notna(),
        (DF['data_ultima_cv'].dt.year == ano) & (DF['data_ultima_cv'].notna()) & (DF["Falha_form"] == "Falha")
        ]
    
    escolha = [
        DF['Data_falha1'].astype(str),
        DF['Data_falha2'].astype(str),
        DF['data_ultima_cv'].astype(str)
    ]    

    DF['Data_falha'] = np.select(condicao, escolha, default=pd.NaT)
    DF['Data_falha'] = pd.to_datetime(DF['Data_falha'], errors='coerce')

    # Cria coluna com binário se falhou ou não no ano
    DF["Falha"] = np.where((DF['Data_falha'].isnull() == False), "Falha", "Sem falha")

    DF["ano"] = ano

    DF = fg.marcar_ultimas_antes(DF, data_ref= "Data_falha", data_busca = "data_dispensa", n_busca = count_disp,
                        lista_cols = ["duracao_sum"],
                        n_salvos = 1, nomes_seq = ["ult_falha"])

    DF["TARV_falha"] = fg.TARV(DF, col_data_ref = "Data_falha", col_dt_disp = "data_dispensa_ult_falha", col_duracao = "duracao_sum_ult_falha")

    del Geno1, Geno_linha, Disp1, Disp_linha, Disp_geno, Disp_geno_pac, CV1, CV_linha
    gc.collect()

    return DF, count_geno, count_cv



def gerar_excel(hoje: datetime.date, dict_bases: dict, salvar_cascata:bool = True, salvar_spectrum:bool = True):    

    ano = hoje.year
    mes = hoje.month

    pasta = f"V:/{ano}/Monitoramento e Avaliação/DOCUMENTOS/Clinico/Arquivos Excel/"

    # --- Variáveis sociodemográficas ---
    socio = [
        "Sexo_cat",
        "Raca_cat2",
        "Escol_cat",
        "reg_res",
        "Pop_genero"
    ]

    list_cd4 = [
        "CD4_cat",
        "CD4_cat_Prim",
    ]

    list_CV = [
         "Deteccao_CV50",
         "Deteccao_CV200",
         "Deteccao_CV1000"
    ]

    list_disp = [
        "Status_ano",
        "Atraso_ano"
    ]

    # --- Colunas de referência para todas as bases ---
    Estrat_var = {
        'PVHA_ano':{"ano":socio+["Idade_cascata"]},
        'Vinc10_ano':{"ano":socio+["Idade_cascata"]},
        'Vinc5_ano':{"ano":socio+["Idade_cascata"]},
        'Vinc_ano':{"ano":socio+["Idade_cascata"]},
        'PVHA_atual':{
            "reg_res":socio+["Idade_cascata"]+list_cd4+list_CV+list_disp,
            "CD4_cat":socio+["Idade_cascata"]+list_CV+list_disp,
            "Deteccao_CV50":socio+["Idade_cascata"]+list_cd4+list_disp,
            "Deteccao_CV200":socio+["Idade_cascata"]+list_cd4+list_disp,
            "Deteccao_CV1000":socio+["Idade_cascata"]+list_cd4+list_disp,
            "Raca_cat2":socio+["Idade_cascata"]+list_cd4+list_CV+list_disp,
            "Idade_cascata":socio+["Idade_cascata"]+list_cd4+list_CV+list_disp,
            "Status_ano":socio+["Idade_cascata"]
            },
        'Vinc10_atual':{"reg_res":socio+["Idade_cascata"]+list_cd4+list_CV+list_disp},
        'Vinc5_atual':{"reg_res":socio+["Idade_cascata"]+list_cd4+list_CV+list_disp},
        'Vinc_atual':{"reg_res":socio+["Idade_cascata"]+list_cd4+list_CV+list_disp},
        'Retidos_ano':{"ano":socio+["Idade_cascata"]},
        'Vinc_ano_gap':{"ano":socio+["Idade_cascata"]},
        'Vinc_ano_perda':{"ano":socio+["Idade_cascata"]},
        'GAP_completo_ano':{"ano":socio+["Idade_cascata"]},
        'GAP_ano':{"ano":socio+["Idade_cascata"]},
        'GAP_sinan_ano':{"ano":socio+["Idade_cascata"]},
        'GAP_completo_atual':{"reg_res":socio+["Idade_cascata"]+list_cd4},
        'GAP_atual':{"reg_res":socio+["Idade_cascata"]+list_cd4},
        'GAP_sinal_atual':{"reg_res":socio+["Idade_cascata"]},
        'TARV_ano':{
            "ano":socio+["Idade_cascata","esquema_cat"]
            },
        'TARV_atual':{
            "reg_res":socio+["Idade_cascata"]+list_cd4+list_CV+list_disp,
            "CD4_cat":socio+["Idade_cascata"]+list_CV+list_disp,
            "Deteccao_CV50":socio+["Idade_cascata"]+list_cd4+list_disp,
            "Deteccao_CV200":socio+["Idade_cascata"]+list_cd4+list_disp,
            "Deteccao_CV1000":socio+["Idade_cascata"]+list_cd4+list_disp,
            "Raca_cat2":socio+["Idade_cascata"]+list_cd4+list_CV+list_disp,
            "Idade_cascata":socio+["Idade_cascata"]+list_cd4+list_CV+list_disp,
            "PVHA":["esquema_AMA"]
            },
        'Supressao_ano':{"ano":socio+["Idade_cascata"]},
        'Supressao_atual':{"reg_res":socio+["Idade_cascata"]+list_cd4+list_CV+list_disp},
        'Perda_ano':{"ano":socio+["Idade_cascata"]},
        'Perda_atual':{
            "reg_res":socio+["Idade_cascata"]+list_cd4+list_CV+list_disp,
            "CD4_cat":socio+["Idade_cascata"]+list_CV+list_disp,
            "Raca_cat2":socio+["Idade_cascata"]+list_cd4+list_CV+list_disp,
            "Idade_cascata":socio+["Idade_cascata"]+list_cd4+list_CV+list_disp
            },
        'PVHA_prim':{
            "ano_min":socio+["Idade_vinc_cat","Status_atual_completo", "CD4_diag_prim_cat"],
            "ano_coleta_cd4_prim":["CD4_cat_prim"]
            },
        'Vinc_novos_ano_cd200menos':{"ano_min":socio+["Idade_vinc_cat","Status_atual_completo"]},
        'Vinc_novos_ano_cd200mais':{"ano_min":socio+["Idade_vinc_cat","Status_atual_completo"]},
        'GAP_prim':{"ano_min":socio+["Idade_vinc_cat"]},
        'Disp_prim':{"ano_disp_prim":socio+["Idade_vinc_cat","Dias_diag_TARV_cat","esquema_cat_prim"]},
        'Disp_prim_atual':{"Dias_diag_TARV_cat":socio+["Idade_vinc_cat"]}
    }

    path = f"{pasta}Monitoramento_HIV_{mes}_{ano}.xlsx"
    # cria arquivo vazio sem aba padrão se não existir
    if not os.path.isfile(path):
        book = openpyxl.Workbook()
        book.save(path)

    writer = pd.ExcelWriter(
        path,
        engine="openpyxl",
        mode="a",
        if_sheet_exists="overlay"
    )

    # processa cada base
    for nome, base in dict_bases.items():
        estratificacoes = Estrat_var[nome]

        for index, colunas in estratificacoes.items():
            tabelas = []
            tabelas_decimal = []

            for var in colunas:
                df = pd.crosstab(
                    index=base[index],
                    columns=base[var],
                    margins=True,
                    margins_name="Total"
                )
                df.drop("Total", axis=0, inplace=True, errors="ignore")
                df = df.sort_index(ascending=False)

                # se não existe coluna Total, pula
                if "Total" not in df.columns:
                    continue

                df_decimal = (df.div(df["Total"], axis=0)).round(3)
                df_decimal.drop(columns="Total", inplace=True)

                if var in ["Sexo_cat", "Idade_vinc_cat"]:
                    df_decimal.drop(columns="Não Informado", inplace=True, errors="ignore")

                tabelas.append(df)
                tabelas_decimal.append(df_decimal)

            # salva no Excel com controle de espaçamento
            linha = 0
            for n in range(len(tabelas)):
                df = tabelas[n]
                df_dec = tabelas_decimal[n]

                df.to_excel(writer, sheet_name=f"{nome}_{index}", startrow=linha)
                df_dec.to_excel(
                    writer,
                    sheet_name=f"{nome}_{index}",
                    startrow=linha,
                    startcol=df.shape[1] + 3
                )

                linha += len(df) + 5  # espaçamento entre tabelas


                                                        ## Incluir cascata no excel
    if salvar_cascata:
        DFs_cascata = [
            dict_bases["DF_ult_ano"],
            dict_bases["DF_ult_ano5"],
            dict_bases["Vinc_ano"],
            dict_bases["Retidos"],
            dict_bases["TARV_ano"],
            dict_bases["TARV_ano"][dict_bases["TARV_ano"]["Abandono_sum"] == 0],
            dict_bases["Supressao_ano"][dict_bases["Supressao_ano"]["CV_cat1000"] == "<1.000"],
            dict_bases["Supressao_ano"][dict_bases["Supressao_ano"]["CV_cat50"] == "<50"],
            dict_bases["GAP_ano"],
            dict_bases["Vinc_ano_gap"],
            dict_bases["Perda_ano"],
            dict_bases["Vinc_ano_perda"]
        ]

        nomes_cascata = [
            "Vinculados 10 anos",
            "Vinculados 5 anos",
            "Vinculados no ano",
            "Retidos",
            "Tarv",
            "Adesão",
            "Supressão 1000",
            "Supressão 50",
            "Vinculados 10 anos sem TARV",
            "Vinculados no ano sem TARV",
            "Vinculados 10 anos em Perda",
            "Vinculados no ano em Perda"
        ]

        estratificacoes = [
            "Raca_cat2",
            "Escol_cat",
            "Sexo_cat",
            "reg_res",
            "Idade_cascata",
            "Orientacao_cat",
            "Genero_cat"
        ]

        lista_linhas = []
        lista_cascata = []
        for DF, nomes_df in zip(DFs_cascata,nomes_cascata):
            
            df_temp = pd.DataFrame(DF["ano"].value_counts()).rename(columns = {"count":nomes_df}).sort_index().T
            lista_linhas.append(df_temp)


        # Adicionando a supressão em decimal
        df_temp = pd.DataFrame(pd.crosstab(columns=dict_bases["Supressao_ano"]["CV_cat1000"],
                    index=dict_bases["Supressao_ano"]['ano'],
                    normalize = "index")["<1.000"]).rename(columns = {"<1.000":"Supressão 1000cp %"}).T
        df_temp = round(df_temp,3)
        df_temp = df_temp.astype(str)
        lista_linhas.append(df_temp)

        df_temp = pd.DataFrame(pd.crosstab(columns=dict_bases["Supressao_ano"]["CV_cat50"],
                    index=dict_bases["Supressao_ano"]['ano'],
                    normalize = "index")["<50"]).rename(columns = {"<50":"Supressão 50cp %"}).T
        df_temp = round(df_temp,3)
        df_temp = df_temp.astype(str)
        lista_linhas.append(df_temp)

        cascata = pd.concat(lista_linhas, axis =0)
        cascata.loc[["Supressão 50cp %","Supressão 1000cp %"]] = cascata.loc[["Supressão 50cp %","Supressão 1000cp %"]].astype(float)
        cascata.index.name = "TOTAL"
        lista_cascata.append(cascata)


        for cols in estratificacoes:
            for categoria in dict_bases["DF_ult_ano"][dict_bases["DF_ult_ano"][cols] != "Não Informado"][cols].unique():
                lista_linhas = []
                for DF, nomes_df in zip(DFs_cascata,nomes_cascata):
                    
                    DF_cat = DF[DF[cols] == categoria]
                    df_temp = pd.DataFrame(DF_cat["ano"].value_counts()).rename(columns = {"count":nomes_df}).sort_index().T
                    lista_linhas.append(df_temp)


                # Adicionando a supressão em decimal
                DF_sup = dict_bases["Supressao_ano"][dict_bases["Supressao_ano"][cols] == categoria]
                df_temp = pd.DataFrame(pd.crosstab(columns=DF_sup["CV_cat1000"],
                            index=DF_sup['ano'],
                            normalize = "index")["<1.000"]).rename(columns = {"<1.000":"Supressão 1000cp %"}).T
                df_temp = round(df_temp,3)
                df_temp = df_temp.astype(str)
                lista_linhas.append(df_temp)
                
                df_temp = pd.DataFrame(pd.crosstab(columns=DF_sup["CV_cat50"],
                            index=DF_sup['ano'],
                            normalize = "index")["<50"]).rename(columns = {"<50":"Supressão 50cp %"}).T
                df_temp = round(df_temp,3)
                df_temp = df_temp.astype(str)
                lista_linhas.append(df_temp)
                
                cascata = pd.concat(lista_linhas, axis =0)
                cascata.loc[["Supressão 50cp %","Supressão 1000cp %"]] = cascata.loc[["Supressão 50cp %","Supressão 1000cp %"]].astype(float)
                cascata.index.name = categoria
                lista_cascata.append(cascata)

        for n in range(len(lista_cascata)):
            lista_cascata[n].to_excel(writer,sheet_name = f"Cascata", startrow = 1 + (len(lista_cascata[n])+3)*n)




    ## Incluir Spectrum no excel
    if salvar_spectrum:
        DFs_spectrum = [
            dict_bases["DF_ult_ano"],
            dict_bases["DF_ult_ano5"],
            dict_bases["Vinc_ano"],
            dict_bases["Vinc_ano_gap"],
            dict_bases["TARV_ano"],
            dict_bases["Vinc_ano_perda"],
            dict_bases["TARV_ano"][dict_bases["TARV_ano"]["Abandono_sum"] == 0],
            dict_bases["Supressao_ano"][dict_bases["Supressao_ano"]["CV_cat1000"] == "<1.000"],
            dict_bases["DF_Prim"],
            dict_bases["Disp_prim"]
        ]

        nomes_spectrum = [
            "Vinculados_10anos",
            "Vinculados_5anos",
            "Vinculados",
            "Vinculados sem TARV",
            "Tarv",
            "Interrompido no ano",
            "Adesão",
            "Supressão 1000",
            "Novos Vinvulados",
            "Iniciaram TARV"
        ]

        tabelas_Spectrum = []
        for ano_spec in dict_bases["DF_ult_ano"]["ano"].unique():
            
            lista_cascata = []            
            lista_linhas = []
            for DF, nomes_df in zip(DFs_spectrum,nomes_spectrum):
                
                DF_cat = DF[(DF["Idade_Spec_crianca"] == "Mais de 15 anos") & (DF["Sexo_cat"] != "Não Informado") &
                            (DF["ano"] == ano_spec)]
                df_temp = pd.DataFrame(DF_cat["Sexo_cat"].value_counts()).rename(columns = {"count":nomes_df}).sort_index().T
                lista_linhas.append(df_temp)
            
            
            # Adicionando a supressão em decimal
            DF_sup = dict_bases["Supressao_ano"][(dict_bases["Supressao_ano"]["Idade_Spec_crianca"] == "Mais de 15 anos") & (dict_bases["Supressao_ano"]["Sexo_cat"] != "Não Informado") &
                            (dict_bases["Supressao_ano"]["ano"] == ano_spec)]
            df_temp = pd.DataFrame(pd.crosstab(columns=DF_sup["CV_cat1000"],
                        index=DF_sup['Sexo_cat'],
                        normalize = "index")["<1.000"]).rename(columns = {"<1.000":"Supressão 1000cp %"}).T
            df_temp = round(df_temp,3)
            df_temp = df_temp.astype(str)
            lista_linhas.append(df_temp)
            
            cascata = pd.concat(lista_linhas, axis =0)
            cascata.loc[["Supressão 1000cp %"]] = cascata.loc[["Supressão 1000cp %"]].astype(float)
            cascata.index.name = ano_spec
            lista_cascata.append(cascata)
            
            
            lista_linhas = []
            for DF, nomes_df in zip(DFs_spectrum,nomes_spectrum):
                
                DF_cat = DF[(DF["Idade_Spec_crianca"] == "Menos de 15 anos") & (DF["Sexo_cat"] != "Não Informado") & (DF["ano"] == ano_spec)]
                df_temp = pd.DataFrame(DF_cat["Idade_Spec_crianca"].value_counts()).rename(columns = {"count":nomes_df}).sort_index().T
                lista_linhas.append(df_temp)
            
            
            # Adicionando a supressão em decimal
            DF_sup = dict_bases["Supressao_ano"][(dict_bases["Supressao_ano"]["Idade_Spec_crianca"] == "Menos de 15 anos") & (dict_bases["Supressao_ano"]["Sexo_cat"] != "Não Informado") & (dict_bases["Supressao_ano"]["ano"] == ano_spec)]
            df_temp = pd.DataFrame(pd.crosstab(columns=DF_sup["CV_cat1000"],
                        index=DF_sup['Idade_Spec_crianca'],
                        normalize = "index")["<1.000"]).rename(columns = {"<1.000":"Supressão 1000cp %"}).T
            df_temp = round(df_temp,3)
            df_temp = df_temp.astype(str)
            lista_linhas.append(df_temp)
            
            cascata = pd.concat(lista_linhas, axis =0)
            cascata.loc[["Supressão 1000cp %"]] = cascata.loc[["Supressão 1000cp %"]].astype(float)
            cascata.index.name = ano_spec
            lista_cascata.append(cascata)
            
            
            Spectrum = pd.concat(lista_cascata, axis =1)
            tabelas_Spectrum.append(Spectrum)

        for n in range(len(tabelas_Spectrum)):
            tabelas_Spectrum[n].to_excel(writer,sheet_name = f"Spectrum", startrow = 1, startcol = 1 + (tabelas_Spectrum[n].shape[1] + 3)*n)


        Bases_spec_estrat = [
            dict_bases["DF_ult_ano"],
            dict_bases["DF_ult_ano5"],
            dict_bases["Vinc_ano"],
            dict_bases["TARV_ano"],
            dict_bases["DF_Prim"],
            dict_bases["Disp_prim"]              
            ]
        
        nome_spec_estrat = [
            "Vinculados 10 anos",
            "Vinculados 5 anos",
            "Vinculados no ano",
            "Em TARV",
            "Novos Vinc",
            "Inicio TARV"
            ]

        for DF_spec, nome_espec in zip(Bases_spec_estrat,nome_spec_estrat):
            tabelas1 = []
            tabelas2 = []
            for ano_spec in dict_bases["DF_ult_ano"]["ano"].unique():
                
                DF = DF_spec[(DF_spec["ano"] == ano_spec) & (DF_spec["Sexo_cat"] != "Não Informado") &
                            (DF_spec["Idade_Spec1"] != "Não Informado")]
                
                df_temp = pd.crosstab(columns=DF["Sexo_cat"],
                            index=DF['Idade_Spec1'],
                            margins=True,
                            margins_name = "Total")
                df_temp.index.name = ano_spec
                tabelas1.append(df_temp)
                
                df_temp = pd.crosstab(columns=DF["Sexo_cat"],
                            index=DF['Idade_Spec2'],
                            margins=True,
                            margins_name = "Total")
                df_temp.index.name = ano_spec
                tabelas2.append(df_temp)
                
            for n in range(len(tabelas1)):
                tabelas1[n].to_excel(writer,sheet_name = f"Spectrum_{nome_espec}", startrow = 1, startcol =  1 + (tabelas1[n].shape[1] + 3)*n)

            for n in range(len(tabelas2)):
                tabelas2[n].to_excel(writer,sheet_name = f"Spectrum_{nome_espec}", startrow = len(tabelas1[n])+3, startcol =  1 + (tabelas2[n].shape[1] + 3)*n)

    writer.close()

    print(f"Arquivo Excel salvo com sucesso na {pasta}")

    

def gerar_resumo(hoje: datetime.date, dict_bases: dict):

    ano = hoje.year
    mes_nome = fg.mes_nome(hoje=hoje)

    PVHA = dict_bases["PVHA_ano"]
    DF_ult_ano = dict_bases["Vinc10_ano"]
    Vinc_ano = dict_bases["Vinc_ano"]
    DF_Prim = dict_bases["PVHA_prim"]
    GAP_completo_atual = dict_bases["GAP_completo_atual"]
    GAP_atual = dict_bases["GAP_atual"]
    Disp_prim = dict_bases["Disp_prim"]
    GAP_prim = dict_bases["GAP_prim"]
    TARV_atual = dict_bases["TARV_atual"]
    Perda_atual = dict_bases["Perda_atual"]
    Supressao_atual = dict_bases["Supressao_atual"]
    TARV_atual = dict_bases["TARV_atual"]


    #PVHA
    PVHA_atual_count = PVHA[PVHA["ano"] == ano].shape[0]

    #Vinculados
    Vinc_atual_count = DF_ult_ano[DF_ult_ano["ano"] == ano].shape[0]
    Vinc_ano_atual_count = Vinc_ano[Vinc_ano["ano"] == ano].shape[0]
    PVHA_novos_12mes_count = DF_Prim[DF_Prim["data_min"] > (hoje.replace(year=hoje.year - 1))].shape[0]
    CD4_200_12mes_count = pd.DataFrame(DF_Prim[DF_Prim["data_min"] > (hoje.replace(year=hoje.year - 1))]["CD4_cat200_prim"].value_counts().sort_index(ascending = False)).T["<200 células / mm³"].iloc[0]

    # GAP de tratamento
    GAP_completo_atual_count = GAP_completo_atual.shape[0]
    GAP_atual_count = GAP_atual.shape[0]
    GAP_novos_12mes_count = GAP_prim[GAP_prim["data_min"] > (hoje.replace(year=hoje.year - 1))].shape[0] # Número de pessoas que entratam em GAP nos últimos 12 meses

    # Início de TARV
    Disp_novos_12mes_count = Disp_prim[Disp_prim["data_dispensa_prim"] > (hoje.replace(year=hoje.year - 1))].shape[0]

    # Em TARV
    EmTARV_count = TARV_atual.shape[0]
    Perda_count = Perda_atual.shape[0]
    ComCV_count = Supressao_atual.shape[0]
    Suprimidos_count = Supressao_atual[Supressao_atual["Deteccao_CV1000"] == "CV suprimida"].shape[0]
    Indetectavel_count = Supressao_atual[Supressao_atual["Deteccao_CV50"] == "CV indetectável"].shape[0]


    # Subnotificação
    Subnotif = pd.DataFrame(DF_Prim[DF_Prim["DT_NOTIFIC_Ad"].isna() &
              DF_Prim["DT_NOTIFIC_cr"].isna()]["ano_min"].value_counts().sort_index())
    
    Total = pd.DataFrame(DF_Prim["ano_min"].value_counts().sort_index())

    print("\033[1mPVHA Diagnosticadas Vivas e/ou Vinculadas\033[0m")
    print()   
    print(f"Em {mes_nome}/{ano}, havia {PVHA_atual_count} PVHA diagnosticadas sem registro de óbito.")
    print()
    print(f"Em {mes_nome}/{ano}, havia {Vinc_atual_count} PVHA vinculadas aos serviços de atenção ao HIV do SUS nos últimos 10 anos.")
    print()
    print(f"Nos últimos 12 meses, {Vinc_ano_atual_count} PVHA mantiveram vínculo com os serviços de atenção ao HIV do SUS.")
    print()
    print(f"Nos últimos 12 meses, {PVHA_novos_12mes_count} novas PVHA foram registradas nos sistemas do HIV.")
    print()
    print(f"Nos últimos 12 meses, {CD4_200_12mes_count} ({(CD4_200_12mes_count/PVHA_novos_12mes_count)*100:.1f}%) das novas PVHA, tinham contagem de CD4 <200 células / mm³.")
    print()
    print()
    print("\033[1mPVHA em GAP de Tratamento\033[0m")
    print()
    print(f"Em {mes_nome}/{ano}, havia {GAP_completo_atual_count} PVHA em GAP de tratamento no Brasil, sendo que {GAP_atual_count} possuem vínculo no Siscel sem dispensação no Siclom.")
    print()
    print(f"Nos últimos 12 meses, {GAP_novos_12mes_count} novas PVHA entraram em GAP de tratamento.")
    print()
    print()
    print("\033[1mPVHA que iniciaram TARV\033[0m")
    print()
    print(f"Nos últimos 12 meses, {Disp_novos_12mes_count} PVHA iniciaram TARV no Brasil.")
    print()
    print()
    print("\033[1mPVHA em TARV\033[0m")
    print()
    print(f"Em {mes_nome}/{ano}, havia {EmTARV_count} PVHA em TARV e {Perda_count} em perda de seguimento no Brasil.")
    print()
    print(f"Das pessoas em TARV, {ComCV_count} ({round(ComCV_count/EmTARV_count*100,1)}%) realizaram exame de CV no SUS após 6 meses do incício da TARV " \
        f"e desses, {Suprimidos_count} ({round(Suprimidos_count/ComCV_count*100,1)}%) apresentam CV suprimida (<1000 cp) e " \
        f"{Indetectavel_count} ({round(Indetectavel_count/ComCV_count*100,1)}%) apresentam CV indetectável (<50 cp).")
    print()
    print()
    print("\033[1mSubnotificação\033[0m")
    print()
    print("Subnotificação global de todas as PVHA é de",
        round(len(DF_Prim[DF_Prim["DT_NOTIFIC_Ad"].isna() &
        DF_Prim["DT_NOTIFIC_cr"].isna()])/len(DF_Prim)*100,1),"%")
    print()
    print(f"Subnotificação em {ano-2} é de {round((Subnotif/Total)*100,1).loc[ano-2][0]}%")
    print(f"Subnotificação em {ano-1} é de {round((Subnotif/Total)*100,1).loc[ano-1][0]}%")
    print()



def gerar_tabelas_simc(hoje: datetime.date, Perda_atual:pd.DataFrame, GAP_atual:pd.DataFrame, Supressao_atual:pd.DataFrame,
                       var_perda:list = [
                            "codigo_paciente",
                            "data_nascimento",
                            "Idade_Hoje",
                            "Sexo_cat",
                            "Raca_cat",
                            "Escol_cat",
                            "codigo_ibge_resid",
                            "cd_inst_sol_cv",
                            "cd_inst_sol_cd4",
                            "codigo_udm",
                            "data_dispensa",
                            "duracao_sum",
                            "dt_lim_prox_disp",
                            "Atraso",
                            "Status_ano",
                            "data_hora_coleta_cd4",
                            "contagem_cd4",
                            "cd_crm",
                            "uf_crm"
                        ],
                        var_GAP:list = [
                            "codigo_paciente",
                            "data_nascimento",
                            "Idade_Hoje",
                            "Sexo_cat",
                            "Raca_cat",
                            "Escol_cat",
                            "codigo_ibge_resid",
                            "cd_inst_sol_cv",
                            "cd_inst_sol_cd4",
                            "data_hora_coleta_cd4",
                            "contagem_cd4"
                        ],
                        var_Sup:list = [
                            "codigo_paciente",
                            "data_nascimento",
                            "Idade_Hoje",
                            "Sexo_cat",
                            "Raca_cat",
                            "Escol_cat",
                            "codigo_ibge_resid",
                            "cd_inst_sol_cv",
                            "cd_inst_sol_cd4",
                            "codigo_udm",
                            "data_dispensa",
                            "CV_cat",
                            "data_hora_coleta_cv",
                            "comentario_copias",
                            "copias",
                            "cd_crm",
                            "uf_crm"
                        ]):
    
    ano = hoje.year
    mes = hoje.month
    mes_nome = fg.mes_nome_completo(hoje)

    pasta = f"V:/{ano}/Monitoramento e Avaliação/COMPARTILHADO/AMA - Banco de Dados/Tabelas SIMC/"

     # Perda de Seguimento
    Tab_perda = Perda_atual[var_perda].copy()

    Tab_perda["codigo_paciente"] = Tab_perda["codigo_paciente"].astype(int)
    Tab_perda["codigo_paciente"] = Tab_perda["codigo_paciente"].apply(lambda x: str(x).zfill(8))

    Tab_perda["cd_inst_sol_cv"] = Tab_perda["cd_inst_sol_cv"].astype("Int64")

    Tab_perda["cd_inst_sol_cd4"] = Tab_perda["cd_inst_sol_cd4"].astype("Int64")

    Tab_perda["duracao_sum"] = Tab_perda["duracao_sum"].astype(int)
    Tab_perda["Atraso"] = Tab_perda["Atraso"].astype(int)

    display(Tab_perda)
    print()

    # GAP de tratamento
    Tab_gap = GAP_atual[var_GAP].copy()

    Tab_gap["codigo_paciente"] = Tab_gap["codigo_paciente"].astype(int)
    Tab_gap["codigo_paciente"] = Tab_gap["codigo_paciente"].apply(lambda x: str(x).zfill(8))

    Tab_gap["cd_inst_sol_cv"] = Tab_gap["cd_inst_sol_cv"].astype("Int64")

    Tab_gap["cd_inst_sol_cd4"] = Tab_gap["cd_inst_sol_cd4"].astype("Int64")

    Tab_gap["contagem_cd4"] = Tab_gap["contagem_cd4"].astype("Int64")

    display(Tab_gap)
    print()

    # CV Detectável
    Tab_cv = Supressao_atual[var_Sup].copy()

    Tab_cv["codigo_paciente"] = Tab_cv["codigo_paciente"].astype(int)
    Tab_cv["codigo_paciente"] = Tab_cv["codigo_paciente"].apply(lambda x: str(x).zfill(8))

    Tab_cv["cd_inst_sol_cv"] = Tab_cv["cd_inst_sol_cv"].astype("Int64")

    Tab_cv["cd_inst_sol_cd4"] = Tab_cv["cd_inst_sol_cd4"].astype("Int64")

    display(Tab_cv)
    print()

    # Verificar se o diretório existe, caso contrário, criar
    caminho = f"{pasta}{str((mes-4)+(12*(ano-2025))).zfill(2)} - {mes_nome} {ano}"

    if not os.path.exists(caminho):
        os.makedirs(caminho)


    Tab_perda.to_csv(f"{caminho}/tb_perda_simc.csv", sep = ";", encoding = "latin-1", index = False)
    Tab_gap.to_csv(f"{caminho}/tb_gap_simc.csv", sep = ";", encoding = "latin-1", index = False)
    Tab_cv.to_csv(f"{caminho}/tb_cv_simc.csv", sep = ";", encoding = "latin-1", index = False)

    print(f"Tabelas SIMC salvas com sucesso na {caminho}/")



def gerar_ppt(hoje:datetime.date):

    ano = hoje.year
    mes = hoje.month

    pasta = f"V:/{ano}/Monitoramento e Avaliação/DOCUMENTOS/Clinico/Arquivos PPT/"
     
    def salvar_ppt(figura:str = "Fig.png"):
        slide_layout = ppt.slide_layouts[6]
        slide = ppt.slides.add_slide(slide_layout)
        grafico = figura
        left = Inches(0.5)
        top = Inches(0.5)
        height = Inches(6.6)
        slide.shapes.add_picture(grafico, left, top, height=height)
        
        left = Inches(0.2)
        top = Inches(6.8)
        width = Inches(2.5)
        height = Inches(0.5)
        source_text_box = slide.shapes.add_textbox(left, top, width, height)
        source_text_frame = source_text_box.text_frame
        source_text = source_text_frame.add_paragraph()
        source_text.text = f"Fonte: DATHI/SVS/MS \n Atualizado até {str(hoje)[:10]}"
        source_text.font.size = Pt(10)
        source_text.font.color.rgb = RGBColor(137, 137, 137)

    if not os.path.isfile(f"{pasta}Monitoramento_HIV_{mes}_{ano}.pptx"):

        # Create a presentation object
        ppt = Presentation()
        
        # Set the slide size to Widescreen 16:9
        ppt.slide_width = Inches(13.333)  # 16:9 aspect ratio
        ppt.slide_height = Inches(7.5)  # 16:9 aspect ratio
        
        # Add a title slide
        slide_layout = ppt.slide_layouts[6]
        slide = ppt.slides.add_slide(slide_layout)
        
        title_box = slide.shapes.add_textbox(Inches(3.1), Inches(2), width=Inches(6), height=Inches(1))
        tf = title_box.text_frame
        p = tf.add_paragraph()
        p.text = "Monitoramento Clínico HIV/AIDS"
        p.font.name = 'Calibri Light'
        p.font.size = Pt(40)
        p.font.color.rgb = RGBColor(118, 113, 113)
        p.font.bold = True
        
        title_box = slide.shapes.add_textbox(Inches(2.85), Inches(2.7), width=Inches(4), height=Inches(1))
        tf = title_box.text_frame
        p = tf.add_paragraph()
        p.text = "Assessoria de Monitoramento e Avaliação - DATHI"
        p.font.name = 'Calibri Light'
        p.font.size = Pt(28)
        p.font.color.rgb = RGBColor(118, 113, 113)
        p.font.bold = True
        
        title_box = slide.shapes.add_textbox(Inches(5.5), Inches(5), width=Inches(2), height=Inches(1))
        tf = title_box.text_frame
        p = tf.add_paragraph()
        p.text = "Banco de " + str(hoje)[:10]  
        p.font.name = 'Montserrat'
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(89, 89, 89)
        
        
        # Add a "line" (thin rectangle) above the title "Monitoramento Clínico TARV"
        line1 = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(3), Inches(2) - Pt(1),  # Adjust these values to change the start position
            Inches(7.5), Pt(0.2)  # Adjust these values to change the end position (Pt(1) creates a thin rectangle)
        )
        fill = line1.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(169, 169, 169)  # A light gray color
        line1.line.fill.background()
        
        # Add a "line" (thin rectangle) below the title "HIV/AIDS"
        line2 = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(3), Inches(3.5) + Pt(28),  # Adjust these values to change the start position
            Inches(7.5), Pt(0.2)  # Adjust these values to change the end position (Pt(1) creates a thin rectangle)
        )
        fill = line2.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(169, 169, 169)  # A light gray color
        line2.line.fill.background()

        ppt.save(f'{pasta}Monitoramento_Clínico_{mes}_{ano}.pptx')


    ppt = Presentation(f'{pasta}Monitoramento_Clínico_{mes}_{ano}.pptx')


    for figura in sorted(glob.glob("Figura_*.png")):
        salvar_ppt(figura=figura)
        
    ppt.save(f'{pasta}Monitoramento_Clínico_{mes}_{ano}.pptx')

    print(f"Arquivo PPT salvo com sucesso na {pasta}")



def gerar_bases_mensais(PVHA_ult_ano:pd.DataFrame, PVHA_prim_ult:pd.DataFrame,
                        pasta:str = "//SAP109/Bancos AMA/Arquivos Atuais/Bancos Atuais HIV/Mensais/"):

    PVHA_prim_ult = PVHA_prim_ult.replace('≥', '>=', regex=True)
    PVHA_prim_ult.to_csv(f"{pasta}PVHA_prim_ult.csv", index = False, sep = ";", decimal = ".", encoding = "latin-1")

    PVHA_ult_ano = PVHA_ult_ano.replace('≥', '>=', regex=True)
    PVHA_ult_ano.to_csv(f"{pasta}PVHA_ult_ano.csv", index = False, sep = ";", decimal = ".", encoding = "latin-1")


    PVHA_prim_ult["cd_crm_prim"] = PVHA_prim_ult["cd_crm_prim"].astype(str)
    PVHA_prim_ult["cd_crm_ult"] = PVHA_prim_ult["cd_crm_ult"].astype(str)
    PVHA_prim_ult["codigo_udm_prim"] = PVHA_prim_ult["codigo_udm_prim"].astype(str)
    PVHA_prim_ult["codigo_udm_ult"] = PVHA_prim_ult["codigo_udm_ult"].astype(str)
    PVHA_prim_ult["ibge_inst_exames_prim"] = PVHA_prim_ult["ibge_inst_exames_prim"].astype(str)
    PVHA_prim_ult["ibge_inst_exames_ult"] = PVHA_prim_ult["ibge_inst_exames_ult"].astype(str)

    PVHA_prim_ult.to_parquet(f"{pasta}PVHA_prim_ult.parquet", index = False)

    PVHA_ult_ano["escolaridade"] = PVHA_ult_ano["escolaridade"].astype(str)
    PVHA_ult_ano["codigo_udm"] = PVHA_ult_ano["codigo_udm"].astype(str)
    PVHA_ult_ano["cd_crm"] = PVHA_ult_ano["cd_crm"].astype(str)
    PVHA_ult_ano["ibge_inst_exames"] = PVHA_ult_ano["ibge_inst_exames"].astype(str)
    PVHA_ult_ano.to_parquet(f"{pasta}PVHA_ult_ano.parquet", index = False)

    print(f"Bases salvas com sucesso na {pasta}")



def gerar_bases_PVHA_mes(hoje: datetime.date,
                     col_codigo_pac_uni: str = "Cod_unificado",
                     col_mes = "mes_ano",
                     cols_datas = ["data_dispensa","data_hora_coleta_cv", "data_hora_coleta_cd4"],
                     col_data_nascimento = "data_nascimento",
                     col_copias = "copias",
                     col_comentario_copias = "comentario_copias",
                     col_contagem_cd4 = "contagem_cd4",
                     anos_obito_adm = 10,
                     data_disp_prim = "data_dispensa_Prim"):
    

    # Carregamento dos Bancos
    PVHA, Pac, Disp, CV, CD4, Sinan_A, Sinan_cong, Sinan_G, Sinan_cr, Sim = carregar_bases_PVHA(
        PVHA = True,
        Cadastro = True,
        Disp = True,
        CV = True,
        CD4=True,
        Sinan=True,
        Sinan_G=True,
        Sinan_cr=True,
        Sim=True
    )

    # Organização dos Bancos

    ### Siclom\SISCEL

    # Organiza o cadastro
    PVHA = organizacao_cadastro(
        PVHA,
        Pac,
        hoje
    )

    # Organiza a dispensação
    Disp = organizacao_disp(
        hoje, Disp,
        esquemas = True,
        Prim = True
    )

    CV = organizacao_cv(
        hoje, CV,
        categorizar = True,
        Prim = True
    )

    CD4 = organizacao_cd4(
        hoje, CD4, 
        categorizar = True,
        Prim = True
    )

    ### SINAN

    Sinan_A  = organizacao_sinan(Sinan_A, tipo = "Ad")

    Sinan_G  = organizacao_sinan(Sinan_G, tipo = "gest")

    Sinan_cr = organizacao_sinan(Sinan_cr, tipo = "cr")

    Sinan_cong["CS_ESCOL_N"] = np.nan
    Sinan_cong = organizacao_sinan(Sinan_cong, tipo = "cong")



    ########################### UNIÃO DOS BANCOS ###############################################

    # Pega apenas a entrada mais antiga
    Sinan_A = Sinan_A.sort_values(["DT_NOTIFIC_Ad"], ascending = True)
    Sinan_Prim = Sinan_A.drop_duplicates(subset = [col_codigo_pac_uni], keep="first")
    Sinan_Prim = Sinan_Prim.reset_index(drop = True)

    Sinan_G = Sinan_G.sort_values(["DT_NOTIFIC_gest"], ascending = True)
    Sinan_G_Prim = Sinan_G.drop_duplicates(subset = [col_codigo_pac_uni], keep="first")
    Sinan_G_Prim = Sinan_G_Prim.reset_index(drop = True)

    Sinan_cr = Sinan_cr.sort_values(["DT_NOTIFIC_cr"], ascending = True)
    Sinan_cr_Prim = Sinan_cr.drop_duplicates(subset = [col_codigo_pac_uni], keep="first")
    Sinan_cr_Prim = Sinan_cr_Prim.reset_index(drop = True)

    Sim = Sim.drop_duplicates(subset = [col_codigo_pac_uni], keep="first")
    Sim = Sim.reset_index(drop = True)


    Disp.sort_values(["data_dispensa"], ascending = [False], inplace = True)
    Disp.reset_index(drop = True, inplace = True)

    CV.sort_values(["data_hora_coleta_cv"], ascending = [False], inplace = True)
    CV.reset_index(drop = True, inplace = True)

    CD4.sort_values(["data_hora_coleta_cd4"], ascending = [False], inplace = True)
    CD4.reset_index(drop = True, inplace = True)

    # Organizar as bases pelo ultimo evento de cada ano
    PVHA_mes, Disp_mes, CV_mes, CD4_mes = bases_ult_meses(hoje, DF_lista=[PVHA, Disp, CV, CD4], DF_nome=["PVHA","Disp", "CV", "CD4"], col_data=["data_min","data_dispensa", "data_hora_coleta_cv", "data_hora_coleta_cd4"]).values()

    print("Base PVHA_mes")
    display(PVHA_mes[col_mes].value_counts().sort_index(ascending = False))
    print()

    DF_mes = pd.merge(PVHA_mes,Disp_mes, on = [col_mes, col_codigo_pac_uni], how = "left")
    DF_mes = pd.merge(DF_mes,CD4_mes, on = [col_mes,col_codigo_pac_uni], how = "left")
    DF_mes = pd.merge(DF_mes,CV_mes, on = [col_mes,col_codigo_pac_uni], how = "left")

    DF_mes = pd.merge(DF_mes, Sinan_Prim[[
        col_codigo_pac_uni,'Escol_num_Ad',"ID_MN_RESI_Ad",
        "ID_MUNICIP_Ad","DT_NOTIFIC_Ad"]], on = col_codigo_pac_uni, how = "left")
    
    DF_mes = pd.merge(DF_mes, Sinan_G_Prim[[
        col_codigo_pac_uni,'Escol_num_gest',"ID_MN_RESI_gest",
        "ID_MUNICIP_gest","DT_NOTIFIC_gest"]], on = col_codigo_pac_uni, how = "left")
    
    DF_mes = pd.merge(DF_mes, Sinan_cr_Prim[[
        col_codigo_pac_uni,'Escol_num_cr',"ID_MN_RESI_cr",
        "ID_MUNICIP_cr","DT_NOTIFIC_cr"]], on = col_codigo_pac_uni, how = "left")

    DF_mes = pd.merge(DF_mes, Sinan_cong[[
        col_codigo_pac_uni,'Escol_num_cong',"ID_MN_RESI_cong",
        "ID_MUNICIP_cong","DT_NOTIFIC_cong"]], on = col_codigo_pac_uni, how = "left")

    print("Base DF_mes após o Merge")
    display(DF_mes[col_mes].value_counts().sort_index(ascending = False))
    print()

    del  PVHA_mes, Disp_mes, CV_mes, CD4_mes, PVHA, Disp, CV, CD4, Sinan_Prim, Sinan_G_Prim, Sinan_cr_Prim, Sinan_cong
    gc.collect()

    ######################## PVHA_MES ###########################################

    DF_mes["mes_ano_min"] = DF_mes["data_min"].dt.strftime("%m_%Y")
    DF_mes["mes_ano_obito"] = DF_mes["data_obito"].dt.strftime("%m_%Y")
    DF_mes["mes_ano_disp"] = DF_mes["data_dispensa"].dt.strftime("%m_%Y")
    DF_mes["mes_ano_cv"] = DF_mes["data_hora_coleta_cv"].dt.strftime("%m_%Y")
    DF_mes["mes_ano_cd4"] = DF_mes["data_hora_coleta_cd4"].dt.strftime("%m_%Y")
    DF_mes["ano"] = DF_mes["mes_ano"].str[3:].astype("Int64")


    # Filtra excluindo: linhas de anos a partir do ano do óbito.
    DF_mes = DF_mes[((DF_mes["mes_ano_obito"].isnull()) |
                    (DF_mes["mes_ano"] < DF_mes["mes_ano_obito"]))].copy()
    
    print("Base DF_mes após o excluir os óbitos")
    display(DF_mes["mes_ano"].value_counts().sort_index(ascending = False))
    print()

    DF_mes['Escol_num'] = DF_mes[['Escol_num','Escol_num_Ad','Escol_num_gest','Escol_num_cr','Escol_num_cong']].max(axis=1)

    Escol = [
        (DF_mes["Escol_num"] == 1),
        (DF_mes["Escol_num"] == 2),
        (DF_mes["Escol_num"] == 3)
    ]

    Escol_escolha = ["0 a 7 anos",
                    "8 a 11 anos",
                    "12 e mais anos"]

    DF_mes["Escol_cat"] = np.select(Escol, Escol_escolha, default="Não Informado")
    display(DF_mes["Escol_cat"].value_counts(dropna = False))
    print()

    # Calcula a idade em todas as datas passadas no cols_idade.
    for cols in cols_datas:
        DF_mes[f"Idade_{cols}_cat"] = fg.idade_cat(DF = DF_mes, data_ref=cols, data_nasc=col_data_nascimento)
        display(DF_mes[f"Idade_{cols}_cat"].value_counts().sort_index())
        print()


    # Crie uma série de datas com o último dia do mês para cada mês_ano presente na coluna mes_ano
    DF_mes["Ult_dia_mes"] = pd.to_datetime(DF_mes["mes_ano"], format="%m-%Y") + pd.offsets.MonthEnd(0)

    # Normalize as datas para garantir que apenas a data seja mantida (sem informações de hora)
    DF_mes["Ult_dia_mes"] = DF_mes["Ult_dia_mes"].dt.normalize()

    # Calcula a idade no último dia de cada ano (importante para Cascata e BI)
    DF_mes["Idade_cascata"] = fg.idade_cat(DF = DF_mes, data_ref="Ult_dia_mes", data_nasc=col_data_nascimento)
    display(DF_mes["Idade_cascata"].value_counts().sort_index())
    print()

    data_disp = next((col for col in cols_datas if "disp" in col.lower()), None)

    data_cv = next((col for col in cols_datas if "cv" in col.lower() or "carga" in col.lower()), None)

    data_cd4 = next((col for col in cols_datas if "cd4" in col.lower()), None)

        
    # Organiza o IBGE de residência
    DF_mes = fg.ibge_resid(DF=DF_mes, cols_ibge = [
                        "ID_MN_RESI_Ad","ID_MN_RESI_gest","ID_MN_RESI_cr", "ID_MN_RESI_cong",
                        "cod_ibge_solicitante_cv","cod_ibge_solicitante_cd4","cod_ibge_udm",
                        "ID_MUNICIP_Ad","ID_MUNICIP_gest", "ID_MUNICIP_cr", "ID_MUNICIP_cong"
                        ])

    # Organiza o IBGE das instituições Solicitantes
    DF_mes = fg.ibge_inst_sol_exames(DF_mes)

    ## A Coluna "Status_ano" mostra a situação do paciente em relação ao último dia do mês em que a dispensação foi feita, ou em relação a data atual
    # Calcula o atraso em relação ao último dia do mês da dispensação
    DF_mes["Atraso_mes"] = (DF_mes["Ult_dia_mes"] - DF_mes["dt_lim_prox_disp"]).dt.days

    cond_status = [
        (DF_mes["Atraso_mes"] <= 60),
        (DF_mes["Atraso_mes"] > 60),
        (DF_mes["data_hora_coleta_cv"].isnull() == False) | (DF_mes["data_hora_coleta_cd4"].isnull() == False)
        ]

    escolha_status = ["Tarv", "Interrupção de Tarv", "Gap de tratamento"]

    DF_mes["Status_mes"] = np.select(cond_status, escolha_status, default= None)

    DF_mes["Status_mes"] = DF_mes["Status_mes"].fillna("Gap de vinculação")

    display("Status no último dia do mês",
            pd.DataFrame(DF_mes[["mes_ano","Status_mes"]].value_counts(dropna = False).sort_index(ascending = False)))
    print()

    # Cria uma variável com o vínculo no ano
    Cond_vin = [
        (((DF_mes["Ult_dia_mes"] - DF_mes[data_cv]).dt.days <= 365) |
        ((DF_mes["Ult_dia_mes"] - DF_mes[data_cd4]).dt.days <= 365) |
        ((DF_mes["Ult_dia_mes"] - DF_mes[data_disp]).dt.days <= 365))
    ]

    escolha_vin = [
        "Vinculado", "Vinculado"
    ]

    DF_mes["Vinculado_ano"]  = np.select(Cond_vin,escolha_vin, default = "Não vinculado")
    display("Vinculados nos últimos 12 meses para o mês de referência",
            pd.DataFrame(DF_mes[["mes_ano","Vinculado_ano"]].value_counts(dropna = False).sort_index(ascending = False)))
    print()


    # Categoriza as faixas de Carga Viral
    DF_mes[col_copias] = pd.to_numeric(DF_mes[col_copias], errors='coerce')
    DF_mes[col_comentario_copias] = pd.to_numeric(DF_mes[col_comentario_copias], errors='coerce')

    cond_cv = [(
            ((DF_mes["Ult_dia_mes"] - DF_mes[data_cv]).dt.days <= 365) &
            (((DF_mes[col_copias].isnull() == True) & (DF_mes[col_comentario_copias] == 0)) |
            ((DF_mes[col_copias].isnull() == True) & (DF_mes[col_comentario_copias] == 1)) |
            (DF_mes[col_copias] < 50))
            ),
            (
            ((DF_mes["Ult_dia_mes"] - DF_mes[data_cv]).dt.days <= 365) &
            (((DF_mes[col_copias].isnull() == True) & (DF_mes[col_comentario_copias] == 2)) |
            (DF_mes[col_copias] >= 10000))
            ),
            (
            ((DF_mes["Ult_dia_mes"] - DF_mes[data_cv]).dt.days <= 365) &
            (DF_mes[col_copias] >= 50) &
            (DF_mes[col_copias] < 200)
            ),
            (
            ((DF_mes["Ult_dia_mes"] - DF_mes[data_cv]).dt.days <= 365) &
            (DF_mes[col_copias] >= 200) &
            (DF_mes[col_copias] < 1000)
            ),
            (
            ((DF_mes["Ult_dia_mes"] - DF_mes[data_cv]).dt.days <= 365) &
            (DF_mes[col_copias] >= 1000) &
            (DF_mes[col_copias] < 10000)
            )
    ]

    escolhas_cv = [
            "<50","10.000+","50-199","200-999","1.000-9.999"
    ]

    DF_mes["CV_cat"] = np.select(cond_cv,escolhas_cv, default = "Não disponível")
    display(DF_mes["CV_cat"].value_counts())
    print()

    # Cria a categoria com 1000 copias de CV   
    cond_cv = [
            (DF_mes["CV_cat"] == "1.000-9.999") | (DF_mes["CV_cat"] == "10.000+"),
            (DF_mes["CV_cat"] == "<50") | (DF_mes["CV_cat"] == "50-199") | (DF_mes["CV_cat"] == "200-999")   
    ]

    escolhas_cv = [
            ">=1.000",
            "<1.000"
    ]

    DF_mes["CV_cat1000"] = np.select(cond_cv,escolhas_cv, default = "Não disponível")
    display(DF_mes["CV_cat1000"].value_counts())
    print()

    # Cria a categoria com 200 copias de CV
    cond_cv = [
            (DF_mes["CV_cat"] == "<50") | (DF_mes["CV_cat"] == "50-199"),
            (DF_mes["CV_cat"] == "1.000-9.999") | (DF_mes["CV_cat"] == "10.000+") |
            (DF_mes["CV_cat"] == "200-999")   
    ]

    escolhas_cv = [
            "<200",
            ">=200"
    ]

    DF_mes["CV_cat200"] = np.select(cond_cv,escolhas_cv, default = "Não disponível")
    display(DF_mes["CV_cat200"].value_counts())
    print()

    # Cria a categoria com 50 copias de CV   
    cond_cv = [
            (DF_mes["CV_cat"] == "<50"),
            (DF_mes["CV_cat"] == "1.000-9.999") | (DF_mes["CV_cat"] == "10.000+") |
            (DF_mes["CV_cat"] == "50-199") | (DF_mes["CV_cat"] == "200-999")   
    ]

    escolhas_cv = [
            "<50",
            ">=50"
    ]

    DF_mes["CV_cat50"] = np.select(cond_cv,escolhas_cv, default = "Não disponível")
    display(DF_mes["CV_cat50"].value_counts())
    print()

    # Cria uma variável com a indetecção viral em 50 cópias
    Cond_sup = [
        (
        (DF_mes["Status_mes"] == "Tarv") & 
        ((DF_mes[data_cv] - DF_mes[data_disp_prim]).dt.days > 180) &
        (DF_mes["CV_cat50"] == "<50")
        ),
        (
        (DF_mes["Status_mes"] == "Tarv") & 
        ((DF_mes[data_cv] - DF_mes[data_disp_prim]).dt.days > 180) &
        (DF_mes["CV_cat50"] == ">=50")
        )
    ]

    escolha_sup = [
    "CV indetectável","CV detectável"
    ]

    DF_mes["Deteccao_CV50"]  = np.select(Cond_sup,escolha_sup, default = "Sem dados de CV")
    display("Indetecção Viral pós TARV na última CV do mês (<50 cp)",
            DF_mes[["mes_ano","Deteccao_CV50"]].value_counts(dropna = False).sort_index())

    # Cria uma variável com a supressão viral em 200 cópias
    Cond_sup = [
        (
        (DF_mes["Status_mes"] == "Tarv") & 
        ((DF_mes[data_cv] - DF_mes[data_disp_prim]).dt.days > 180) &
        (DF_mes["CV_cat200"] == "<200")
        ),
        (
        (DF_mes["Status_mes"] == "Tarv") & 
        ((DF_mes[data_cv] - DF_mes[data_disp_prim]).dt.days > 180) &
        (DF_mes["CV_cat200"] == ">=200")
        )
    ]

    escolha_sup = [
        "CV suprimida","CV não suprimida"
    ]

    DF_mes["Deteccao_CV200"]  = np.select(Cond_sup,escolha_sup, default = "Sem dados de CV")
    display("Supressão Viral pós TARV na última CV do mês (<200 cp)",
            DF_mes[["mes_ano","Deteccao_CV200"]].value_counts(dropna = False).sort_index())


    # Cria uma variável com a supressão viral em 1000 cópias
    Cond_sup = [
        (
        (DF_mes["Status_mes"] == "Tarv") & 
        ((DF_mes[data_cv] - DF_mes[data_disp_prim]).dt.days > 180) &
        (DF_mes["CV_cat1000"] == "<1.000")
        ),
        (
        (DF_mes["Status_mes"] == "Tarv") & 
        ((DF_mes[data_cv] - DF_mes[data_disp_prim]).dt.days > 180) &
        (DF_mes["CV_cat1000"] == ">=1.000")
        )
    ]

    escolha_sup = [
        "CV suprimida","CV não suprimida"
    ]

    DF_mes["Deteccao_CV1000"]  = np.select(Cond_sup,escolha_sup, default = "Sem dados de CV")
    display("Supressão Viral pós TARV na última CV do ano (<1000 cp)",
            DF_mes[["mes_ano","Deteccao_CV1000"]].value_counts(dropna = False).sort_index())

    DF_mes[col_contagem_cd4] = pd.to_numeric(DF_mes[col_contagem_cd4], errors="coerce")

    # criando variável de categorias do valor do CD4
    Cond_CD4cat = [
        (DF_mes[col_contagem_cd4] > 0) & (DF_mes[col_contagem_cd4] < 200),
        (DF_mes[col_contagem_cd4] >= 200) & (DF_mes[col_contagem_cd4] < 350),
        (DF_mes[col_contagem_cd4] >= 350) & (DF_mes[col_contagem_cd4] < 500),
        (DF_mes[col_contagem_cd4] >= 500) & (DF_mes[col_contagem_cd4] < 2500)
        ]

    Escolha_CD4cat = ["0-199", "200-349", "350-499", "500+"]

    DF_mes["CD4_cat"] = np.select(Cond_CD4cat, Escolha_CD4cat, default="Sem resultado de CD4")

    display(DF_mes["CD4_cat"].value_counts(dropna = False).sort_index())
    print()

    # criando variável de categorias do valor do CD4
    Cond_CD4cat2 = [
        (DF_mes[col_contagem_cd4] > 0) & (DF_mes[col_contagem_cd4] < 200),
        (DF_mes[col_contagem_cd4] >= 200) & (DF_mes[col_contagem_cd4] < 2500)
        ]

    Escolha_CD4cat2 = ["<200 células / mm³", "≥200 células / mm³"]

    DF_mes["CD4_cat200"] = np.select(Cond_CD4cat2, Escolha_CD4cat2, default="Sem resultado de CD4")

    display(DF_mes["CD4_cat200"].value_counts(dropna = False).sort_index())
    print()

    # Calcular o último ano de registro de evento de cada pessoa
    ultimo_evento = DF_mes.groupby(col_codigo_pac_uni).agg({
        "data_min":"max",
        "data_dispensa": 'max',
        "data_hora_coleta_cv": 'max',
        "data_hora_coleta_cd4": 'max'
    }).max(axis=1).reset_index()
    ultimo_evento.columns = [col_codigo_pac_uni, 'ultimo_evento']

    ultimo_vinculo = DF_mes.groupby(col_codigo_pac_uni).agg({
        "data_dispensa": 'max',
        "data_hora_coleta_cv": 'max',
        "data_hora_coleta_cd4": 'max'
    }).max(axis=1).reset_index()
    ultimo_vinculo.columns = [col_codigo_pac_uni, 'ultimo_vinculo']

    # Mesclar essa informação de último evento com o DataFrame original
    DF_mes = DF_mes.merge(ultimo_evento, on=col_codigo_pac_uni)
    DF_mes = DF_mes.merge(ultimo_vinculo, on=col_codigo_pac_uni)


    # Filtrar os dados onde a diferença entre o ano atual e o último evento é menor ou igual a 10 anos
    DF_mes = DF_mes[(DF_mes["Ult_dia_mes"] - DF_mes['ultimo_evento']).dt.days <= anos_obito_adm*365]

    print("Base DF_mes após o excluir os óbitos ADM")
    display(DF_mes["mes_ano"].value_counts().sort_index(ascending = False))
    print()

    # Adiciona nova coluna para os vinculados nos últimos 5 anos
    DF_mes['Vinculado_5anos'] = np.where((DF_mes['ultimo_vinculo'].notna()) &
                                        ((DF_mes["Ult_dia_mes"] - DF_mes['ultimo_vinculo']).dt.days <= 5*365), "Vinculado", "Não Vinculado")

    # Adiciona nova coluna para os vinculados nos últimos 5 anos
    DF_mes['Vinculado_10anos'] = np.where((DF_mes['ultimo_vinculo'].notna()) &
                                        ((DF_mes["Ult_dia_mes"] - DF_mes['ultimo_vinculo']).dt.days <= 10*365), "Vinculado", "Não Vinculado")

    # Filtra os cadastros inválidos do Siclom
    DF_mes = DF_mes[
        ~((DF_mes["st_num"].isin([2,3])) &
          (DF_mes["Status_mes"].isin(["Interrupção de Tarv","Gap de tratamento"])) &
          ((DF_mes["Ult_dia_mes"] - DF_mes['ultimo_vinculo']).dt.days >= 2*365))
        ]

    print("Base DF_mes após o excluir os cadastros inválidos do SICLOM")
    display(DF_mes["mes_ano"].value_counts().sort_index(ascending = False))
    print()


    DF_mes["Raca_cat2"] = DF_mes["Raca_cat2"].fillna("Não Informado")
    DF_mes["Sexo_cat"] = DF_mes["Sexo_cat"].fillna("Não Informado")

    return DF_mes



