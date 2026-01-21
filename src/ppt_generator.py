from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os
import shutil
import pandas as pd

def generate_ppt(output_dir, metrics, data_fechamento):
    """
    Gera a apresentação PowerPoint com os slides e gráficos.
    """
    print("Gerando apresentação PowerPoint...")
    
    ppt = Presentation()
    ppt.slide_width = Inches(13.333)
    ppt.slide_height = Inches(7.5)
    
    # Variáveis do metrics
    hoje2 = metrics.get('hoje2', '')
    latest_month_year_count = metrics.get('latest_month_year_count', '0')
    mediana_uso = metrics.get('mediana_uso', 0)
    
    emPrEP_porcent = metrics.get('emPrEP_porcent', '0')
    formatted_em_prep_count = metrics.get('formatted_em_prep_count', '0')
    descontinuados_porcent = metrics.get('descontinuados_porcent', '0')
    formatted_discontinued_count = metrics.get('formatted_discontinued_count', '0')
    
    highest_percentage_pop = metrics.get('highest_percentage_pop', 0)
    formatted_highest_count_pop = metrics.get('formatted_highest_count_pop', '0')
    highest_category_pop = metrics.get('highest_category_pop', '')
    
    highest_percentage_fetar = metrics.get('highest_percentage_fetar', 0)
    formatted_highest_count_fetar = metrics.get('formatted_highest_count_fetar', '0')
    highest_category_fetar = metrics.get('highest_category_fetar', '')
    
    formatted_young_percentage = metrics.get('formatted_young_percentage', '0')
    formatted_young_counts = metrics.get('formatted_young_counts', '0')
    median_age = metrics.get('median_age', 0)
    
    highest_percentage_escol = metrics.get('highest_percentage_escol', 0)
    formatted_highest_count_escol = metrics.get('formatted_highest_count_escol', '0')
    highest_category_escol = metrics.get('highest_category_escol', '')
    
    highest_percentage_raca = metrics.get('highest_percentage_raca', 0)
    formatted_highest_count_raca = metrics.get('formatted_highest_count_raca', '0')
    highest_category_raca = metrics.get('highest_category_raca', '')
    
    formatted_raca_negra_percentage = metrics.get('formatted_raca_negra_percentage', '0')
    formatted_raca_negra_counts = metrics.get('formatted_raca_negra_counts', '0')
    
    denominator_IST = metrics.get('denominator_IST', 0)
    highest_percentage_IST = metrics.get('highest_percentage_IST', 0)
    formatted_highest_count_IST = metrics.get('formatted_highest_count_IST', '0')
    highest_category_IST = metrics.get('highest_category_IST', '')
    
    prep_diaria_percent = metrics.get('prep_diaria_percent', '0')
    prep_diaria_count_formatted = metrics.get('prep_diaria_count_formatted', '0')
    prep_demand_percent = metrics.get('prep_demand_percent', '0')
    prep_demand_count_formatted = metrics.get('prep_demand_count_formatted', '0')

    # SLIDE 1: Título
    slide_layout = ppt.slide_layouts[6]
    slide = ppt.slides.add_slide(slide_layout)

    title_box = slide.shapes.add_textbox(Inches(4), Inches(2), width=Inches(6), height=Inches(1))
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = "Monitoramento da PrEP"
    p.font.name = 'Calibri Light'
    p.font.size = Pt(40)
    p.font.color.rgb = RGBColor(118, 113, 113)
    p.font.bold = True

    title_box = slide.shapes.add_textbox(Inches(4.8), Inches(2.7), width=Inches(4), height=Inches(1))
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = "Profilaxia Pré-Exposição"
    p.font.name = 'Calibri Light'
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(118, 113, 113)
    p.font.bold = True

    title_box = slide.shapes.add_textbox(Inches(5.5), Inches(5), width=Inches(2), height=Inches(1))
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = "Banco de " + str(hoje2)  
    p.font.name = 'Montserrat'
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(89, 89, 89)

    line1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3), Inches(2) - Pt(1), Inches(7.5), Pt(0.2))
    line1.fill.solid()
    line1.fill.fore_color.rgb = RGBColor(169, 169, 169)
    line1.line.fill.background()

    line2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3), Inches(3.5) + Pt(28), Inches(7.5), Pt(0.2))
    line2.fill.solid()
    line2.fill.fore_color.rgb = RGBColor(169, 169, 169)
    line2.line.fill.background()

    # SLIDE 2: Novos Usuários
    slide = ppt.slides.add_slide(ppt.slide_layouts[6])
    
    # Imagem (Fundo) - Restaurado posicionamento original
    chart_img_path = os.path.join(output_dir, 'PrEP_novosusuarios.png')
    if os.path.exists(chart_img_path):
        slide.shapes.add_picture(chart_img_path, Inches(1.5), Inches(1.0), height=Inches(6.0))
    
    # Texto (Frente)
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), width=Inches(8), height=Inches(4))
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = "Novos(as) usuários(as)"
    p.font.name = 'Montserrat'
    p.font.size = Pt(30)
    p.font.color.rgb = RGBColor(89, 89, 89)
    p.font.bold = True

    text = f"Apesar do impacto da pandemia de COVID-19 \nno número de novos(as) usuários(as) em 2020, \nobservou-se uma retomada com aumento \nexpressivo a partir de agosto de 2021.\n\nEm {hoje2}, {latest_month_year_count} novos(as) usuários(as) \nentraram em PrEP."
    text_box = slide.shapes.add_textbox(Inches(0.2), Inches(1.2), width=Inches(5), height=Inches(3))
    tf = text_box.text_frame
    p = tf.add_paragraph()
    p.text = text
    p.font.name = 'Calibri'
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(89, 89, 89)

    # SLIDE 3: Em PrEP por Ano
    slide = ppt.slides.add_slide(ppt.slide_layouts[6])
    
    # Imagem (Fundo) - Ajustado Left=3.0 (Direita)
    chart_img_path = os.path.join(output_dir, 'PrEP_emprep.png')
    if os.path.exists(chart_img_path):
        slide.shapes.add_picture(chart_img_path, Inches(3.0), Inches(1.5), height=Inches(6.0))
    
    # Texto
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), width=Inches(8), height=Inches(4))
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = "Usuários(as) com pelo menos uma dispensação \ne usuários(as) em PrEP por ano"
    p.font.name = 'Montserrat'
    p.font.size = Pt(30)
    p.font.color.rgb = RGBColor(89, 89, 89)
    p.font.bold = True

    text = f"A mediana de tempo de uso \nda PrEP foi de {mediana_uso:.0f} dias"
    text_box = slide.shapes.add_textbox(Inches(1.8), Inches(3.5), width=Inches(5), height=Inches(3))
    tf = text_box.text_frame
    p = tf.add_paragraph()
    p.text = text
    p.font.name = 'Calibri'
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(89, 89, 89)

    # SLIDE 4: Cascata
    slide = ppt.slides.add_slide(ppt.slide_layouts[6])
    
    # Imagem
    chart_img_path = os.path.join(output_dir, 'PrEP_cascata.png')
    if os.path.exists(chart_img_path):
        slide.shapes.add_picture(chart_img_path, Inches(0.3), Inches(1), height=Inches(6.0))
    
    # Texto
    title_box = slide.shapes.add_textbox(Inches(8), Inches(0.2), width=Inches(5), height=Inches(1))
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = "Usuários(as) de PreP"
    p.font.name = 'Montserrat'
    p.font.size = Pt(30)
    p.font.color.rgb = RGBColor(89, 89, 89)
    p.font.bold = True

    text = f"{emPrEP_porcent}% ({formatted_em_prep_count}) das pessoas que tiveram pelo menos \numa dispensação nos últimos 12 meses, \nestavam em PrEP em {hoje2}. \n\n{descontinuados_porcent}% ({formatted_discontinued_count}) das pessoas que tiveram pelo menos \numa dispensação nos últimos 12 meses, \nestavam descontinuadas em {hoje2}."
    text_box = slide.shapes.add_textbox(Inches(8), Inches(1.2), width=Inches(5), height=Inches(3))
    tf = text_box.text_frame
    p = tf.add_paragraph()
    p.text = text
    p.font.name = 'Calibri'
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(89, 89, 89)

    # SLIDE 5: Populações
    slide = ppt.slides.add_slide(ppt.slide_layouts[6])
    
    # Imagem
    chart_img_path = os.path.join(output_dir, 'PrEP_pop.png')
    if os.path.exists(chart_img_path):
        slide.shapes.add_picture(chart_img_path, Inches(0.3), Inches(1), height=Inches(6.0))
    
    # Texto
    title_box = slide.shapes.add_textbox(Inches(8), Inches(0.2), width=Inches(5), height=Inches(1))
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = "Usuários(as) em PreP"
    p.font.name = 'Montserrat'
    p.font.size = Pt(30)
    p.font.color.rgb = RGBColor(89, 89, 89)
    p.font.bold = True

    title_box = slide.shapes.add_textbox(Inches(8), Inches(0.8), width=Inches(5), height=Inches(1))
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = "Populações"
    p.font.name = 'Montserrat'
    p.font.size = Pt(30)
    p.font.color.rgb = RGBColor(166, 166, 166)
    p.font.bold = True

    text = f"Em {hoje2}, {highest_percentage_pop:.0f}% ({formatted_highest_count_pop}) das pessoas em PrEP \neram {highest_category_pop}."
    text_box = slide.shapes.add_textbox(Inches(8), Inches(1.8), width=Inches(5), height=Inches(3))
    tf = text_box.text_frame
    p = tf.add_paragraph()
    p.text = text
    p.font.name = 'Calibri'
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(89, 89, 89)

    # SLIDE 6: Faixa Etária
    slide = ppt.slides.add_slide(ppt.slide_layouts[6])
    
    # Imagem (Ajuste de altura reduzida)
    chart_img_path = os.path.join(output_dir, 'PrEP_fetar.png')
    if os.path.exists(chart_img_path):
        slide.shapes.add_picture(chart_img_path, Inches(0.3), Inches(0.8), height=Inches(5.5))
    
    title_box = slide.shapes.add_textbox(Inches(8), Inches(0.2), width=Inches(5), height=Inches(1))
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = "Usuários(as) em PreP"
    p.font.name = 'Montserrat'
    p.font.size = Pt(30)
    p.font.color.rgb = RGBColor(89, 89, 89)
    p.font.bold = True

    title_box = slide.shapes.add_textbox(Inches(8), Inches(0.8), width=Inches(5), height=Inches(1))
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = "Faixa etária"
    p.font.name = 'Montserrat'
    p.font.size = Pt(30)
    p.font.color.rgb = RGBColor(166, 166, 166)
    p.font.bold = True

    text = f"Em {hoje2}, {highest_percentage_fetar:.0f}% ({formatted_highest_count_fetar}) das pessoas em PrEP \ntinham de {highest_category_fetar} anos de idade. \n\nOs jovens (18 a 29 anos) representavam \n{formatted_young_percentage}% dos(as) usuários(as) ({formatted_young_counts}). \n\nA mediana de idade é de {median_age} anos."
    text_box = slide.shapes.add_textbox(Inches(8), Inches(1.8), width=Inches(5), height=Inches(3))
    tf = text_box.text_frame
    p = tf.add_paragraph()
    p.text = text
    p.font.name = 'Calibri'
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(89, 89, 89)

    # SLIDE 7: Escolaridade
    slide = ppt.slides.add_slide(ppt.slide_layouts[6])
    
    # Imagem
    chart_img_path = os.path.join(output_dir, 'PrEP_escol4.png')
    if os.path.exists(chart_img_path):
        slide.shapes.add_picture(chart_img_path, Inches(0.3), Inches(0.8), height=Inches(6.3))
    
    # Título e Texto (Mover para Extrema Esquerda: Inches(0.3))
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), width=Inches(5), height=Inches(1))
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = "Usuários(as) em PreP"
    p.font.name = 'Montserrat'
    p.font.size = Pt(30)
    p.font.color.rgb = RGBColor(89, 89, 89)
    p.font.bold = True

    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.8), width=Inches(5), height=Inches(1))
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = "Escolaridade"
    p.font.name = 'Montserrat'
    p.font.size = Pt(30)
    p.font.color.rgb = RGBColor(166, 166, 166)
    p.font.bold = True

    text = f"Em {hoje2}, {highest_percentage_escol:.0f}% ({formatted_highest_count_escol}) das pessoas em PrEP \ntinham {highest_category_escol} de estudo."
    text_box = slide.shapes.add_textbox(Inches(0.3), Inches(1.8), width=Inches(5), height=Inches(3))
    tf = text_box.text_frame
    p = tf.add_paragraph()
    p.text = text
    p.font.name = 'Calibri'
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(89, 89, 89)

    # SLIDE 8: Raça
    slide = ppt.slides.add_slide(ppt.slide_layouts[6])
    
    # Imagem
    chart_img_path = os.path.join(output_dir, 'PrEP_raca.png')
    if os.path.exists(chart_img_path):
        slide.shapes.add_picture(chart_img_path, Inches(0.3), Inches(0.8), height=Inches(6.3))
    
    title_box = slide.shapes.add_textbox(Inches(8), Inches(0.2), width=Inches(5), height=Inches(1))
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = "Usuários(as) em PreP"
    p.font.name = 'Montserrat'
    p.font.size = Pt(30)
    p.font.color.rgb = RGBColor(89, 89, 89)
    p.font.bold = True

    title_box = slide.shapes.add_textbox(Inches(8), Inches(0.8), width=Inches(5), height=Inches(1))
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = "Raça/cor"
    p.font.name = 'Montserrat'
    p.font.size = Pt(30)
    p.font.color.rgb = RGBColor(166, 166, 166)
    p.font.bold = True

    text = f"Em {hoje2}, {highest_percentage_raca:.0f}% ({formatted_highest_count_raca}) das pessoas em PrEP \nautodeclararam-se {highest_category_raca}. \n\nPessoas negras (pretas + pardas) \nrepresentaram {formatted_raca_negra_percentage}% ({formatted_raca_negra_counts}) dos(as) usuários(as)."
    text_box = slide.shapes.add_textbox(Inches(8), Inches(1.8), width=Inches(5), height=Inches(3))
    tf = text_box.text_frame
    p = tf.add_paragraph()
    p.text = text
    p.font.name = 'Calibri'
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(89, 89, 89)

    # SLIDE 9: IST
    slide = ppt.slides.add_slide(ppt.slide_layouts[6])
    
    # Imagem
    chart_img_path = os.path.join(output_dir, 'PrEP_IST.png')
    if os.path.exists(chart_img_path):
        slide.shapes.add_picture(chart_img_path, Inches(0.3), Inches(1.5), height=Inches(6.0))
    
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), width=Inches(8), height=Inches(4))
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = "Infecção Sexualmente Transmissível (IST)"
    p.font.name = 'Montserrat'
    p.font.size = Pt(30)
    p.font.color.rgb = RGBColor(89, 89, 89)
    p.font.bold = True

    text_box = slide.shapes.add_textbox(Inches(8), Inches(2.5), width=Inches(5), height=Inches(1))
    tf = text_box.text_frame
    p = tf.add_paragraph()
    p.text = "Nos últimos 3 meses, \no usuário(a) tem ou teve \nalgum sinal/sintoma \nou foi diagnosticado(a) com IST?"
    p.font.name = 'Montserrat'
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(89, 89, 89)
    p.font.bold = True

    # Texto corrigido para Slide 9
    text = f"Das {denominator_IST} pessoas com informação preenchida, \n{highest_percentage_IST:.0f}% ({formatted_highest_count_IST}) não tinham {highest_category_IST}."
    text_box = slide.shapes.add_textbox(Inches(8), Inches(4), width=Inches(5), height=Inches(3))
    tf = text_box.text_frame
    p = tf.add_paragraph()
    p.text = text
    p.font.name = 'Calibri'
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(89, 89, 89)

    # SLIDE 10: Modalidades
    slide = ppt.slides.add_slide(ppt.slide_layouts[6])
    
    # Imagem
    chart_img_path = os.path.join(output_dir, 'PrEP_modalidades.png')
    if os.path.exists(chart_img_path):
        slide.shapes.add_picture(chart_img_path, Inches(3.0), Inches(3), height=Inches(4))
    
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), width=Inches(8), height=Inches(4))
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = "Modalidade da PrEP" 
    p.font.name = 'Montserrat'
    p.font.size = Pt(30)
    p.font.color.rgb = RGBColor(89, 89, 89)
    p.font.bold = True

    text = "*Considera apenas quem teve 2 dispensações ou mais"
    text_box = slide.shapes.add_textbox(Inches(0.2), Inches(6.5), width=Inches(5), height=Inches(0.5))
    tf = text_box.text_frame
    p = tf.add_paragraph()
    p.text = text
    p.font.name = 'Calibri'
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(137, 137, 137)

    text = "Desde a última dispensa, em média, como você tomou a PrEP*?"
    text_box = slide.shapes.add_textbox(Inches(3.0), Inches(2.3), width=Inches(5), height=Inches(3))
    tf = text_box.text_frame
    p = tf.add_paragraph()
    p.text = text
    p.font.name = 'Calibri'
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(31, 73, 125)

    text = "Em dezembro de 2022, por meio da Nota Técnica nº \n563/2022-CGAHV/.DCCI/SVS/MS, o MS incluiu a \nmodalidade de “PrEP sob demanda” como \nalternativa de uso."
    text_box = slide.shapes.add_textbox(Inches(0.2), Inches(1.0), width=Inches(5), height=Inches(3))
    tf = text_box.text_frame
    p = tf.add_paragraph()
    p.text = text
    p.font.name = 'Calibri'
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(89, 89, 89)

    text_box = slide.shapes.add_textbox(Inches(7.0), Inches(1), width=Inches(2.5), height=Inches(2.0))
    tf = text_box.text_frame
    p1 = tf.add_paragraph()
    p1.text = f"{prep_diaria_percent}% ({prep_diaria_count_formatted}) das dispensações foram para PrEP diária"
    p1.font.name = 'Calibri'
    p1.font.size = Pt(16)
    p1.font.color.rgb = RGBColor(89, 89, 89)
    tf.add_paragraph()
    p2 = tf.add_paragraph()
    p2.text = f"{prep_demand_percent}% ({prep_demand_count_formatted}) das dispensações foram para PrEP sob demanda"
    p2.font.name = 'Calibri'
    p2.font.size = Pt(16)
    p2.font.color.rgb = RGBColor(89, 89, 89)

    # Footer
    for slide in ppt.slides:
        source_text_box = slide.shapes.add_textbox(Inches(0.2), Inches(6.8), width=Inches(2.5), height=Inches(0.5))
        source_text = source_text_box.text_frame.add_paragraph()
        source_text.text = "Fonte: DATHI/SVS/MS"
        source_text.font.size = Pt(10)
        source_text.font.color.rgb = RGBColor(137, 137, 137)

    # Save
    data_dt = pd.to_datetime(data_fechamento)
    mes = int(data_dt.month)
    ano = int(data_dt.year)
    
    filename = "Monitoramento_PrEP_{:02d}_{}.pptx".format(mes, ano)
    save_path = os.path.join(output_dir, filename)
    
    print(f"Salvando PPT em: {save_path}")
    ppt.save(save_path)
    print("PPT Salvo com sucesso.")